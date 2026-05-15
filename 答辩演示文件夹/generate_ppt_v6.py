from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(0x0C, 0x0C, 0x0C)
BG2 = RGBColor(0x14, 0x14, 0x14)
WHITE = RGBColor(0xF0, 0xF0, 0xF0)
MUTED = RGBColor(0x7A, 0x7A, 0x8A)
DIM = RGBColor(0x4A, 0x4A, 0x5A)
RED = RGBColor(0xC4, 0x1E, 0x3A)
GOLD = RGBColor(0xD4, 0xA5, 0x6E)
DARK_RED = RGBColor(0x3A, 0x0A, 0x14)

SERIF = "Source Han Serif SC"
SANS = "Source Han Sans SC"
MONO = "Consolas"

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "defense_v6.pptx")


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        if line_width:
            sh.line.width = Pt(line_width)
    else:
        sh.line.fill.background()
    return sh


def add_text(slide, left, top, width, height, text, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font=SANS, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    try:
        rPr = p.runs[0]._r.get_or_add_rPr()
        rPr.set(qn('a:lang'), 'zh-CN')
    except:
        pass
    return tb


def add_multiline(slide, left, top, width, height, lines, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font=SANS, spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(size * (spacing - 1))
        try:
            rPr = p.runs[0]._r.get_or_add_rPr()
            rPr.set(qn('a:lang'), 'zh-CN')
        except:
            pass
    return tb


def gradient_line(slide, left, top, width, height=Pt(3)):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = sh.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RED
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = GOLD
    fill.gradient_stops[1].position = 1.0
    sh.line.fill.background()
    return sh


def page_num(slide, num, total=22):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f"{num} / {total}", size=10, color=DIM, align=PP_ALIGN.RIGHT, font=MONO)


def section_tag(slide, text):
    add_text(slide, Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.35),
             text, size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=MONO)


def content_slide(prs, section, title, items, page_idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    section_tag(slide, section)
    add_text(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             title, size=36, color=WHITE, bold=True, font=SERIF)
    gradient_line(slide, Inches(1.2), Inches(1.65), Inches(2.5))
    y = Inches(2.2)
    for item in items:
        if item.startswith("##"):
            add_text(slide, Inches(1.2), y, Inches(10), Inches(0.5),
                     item[2:].strip(), size=20, color=GOLD, bold=True, font=SANS)
            y += Inches(0.55)
        elif item.startswith("**"):
            label = item.replace("**", "").strip()
            add_text(slide, Inches(1.2), y, Inches(10), Inches(0.4),
                     label, size=15, color=RED, bold=True, font=SANS)
            y += Inches(0.45)
        elif item.startswith(">>"):
            val = item[2:].strip()
            add_text(slide, Inches(1.5), y, Inches(10), Inches(0.4),
                     val, size=14, color=MUTED, font=SANS)
            y += Inches(0.4)
        elif item.startswith("::"):
            big = item[2:].strip()
            add_text(slide, Inches(1.2), y, Inches(10), Inches(0.9),
                     big, size=48, color=RED, bold=True, font=SERIF, align=PP_ALIGN.LEFT)
            y += Inches(0.95)
        else:
            add_text(slide, Inches(1.5), y, Inches(9.5), Inches(0.4),
                     item, size=14, color=MUTED, font=SANS)
            y += Inches(0.38)
    page_num(slide, page_idx)
    return slide


def two_col_slide(prs, section, title, left_items, right_items, left_title="", right_title="", page_idx=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    section_tag(slide, section)
    add_text(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             title, size=36, color=WHITE, bold=True, font=SERIF)
    gradient_line(slide, Inches(1.2), Inches(1.65), Inches(2.5))
    if left_title:
        add_text(slide, Inches(1.2), Inches(2.1), Inches(5), Inches(0.4),
                 left_title, size=18, color=GOLD, bold=True, font=SANS)
    if right_title:
        add_text(slide, Inches(7.2), Inches(2.1), Inches(5), Inches(0.4),
                 right_title, size=18, color=GOLD, bold=True, font=SANS)
    add_shape(slide, Inches(6.5), Inches(2.3), Pt(1), Inches(4.2), DIM)
    y = Inches(2.6)
    for item in left_items:
        if item.startswith("**"):
            add_text(slide, Inches(1.2), y, Inches(5), Inches(0.35),
                     item.replace("**", "").strip(), size=14, color=RED, bold=True, font=SANS)
            y += Inches(0.4)
        else:
            add_text(slide, Inches(1.5), y, Inches(4.8), Inches(0.35),
                     item, size=13, color=MUTED, font=SANS)
            y += Inches(0.38)
    y = Inches(2.6)
    for item in right_items:
        if item.startswith("**"):
            add_text(slide, Inches(7.2), y, Inches(5), Inches(0.35),
                     item.replace("**", "").strip(), size=14, color=RED, bold=True, font=SANS)
            y += Inches(0.4)
        else:
            add_text(slide, Inches(7.5), y, Inches(4.8), Inches(0.35),
                     item, size=13, color=MUTED, font=SANS)
            y += Inches(0.38)
    page_num(slide, page_idx)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # === 1. COVER ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_shape(s, Inches(0), Inches(0), Inches(0.15), H, RED)
    add_text(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "基于CRS推荐与AI数字人的", size=22, color=MUTED, font=SANS)
    add_text(s, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
             "非遗文化传播系统实现", size=52, color=WHITE, bold=True, font=SERIF)
    gradient_line(s, Inches(1.5), Inches(4.1), Inches(4))
    add_text(s, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
             "An Intangible Cultural Heritage Communication System Based on CRS Recommendation and AI Digital Human",
             size=12, color=DIM, font=SANS)
    add_text(s, Inches(1.5), Inches(5.5), Inches(5), Inches(0.35),
             "答辩人：王子轩", size=16, color=WHITE, font=SANS)
    add_text(s, Inches(1.5), Inches(5.95), Inches(5), Inches(0.35),
             "指导教师：徐龙琴 教授 / 高静 高级工程师", size=13, color=MUTED, font=SANS)
    add_text(s, Inches(1.5), Inches(6.4), Inches(5), Inches(0.35),
             "信息科学与技术学院 · 数据科学与大数据技术222班", size=13, color=MUTED, font=SANS)
    add_text(s, Inches(1.5), Inches(6.85), Inches(5), Inches(0.35),
             "仲恺农业工程学院", size=13, color=MUTED, font=SANS)
    add_text(s, Inches(9), Inches(6.85), Inches(3.5), Inches(0.35),
             "2026.05.09", size=13, color=DIM, align=PP_ALIGN.RIGHT, font=MONO)

    # === 2. OUTLINE ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_text(s, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             "汇报提纲", size=36, color=WHITE, bold=True, font=SERIF)
    gradient_line(s, Inches(1.2), Inches(1.65), Inches(2.5))
    outline = [
        ("01", "绪论", "研究背景、意义、现状"),
        ("02", "相关技术与理论", "技术栈与理论基础"),
        ("03", "系统分析", "需求分析、可行性、数据流程"),
        ("04", "系统设计", "架构、推荐、AI数字人、知识图谱"),
        ("05", "系统实现", "小程序端、Web管理端"),
        ("06", "系统测试", "功能、性能、兼容性测试"),
        ("07", "实验与分析", "AI对话质量评估"),
        ("08", "结论与展望", "总结与未来方向"),
    ]
    y = Inches(2.3)
    for num, title, desc in outline:
        add_text(s, Inches(1.5), y, Inches(0.8), Inches(0.5),
                 num, size=28, color=RED, bold=True, font=MONO)
        add_text(s, Inches(2.5), y, Inches(3), Inches(0.5),
                 title, size=18, color=WHITE, bold=True, font=SANS)
        add_text(s, Inches(5.8), y, Inches(5), Inches(0.5),
                 desc, size=13, color=MUTED, font=SANS)
        y += Inches(0.6)
    page_num(s, 2)

    # === 3. RESEARCH BACKGROUND ===
    content_slide(prs, "01 绪论", "研究背景", [
        "::2.2亿+",
        "**抖音平台国家级非遗相关视频数量",
        "累计获赞138亿次，转发45亿次",
        ">>87万+万粉创作者投稿非遗内容",
        "",
        "::1000+",
        "**国家级非遗代表性项目",
        "中国列入联合国非遗名录数量世界第一",
        ">>但大量非遗项目面临传承困境",
        "",
        "##三大传播趋势",
        "**短视频化 — 抖音非遗内容爆发式增长",
        "**资源数字化 — 动画、游戏、电影均需数字化非遗资源",
        "**传播社区化 — 以兴趣为纽带的论坛传承，星火燎原",
    ], 3)

    # === 4. RESEARCH SIGNIFICANCE ===
    content_slide(prs, "01 绪论", "研究意义", [
        "##理论意义",
        "探索CRS会话式推荐在文化领域的应用范式",
        "构建非遗知识图谱与推荐系统的融合框架",
        "丰富AI数字人在文化传播场景中的理论依据",
        "",
        "##实践意义",
        "以微信小程序为载体，降低非遗传播的触达门槛",
        "个性化推荐提升用户参与度与内容匹配效率",
        "AI数字人'黑塔'提供沉浸式导览与多轮问答体验",
        "为非遗数字化传播与个性化服务提供实践参考",
    ], 4)

    # === 5. RESEARCH STATUS ===
    two_col_slide(prs, "01 绪论", "国内外研究现状",
        [
            "**国内研究",
            "非遗数字化：从保存记录转向传播创作",
            "推荐系统：协同过滤→深度学习→会话式推荐",
            "AI数字人：虚拟主播、智能客服逐步落地",
            "微信生态：小程序成为轻量级应用首选",
            "",
            "**研究空白",
            "CRS推荐在非遗领域应用极少",
            "AI数字人+推荐系统缺乏深度整合",
            "非遗传播缺少社区化+个性化闭环",
        ],
        [
            "**国外研究",
            "CRS：基于对话的推荐系统成为热点方向",
            "知识图谱：Google Knowledge Graph推动语义理解",
            "数字人：NLP+TTS+Avatar技术日趋成熟",
            "LLM：GPT系列推动对话质量飞跃提升",
            "",
            "**发展趋势",
            "从静态展示到动态交互",
            "从通用推荐到会话式个性化推荐",
            "从单向传播到双向互动社区",
        ],
        left_title="国内", right_title="国外", page_idx=5)

    # === 6. PAPER STRUCTURE ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    section_tag(s, "01 绪论")
    add_text(s, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             "论文组织结构", size=36, color=WHITE, bold=True, font=SERIF)
    gradient_line(s, Inches(1.2), Inches(1.65), Inches(2.5))
    chapters = [
        ("1", "绪论", "背景、意义、现状"),
        ("2", "相关技术与理论", "Python / FastAPI / LLM / 提示工程"),
        ("3", "系统分析", "需求分析、可行性、数据流程"),
        ("4", "系统设计", "架构、推荐、AI数字人、知识图谱、数据库"),
        ("5", "系统实现", "小程序端 + Web管理端"),
        ("6", "系统测试", "功能 / 性能 / 兼容性"),
        ("7", "实验与分析", "AI对话质量评估"),
        ("8", "结论与展望", "总结与未来工作"),
    ]
    y = Inches(2.2)
    for i, (ch, title, desc) in enumerate(chapters):
        x_base = Inches(1.2) if i < 4 else Inches(7.0)
        y_pos = Inches(2.2) + Inches(1.1) * (i % 4)
        add_shape(s, x_base, y_pos, Inches(0.5), Inches(0.5), DARK_RED)
        add_text(s, x_base + Inches(0.05), y_pos + Inches(0.05), Inches(0.4), Inches(0.4),
                 ch, size=20, color=RED, bold=True, font=MONO, align=PP_ALIGN.CENTER)
        add_text(s, x_base + Inches(0.7), y_pos, Inches(4.5), Inches(0.35),
                 title, size=17, color=WHITE, bold=True, font=SANS)
        add_text(s, x_base + Inches(0.7), y_pos + Inches(0.35), Inches(4.5), Inches(0.3),
                 desc, size=12, color=MUTED, font=SANS)
    page_num(s, 6)

    # === 7. TECH STACK ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    section_tag(s, "02 相关技术与理论")
    add_text(s, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             "相关技术与理论", size=36, color=WHITE, bold=True, font=SERIF)
    gradient_line(s, Inches(1.2), Inches(1.65), Inches(2.5))
    techs = [
        ("Python", "后端核心语言，FastAPI框架构建RESTful API"),
        ("FastAPI", "高性能异步Web框架，自动生成OpenAPI文档"),
        ("SQLAlchemy", "ORM框架，实现数据库模型映射与查询"),
        ("SQLite", "轻量级嵌入式数据库，零配置部署"),
        ("WXML/WXSS/JS", "微信小程序前端开发技术栈"),
        ("LLM + 提示工程", "大语言模型驱动AI数字人对话能力"),
        ("知识图谱", "Neo4j构建非遗实体关系网络"),
        ("CRS推荐", "会话式推荐系统，动态调整推荐策略"),
    ]
    y = Inches(2.2)
    for i, (name, desc) in enumerate(techs):
        col = 0 if i < 4 else 1
        row = i % 4
        x = Inches(1.2) + Inches(5.8) * col
        yp = Inches(2.2) + Inches(1.15) * row
        add_text(s, x, yp, Inches(5), Inches(0.35),
                 name, size=16, color=GOLD, bold=True, font=MONO)
        add_text(s, x, yp + Inches(0.38), Inches(5.2), Inches(0.35),
                 desc, size=13, color=MUTED, font=SANS)
    page_num(s, 7)

    # === 8. REQUIREMENTS ===
    content_slide(prs, "03 系统分析", "功能需求分析", [
        "##核心功能模块",
        "**非遗内容浏览 — 分类展示、详情查看、搜索筛选",
        "**智能导览 — AI数字人多轮问答、语音播报",
        "**个性化推荐 — CRS会话式推荐、冷启动→精准推荐",
        "**活动报名 — 非遗活动发布、在线报名、签到",
        "**社区互动 — 帖子发布、评论点赞、兴趣圈子",
        "**用户中心 — 偏好管理、历史记录、收藏关注",
        "",
        "##管理端功能",
        "**内容管理 / 活动管理 / 帖子管理 / 用户管理 / 统计 / 知识库",
    ], 8)

    # === 9. DATA FLOW ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    section_tag(s, "03 系统分析")
    add_text(s, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             "系统数据流程", size=36, color=WHITE, bold=True, font=SERIF)
    gradient_line(s, Inches(1.2), Inches(1.65), Inches(2.5))
    flow_steps = [
        ("1", "用户操作", "浏览、搜索、对话、报名"),
        ("2", "行为采集", "点击流、停留时长、CRS对话记录"),
        ("3", "画像构建", "偏好标签、兴趣向量、场景信号"),
        ("4", "推荐引擎", "冷启动→混合推荐→精准推荐"),
        ("5", "内容分发", "个性化列表、延伸推荐、行动清单"),
        ("6", "反馈闭环", "推荐反馈→ASK追问→策略调整"),
    ]
    y = Inches(2.3)
    for num, title, desc in flow_steps:
        add_shape(s, Inches(1.2), y, Inches(0.55), Inches(0.55), DARK_RED)
        add_text(s, Inches(1.25), y + Inches(0.05), Inches(0.45), Inches(0.45),
                 num, size=20, color=RED, bold=True, font=MONO, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.0), y + Inches(0.02), Inches(3), Inches(0.35),
                 title, size=16, color=WHITE, bold=True, font=SANS)
        add_text(s, Inches(2.0), y + Inches(0.35), Inches(8), Inches(0.3),
                 desc, size=13, color=MUTED, font=SANS)
        if num != "6":
            add_shape(s, Inches(1.45), y + Inches(0.6), Pt(2), Inches(0.5), DIM)
        y += Inches(0.85)
    page_num(s, 9)

    # === 10. ARCHITECTURE ===
    content_slide(prs, "04 系统设计", "系统整体架构", [
        "##四层架构设计",
        "",
        "**表现层（Presentation）",
        "微信小程序端 + Web管理端，WXML/WXSS/JS渲染",
        "",
        "**接口层（API Gateway）",
        "FastAPI构建RESTful API，JWT认证，请求路由分发",
        "",
        "**业务逻辑层（Business Logic）",
        "推荐引擎 / AI数字人 / 知识图谱 / 内容治理 / 社区互动",
        "",
        "**数据层（Data）",
        "SQLite关系数据 + Neo4j知识图谱 + LLM API",
    ], 10)

    # === 11. RECOMMENDATION MODULE ===
    content_slide(prs, "04 系统设计", "推荐模块设计", [
        "##CRS会话式推荐三阶段",
        "",
        "::冷启动",
        "**新用户无历史行为时",
        "基于热门内容+分类偏好初始化推荐",
        ">>用户首次选择兴趣标签触发推荐",
        "",
        "::混合推荐",
        "**结合协同过滤与内容特征",
        "用户画像+行为记录+场景信号综合计算",
        ">>ASK追问机制收集隐性偏好",
        "",
        "::精准推荐",
        "**CRS对话反馈持续优化",
        "推荐反馈→策略调整→延伸推荐→行动清单",
    ], 11)

    # === 12. AI DIGITAL HUMAN ===
    content_slide(prs, "04 系统设计", "AI数字人设计", [
        "##AI数字人'黑塔'",
        "",
        "**多轮问答能力",
        "基于LLM+知识图谱，理解非遗领域问题并生成回答",
        "",
        "**语音播报",
        "TTS文字转语音，支持方言与普通话切换",
        "",
        "**CRS状态显示",
        "实时展示推荐推理过程，增强透明度与信任感",
        "",
        "**延伸推荐",
        "对话中识别兴趣点，主动推送相关非遗内容",
        "",
        "**行动清单生成",
        "根据对话内容生成可执行的学习/体验建议",
    ], 12)

    # === 13. KNOWLEDGE GRAPH ===
    two_col_slide(prs, "04 系统设计", "知识图谱与数据库设计",
        [
            "**知识图谱",
            "Neo4j图数据库存储非遗实体关系",
            "节点：非遗项目、传承人、技艺、地区",
            "边：传承关系、地域关联、技艺相似",
            "支撑AI数字人语义理解与推荐推理",
            "",
            "**核心数据表",
            "用户表（偏好、画像）",
            "非遗内容表（分类、标签）",
            "活动表（时间、地点、名额）",
            "帖子/评论表（社区互动）",
            "对话记录表（CRS会话历史）",
        ],
        [
            "**数据库设计原则",
            "SQLite轻量部署，零配置运维",
            "SQLAlchemy ORM统一数据访问层",
            "外键约束保证数据一致性",
            "索引优化高频查询性能",
            "",
            "**接口设计",
            "RESTful API规范",
            "JWT Token认证鉴权",
            "OpenAPI自动文档生成",
            "统一响应格式与错误码",
            "请求频率限制与日志记录",
        ],
        left_title="知识图谱 & 数据", right_title="接口设计", page_idx=13)

    # === 14. MINI PROGRAM ===
    content_slide(prs, "05 系统实现", "小程序端功能模块", [
        "##首页模块",
        "轮播推荐、分类导航、热门内容、活动入口",
        "",
        "##非遗内容模块",
        "分类浏览、详情页、搜索筛选、收藏分享",
        "",
        "##活动报名模块",
        "活动列表、在线报名、签到打卡、名额管理",
        "",
        "##社区讨论模块",
        "帖子发布、评论互动、点赞收藏、兴趣圈子",
        "",
        "##用户中心模块",
        "偏好设置、浏览历史、收藏关注、消息通知",
    ], 14)

    # === 15. AI DIGITAL HUMAN MODULE ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    section_tag(s, "05 系统实现")
    add_text(s, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             "AI数字人模块实现", size=36, color=WHITE, bold=True, font=SERIF)
    gradient_line(s, Inches(1.2), Inches(1.65), Inches(2.5))

    features = [
        ("对话界面", "聊天气泡式交互，支持文字输入与语音识别"),
        ("语音播报", "TTS引擎将回答转为语音，可调节语速音色"),
        ("CRS状态", "可视化展示推荐推理路径与置信度"),
        ("延伸推荐", "对话中智能插入相关非遗内容卡片"),
        ("行动清单", "一键生成学习计划与体验活动推荐"),
    ]
    y = Inches(2.2)
    for i, (name, desc) in enumerate(features):
        add_shape(s, Inches(1.2), y, Inches(10.8), Inches(0.7), BG2)
        add_text(s, Inches(1.5), y + Inches(0.08), Inches(2.5), Inches(0.3),
                 name, size=15, color=GOLD, bold=True, font=SANS)
        add_text(s, Inches(4.2), y + Inches(0.08), Inches(7.5), Inches(0.55),
                 desc, size=13, color=MUTED, font=SANS)
        y += Inches(0.85)

    add_text(s, Inches(1.2), y + Inches(0.3), Inches(10), Inches(0.5),
             "技术实现：LLM API + 提示工程 + 知识图谱查询 + TTS + 前端状态管理",
             size=13, color=DIM, font=SANS)
    page_num(s, 15)

    # === 16. WEB ADMIN ===
    two_col_slide(prs, "05 系统实现", "Web管理端功能模块",
        [
            "**内容管理",
            "非遗内容CRUD、分类标签管理、上下架控制",
            "",
            "**活动管理",
            "活动发布编辑、报名审核、签到管理",
            "",
            "**帖子管理",
            "社区帖子审核、评论管理、违规处理",
        ],
        [
            "**用户管理",
            "用户列表、角色权限、封禁解封",
            "",
            "**统计模块",
            "访问量、活跃度、推荐效果、内容热度",
            "",
            "**知识库管理",
            "知识图谱节点编辑、关系维护、语料导入",
        ],
        left_title="内容 & 运营", right_title="用户 & 数据", page_idx=16)

    # === 17. TESTING ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    section_tag(s, "06 系统测试")
    add_text(s, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             "系统测试", size=36, color=WHITE, bold=True, font=SERIF)
    gradient_line(s, Inches(1.2), Inches(1.65), Inches(2.5))

    test_sections = [
        ("功能测试", "全部核心功能通过，覆盖用户登录、内容浏览、AI对话、推荐、活动报名、社区互动等场景"),
        ("性能测试", "API平均响应<200ms，并发50用户无异常，页面加载<2s"),
        ("兼容性测试", "iOS/Android双端适配，主流机型与微信版本均正常运行"),
    ]
    y = Inches(2.3)
    for name, desc in test_sections:
        add_text(s, Inches(1.2), y, Inches(3), Inches(0.4),
                 name, size=18, color=GOLD, bold=True, font=SANS)
        add_text(s, Inches(1.2), y + Inches(0.45), Inches(10.5), Inches(0.6),
                 desc, size=14, color=MUTED, font=SANS)
        y += Inches(1.3)

    add_shape(s, Inches(1.2), y, Inches(10.8), Pt(1), DIM)
    add_text(s, Inches(1.2), y + Inches(0.2), Inches(10), Inches(0.4),
             "测试工具：Postman / Chrome DevTools / 微信开发者工具", size=12, color=DIM, font=MONO)
    page_num(s, 17)

    # === 18. EXPERIMENTS ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    section_tag(s, "07 实验与分析")
    add_text(s, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             "实验与分析", size=36, color=WHITE, bold=True, font=SERIF)
    gradient_line(s, Inches(1.2), Inches(1.65), Inches(2.5))

    add_text(s, Inches(1.2), Inches(2.2), Inches(5), Inches(0.4),
             "AI对话质量评估", size=20, color=GOLD, bold=True, font=SANS)

    metrics = [
        ("相关性", "回答与问题的语义相关程度"),
        ("准确性", "非遗知识的事实正确性"),
        ("流畅性", "语言表达的自然与连贯程度"),
        ("有用性", "回答对用户实际需求的满足程度"),
    ]
    y = Inches(2.8)
    for name, desc in metrics:
        add_text(s, Inches(1.5), y, Inches(1.8), Inches(0.35),
                 name, size=15, color=RED, bold=True, font=SANS)
        add_text(s, Inches(3.5), y, Inches(8), Inches(0.35),
                 desc, size=13, color=MUTED, font=SANS)
        y += Inches(0.5)

    add_text(s, Inches(1.2), y + Inches(0.4), Inches(5), Inches(0.4),
             "实验环境", size=20, color=GOLD, bold=True, font=SANS)
    add_text(s, Inches(1.5), y + Inches(0.9), Inches(10), Inches(0.35),
             "数据集：自建非遗知识库 + 真实用户对话记录", size=13, color=MUTED, font=SANS)
    add_text(s, Inches(1.5), y + Inches(1.25), Inches(10), Inches(0.35),
             "评估方法：人工标注 + 自动指标（BLEU/ROUGE）结合", size=13, color=MUTED, font=SANS)
    page_num(s, 18)

    # === 19. CONCLUSION ===
    content_slide(prs, "08 结论与展望", "结论与展望", [
        "##主要成果",
        "构建了基于CRS推荐与AI数字人的非遗文化传播系统",
        "实现冷启动→混合推荐→精准推荐的推荐策略过渡",
        "设计AI数字人'黑塔'，支持多轮问答、语音播报、延伸推荐",
        "集成知识图谱，提升AI对话的领域理解能力",
        "以微信小程序为载体，降低非遗传播触达门槛",
        "",
        "##未来展望",
        "引入更多非遗门类的知识图谱数据",
        "优化CRS推荐算法，提升冷启动阶段推荐精度",
        "增强AI数字人的情感表达与多模态交互能力",
        "拓展AR/VR沉浸式非遗体验场景",
    ], 19)

    # === 20. THANK YOU ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_shape(s, Inches(0), Inches(0), Inches(0.15), H, RED)
    add_text(s, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
             "感谢聆听", size=56, color=WHITE, bold=True, font=SERIF, align=PP_ALIGN.LEFT)
    gradient_line(s, Inches(1.5), Inches(4.2), Inches(4))
    add_text(s, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
             "Thank You", size=24, color=MUTED, font=SANS)
    add_text(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4),
             "答辩人：王子轩 | 指导教师：徐龙琴 教授 / 高静 高级工程师", size=14, color=DIM, font=SANS)
    add_text(s, Inches(1.5), Inches(5.95), Inches(10), Inches(0.4),
             "仲恺农业工程学院 · 信息科学与技术学院", size=13, color=DIM, font=SANS)
    page_num(s, 20)

    prs.save(OUT)
    print(f"saved: {OUT}")
    print(f"size: {os.path.getsize(OUT)} bytes")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
