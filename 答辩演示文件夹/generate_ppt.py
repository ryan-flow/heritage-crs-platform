from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

BG_COLOR = RGBColor(0xF5, 0xF0, 0xE8)
TITLE_COLOR = RGBColor(0x8B, 0x22, 0x22)
SUBTITLE_COLOR = RGBColor(0x5C, 0x3A, 0x21)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
ACCENT_COLOR = RGBColor(0xC4, 0x1E, 0x3A)
LIGHT_ACCENT = RGBColor(0xE8, 0xD5, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x2C, 0x1A, 0x12)
GOLD = RGBColor(0xD4, 0xA5, 0x37)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=BODY_COLOR, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=BODY_COLOR, spacing=Pt(6), bold_prefix=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        if bold_prefix and "——" in item:
            parts = item.split("——", 1)
            run1 = p.add_run()
            run1.text = parts[0] + "——"
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = ACCENT_COLOR
            run1.font.bold = True
            run1.font.name = "微软雅黑"
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.bold = False
            run2.font.name = "微软雅黑"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = False
            run.font.name = "微软雅黑"
    return txBox


def add_decorated_title(slide, text, top=Inches(0.4)):
    add_shape(slide, Inches(0.5), top, Inches(0.08), Inches(0.55), ACCENT_COLOR)
    add_text_box(slide, Inches(0.75), top, Inches(10), Inches(0.6), text, font_size=28, color=TITLE_COLOR, bold=True)
    add_shape(slide, Inches(0.5), top + Inches(0.65), Inches(12.3), Inches(0.02), LIGHT_ACCENT)


def add_page_number(slide, num, total):
    add_text_box(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35),
                 f"{num}/{total}", font_size=11, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.RIGHT)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]
    total_slides = 22

    # ===== Slide 1: Cover =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ACCENT_COLOR)
    add_shape(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), ACCENT_COLOR)

    add_shape(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.02), GOLD)
    add_text_box(slide, Inches(1.5), Inches(1.4), Inches(10.3), Inches(1.2),
                 "基于CRS推荐与AI数字人的\n非遗文化传播系统实现",
                 font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(1.5), Inches(2.9), Inches(10.3), Inches(0.02), GOLD)

    info_items = [
        "学    院：信息科学与技术学院",
        "专    业：数据科学与大数据技术222班",
        "姓    名：王子轩",
        "学    号：202210274225",
        "指导教师：徐龙琴 教授  高静 高级工程师（校外）",
    ]
    for i, item in enumerate(info_items):
        add_text_box(slide, Inches(3.5), Inches(3.3 + i * 0.5), Inches(6.3), Inches(0.45),
                     item, font_size=18, color=RGBColor(0xDD, 0xCC, 0xBB), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.4),
                 "仲恺农业工程学院", font_size=22, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.5), Inches(6.6), Inches(6.3), Inches(0.4),
                 "2026年5月9日", font_size=16, color=RGBColor(0xAA, 0x99, 0x88), alignment=PP_ALIGN.CENTER)

    # ===== Slide 2: TOC =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "汇报提纲")
    add_page_number(slide, 2, total_slides)

    toc_items = [
        ("01", "研究背景与意义", "非遗数字化传播的时代需求与课题定位"),
        ("02", "国内外研究现状", "推荐系统、CRS、AI数字人研究进展与不足"),
        ("03", "相关技术与理论", "FastAPI、微信小程序、LLM、CRS理论"),
        ("04", "系统分析", "需求分析、可行性分析、数据流程"),
        ("05", "系统设计", "架构设计、推荐模块、AI数字人、知识图谱"),
        ("06", "系统实现", "小程序端与管理端核心功能展示"),
        ("07", "系统测试与实验", "功能测试、性能测试、兼容性测试、实验分析"),
        ("08", "结论与展望", "研究成果总结与未来改进方向"),
    ]
    for i, (num, title, desc) in enumerate(toc_items):
        y = Inches(1.3 + i * 0.72)
        card = add_rounded_rect(slide, Inches(1.0), y, Inches(11.3), Inches(0.62), WHITE)
        add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(0.6), Inches(0.5),
                     num, font_size=22, color=ACCENT_COLOR, bold=True)
        add_text_box(slide, Inches(1.9), y + Inches(0.05), Inches(3), Inches(0.5),
                     title, font_size=18, color=TITLE_COLOR, bold=True)
        add_text_box(slide, Inches(5.2), y + Inches(0.08), Inches(6.8), Inches(0.5),
                     desc, font_size=14, color=RGBColor(0x66, 0x66, 0x66))

    # ===== Slide 3: Research Background =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "研究背景——非遗数字化传播的时代需求")
    add_page_number(slide, 3, total_slides)

    left_items = [
        "● 国家级非遗代表性项目1000+，联合国非遗名录数量世界第一",
        "● 短视频化趋势：抖音非遗视频超2.2亿条，累计138亿次点赞",
        "● 资源数字化：动画、游戏、产品设计吸收中国元素需数字化非遗",
        "● 体验化转型：非遗体验从博物馆附属走向独立活动产品",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(6), Inches(3.5), left_items, font_size=15, bold_prefix=False)

    right_items = [
        "推荐能力弱——大多采用热门排序，无个性化推荐",
        "交互方式单一——以图文展示为主，缺少语音、对话交互",
        "场景感知缺失——不能捕捉时间、地点、文化情境信号",
        "冷启动问题严重——用户行为数据稀疏，新用户流失率高",
        "推荐结果不可解释——不给出推荐理由，影响用户信任",
    ]
    add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.4),
                 "现有平台不足", font_size=18, color=ACCENT_COLOR, bold=True)
    add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(4), right_items, font_size=14, bold_prefix=True)

    add_shape(slide, Inches(6.5), Inches(1.3), Inches(0.02), Inches(4.5), LIGHT_ACCENT)

    # ===== Slide 4: Research Significance =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "研究意义")
    add_page_number(slide, 4, total_slides)

    theory_card = add_rounded_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.2), WHITE)
    add_text_box(slide, Inches(0.9), Inches(1.5), Inches(5.2), Inches(0.5),
                 "理论意义", font_size=22, color=ACCENT_COLOR, bold=True)
    theory_items = [
        "● 将CRS对话推荐思想引入非遗文化传播场景，探索SAUR交互模式在文化领域的适用性",
        "● 提出以CRS为引导、知识图谱为增强、AI数字人为载体的智能化服务框架",
        "● 构建三维置信度评估模型（显式偏好+隐式行为+对话语义），实现冷启动到精准推荐的渐进收敛",
        "● 为非遗推荐领域提供可解释的推荐机制设计参考",
    ]
    add_bullet_list(slide, Inches(0.9), Inches(2.1), Inches(5.2), Inches(4), theory_items, font_size=14, bold_prefix=False)

    practice_card = add_rounded_rect(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.2), WHITE)
    add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.2), Inches(0.5),
                 "实践意义", font_size=22, color=ACCENT_COLOR, bold=True)
    practice_items = [
        "● 开发完整的非遗文化传播微信小程序，覆盖浏览、问答、推荐、活动、社区全链路",
        '● AI数字人"黑塔"实现多轮问答、语音播报、CRS状态显示、延伸推荐、行动清单生成',
        "● 本地知识库优先+大模型回退策略，兼顾知识准确性与语义理解能力",
        "● 为非遗数字化传播与个性化服务提供可落地的实践参考",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.2), Inches(4), practice_items, font_size=14, bold_prefix=False)

    # ===== Slide 5: Research Status 1 =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "国内外研究现状——推荐系统与CRS")
    add_page_number(slide, 5, total_slides)

    sections = [
        ("推荐系统研究", [
            "协同过滤：依赖历史行为，冷启动问题突出",
            "内容推荐：特征匹配，但多样性不足",
            "混合推荐：融合多源信号，适合数据稀疏场景",
        ]),
        ("对话推荐系统（CRS）", [
            "Jannach等综述：CRS通过多轮交互获取用户反馈",
            "赵梦媛等：传统推荐三大缺陷——稀疏数据、忽略上下文、假设用户已知偏好",
            "Wang等：知识增强提示学习统一CRS框架",
        ]),
        ("大语言模型与推荐融合", [
            "谢广明等：LLM在推荐中的判别式与生成式范式",
            "Wu等：LLM for Recommendation综述",
            "趋势：从独立推荐模型走向LLM驱动的交互式推荐",
        ]),
    ]
    for i, (title, items) in enumerate(sections):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.3), Inches(3.9), Inches(5.2), WHITE)
        add_text_box(slide, x + Inches(0.2), Inches(1.5), Inches(3.5), Inches(0.5),
                     title, font_size=18, color=ACCENT_COLOR, bold=True)
        add_bullet_list(slide, x + Inches(0.2), Inches(2.2), Inches(3.5), Inches(4),
                        [f"● {it}" for it in items], font_size=13, bold_prefix=False, spacing=Pt(12))

    # ===== Slide 6: Research Status 2 =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "国内外研究现状——AI数字人与现有不足")
    add_page_number(slide, 6, total_slides)

    card1 = add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.5), WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
                 "AI数字人与文化导览", font_size=18, color=ACCENT_COLOR, bold=True)
    ai_items = [
        '● 康宁等"辽小博"：知识图谱+大模型融合微调的AI智慧导览系统',
        "● Carreiro等InHeritage：游戏化+AR的文化遗产传播移动应用",
        "● 趋势：从静态展示走向AI驱动的智能导览与个性化交互",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(1.5), ai_items, font_size=14, bold_prefix=False)

    card2 = add_rounded_rect(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(2.5), WHITE)
    add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.2), Inches(0.5),
                 "知识图谱增强推荐", font_size=18, color=ACCENT_COLOR, bold=True)
    kg_items = [
        "● Guo等综述：知识图谱通过实体关系扩展推荐候选集",
        "● 王敏等：KG+LLM增强推荐系统研究",
        "● 汪天雄等：非遗数字化发展中知识图谱视角的前景",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.2), Inches(1.5), kg_items, font_size=14, bold_prefix=False)

    gap_card = add_rounded_rect(slide, Inches(0.5), Inches(4.2), Inches(12.2), Inches(2.8), RGBColor(0xFD, 0xF0, 0xE8))
    add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11.6), Inches(0.5),
                 "现有研究不足与本课题定位", font_size=18, color=ACCENT_COLOR, bold=True)
    gap_items = [
        "1. 现有非遗平台缺少CRS对话推荐能力，无法在交互中主动发现用户偏好",
        "2. AI导览系统多侧重知识问答，缺少推荐引导与偏好采集的闭环设计",
        "3. 知识图谱在非遗推荐中的应用尚处起步阶段，推荐解释性不足",
        "4. 缺少面向非遗场景的完整系统实践，将CRS、AI数字人、知识图谱三者协同的方案尚未见报道",
        "→ 本课题定位：以CRS为引导、知识图谱为增强、AI数字人为载体，构建非遗文化传播的智能化服务框架",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(11.6), Inches(2), gap_items, font_size=14, bold_prefix=False, spacing=Pt(6))

    # ===== Slide 7: Technologies =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "相关技术与理论")
    add_page_number(slide, 7, total_slides)

    techs = [
        ("前端", "微信小程序\nWXML/WXSS/JS", "原生开发，无需第三方UI库\n微信登录、分享、定位直接调用"),
        ("后端", "FastAPI + SQLAlchemy", "高性能异步Web框架\nORM框架，支持SQLite/MySQL切换"),
        ("数据库", "SQLite", "轻量级关系型数据库\n适合开发和小规模部署"),
        ("AI引擎", "豆包大模型 + 本地知识库", "三级回退：知识库→联网搜索→大模型\nTTS双引擎：豆包TTS + Edge-TTS"),
        ("推荐", "CRS对话推荐 + 知识图谱", "ASK-REC状态机驱动\n规则加权混合推荐"),
    ]
    for i, (cat, name, desc) in enumerate(techs):
        x = Inches(0.4 + i * 2.55)
        card = add_rounded_rect(slide, x, Inches(1.3), Inches(2.35), Inches(5.5), WHITE)
        add_shape(slide, x, Inches(1.3), Inches(2.35), Inches(0.5), ACCENT_COLOR)
        add_text_box(slide, x + Inches(0.1), Inches(1.35), Inches(2.15), Inches(0.4),
                     cat, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(2.0), Inches(2.05), Inches(0.8),
                     name, font_size=15, color=TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(3.0), Inches(2.05), Inches(3.5),
                     desc, font_size=12, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

    # ===== Slide 8: System Analysis - Requirements =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统分析——功能需求分析")
    add_page_number(slide, 8, total_slides)

    roles = [
        ("普通用户", "浏览非遗内容、AI对话问答\n个性化推荐、活动报名\n社区讨论、偏好管理"),
        ("内容管理员", "内容审核发布\n质量评定、数据收集\n基本用户管理"),
        ("系统管理员", "系统运维配置\n运行监控、权限控制"),
    ]
    for i, (role, desc) in enumerate(roles):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.3), Inches(3.9), Inches(1.8), WHITE)
        add_shape(slide, x, Inches(1.3), Inches(3.9), Inches(0.45), ACCENT_COLOR)
        add_text_box(slide, x + Inches(0.2), Inches(1.35), Inches(3.5), Inches(0.35),
                     role, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(1.85), Inches(3.5), Inches(1.1),
                     desc, font_size=13, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

    modules = [
        ("非遗内容展示", "浏览列表、详情\n分类导航、搜索"),
        ("个性化推荐", "首页推荐、冷启动引导\n动态调整、推荐解释"),
        ("AI数字人对话", "知识问答、多轮会话\nASK追问、推荐触发"),
        ("活动报名", "活动浏览、详情\n在线报名、记录管理"),
        ("社区互动", "帖子浏览、发帖\n评论、点赞、话题标签"),
        ("用户中心", "微信登录、偏好设置\n收藏管理、行为信息"),
    ]
    for i, (mod, desc) in enumerate(modules):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(3.5 + row * 1.8)
        card = add_rounded_rect(slide, x, y, Inches(3.9), Inches(1.5), WHITE)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.5), Inches(0.4),
                     mod, font_size=15, color=ACCENT_COLOR, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.5), Inches(0.9),
                     desc, font_size=12, color=BODY_COLOR)

    # ===== Slide 9: System Analysis - Feasibility =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统分析——可行性分析")
    add_page_number(slide, 9, total_slides)

    feasibilities = [
        ("技术可行性", "✓ 微信小程序成熟框架，原生开发\n✓ FastAPI高性能异步框架+SQLAlchemy ORM\n✓ 本地知识库-联网搜索-大模型三级回退\n✓ 规则加权混合推荐，适合数据稀疏场景"),
        ("经济可行性", "✓ SQLite本地数据库，零成本\n✓ 云服务器月费用百元以内\n✓ 豆包大模型提供免费额度\n✓ 微信小程序注册发布平价"),
        ("操作可行性", "✓ 小程序免安装，扫码即用\n✓ AI数字人自然语言交互，零学习成本\n✓ 管理端Web界面，表单图表操作直观"),
        ("社会与法律可行性", "✓ 本地知识库降低AI幻觉\n✓ 符合《个人信息保护法》\n✓ 传播非遗文化有积极社会意义\n✓ 符合《非物质文化遗产法》第三条"),
    ]
    for i, (title, desc) in enumerate(feasibilities):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(1.3 + row * 2.9)
        card = add_rounded_rect(slide, x, y, Inches(6), Inches(2.5), WHITE)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.4), Inches(0.5),
                     title, font_size=18, color=ACCENT_COLOR, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.65), Inches(5.4), Inches(1.7),
                     desc, font_size=14, color=BODY_COLOR)

    # ===== Slide 10: System Architecture =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统设计——整体架构")
    add_page_number(slide, 10, total_slides)

    layers = [
        ("前端交互层", "微信小程序用户端 | Web管理端", RGBColor(0xE8, 0xD5, 0xC4)),
        ("接口服务层", "AI问答接口 | CRS状态接口 | 用户画像接口 | 内容与活动接口 | 讨论区接口 | 登录接口", RGBColor(0xD4, 0xC0, 0xAA)),
        ("核心能力层", "AI服务 | 推荐服务 | 本地知识检索 | 知识图谱增强 | 联网补充 | TTS语音播报", RGBColor(0xC4, 0x1E, 0x3A)),
        ("数据与外部服务层", "SQLite数据库 | 本地知识数据 | 推荐日志 | CRS会话数据 | 豆包API | 联网搜索 | 微信登录", RGBColor(0x8B, 0x22, 0x22)),
    ]
    for i, (name, content, color) in enumerate(layers):
        y = Inches(1.3 + i * 1.45)
        card = add_rounded_rect(slide, Inches(1.5), y, Inches(10.3), Inches(1.2), color)
        text_color = WHITE if i >= 2 else DARK_BG
        add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(2.5), Inches(0.4),
                     name, font_size=18, color=text_color, bold=True)
        add_text_box(slide, Inches(1.8), y + Inches(0.55), Inches(9.7), Inches(0.5),
                     content, font_size=13, color=text_color)

    add_text_box(slide, Inches(0.3), Inches(1.3), Inches(1.1), Inches(0.4),
                 "设计原则：", font_size=14, color=ACCENT_COLOR, bold=True)
    add_text_box(slide, Inches(0.3), Inches(1.7), Inches(1.1), Inches(3),
                 "前后端分离\n主线趋同\n轻量优先\n可扩展", font_size=12, color=BODY_COLOR)

    # ===== Slide 11: CRS Recommendation Design =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统设计——CRS对话推荐机制")
    add_page_number(slide, 11, total_slides)

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
                 "核心问题：现在应该继续提问，还是应该直接回答？——效率与精度的平衡",
                 font_size=16, color=SUBTITLE_COLOR, bold=True)

    left_card = add_rounded_rect(slide, Inches(0.5), Inches(1.9), Inches(5.8), Inches(5), WHITE)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.2), Inches(0.4),
                 "三维置信度评估模型", font_size=18, color=ACCENT_COLOR, bold=True)
    conf_items = [
        "C = 0.40×Se + 0.35×Si + 0.25×Sd",
        "",
        "Se（显式偏好）：用户偏好设置、ASK选项回答",
        "Si（隐式行为）：曝光、点击、浏览、活动报名、讨论互动",
        "Sd（对话语义）：AI提问中的非遗关键词、地区特征、场景偏好",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(5.2), Inches(4), conf_items, font_size=14, bold_prefix=False)

    right_card = add_rounded_rect(slide, Inches(6.9), Inches(1.9), Inches(5.8), Inches(5), WHITE)
    add_text_box(slide, Inches(7.2), Inches(2.0), Inches(5.2), Inches(0.4),
                 "ASK-REC状态机", font_size=18, color=ACCENT_COLOR, bold=True)

    modes = [
        ("cold_start", "C < 28", "优先ASK，采集高价值偏好信息", RGBColor(0x66, 0x99, 0xCC)),
        ("mixed", "28 ≤ C < 62", "边推荐边追问，补充画像", RGBColor(0xDD, 0xAA, 0x44)),
        ("precision", "C ≥ 62", "以REC为主，精准推荐", RGBColor(0x66, 0xBB, 0x66)),
    ]
    for i, (mode, threshold, desc, color) in enumerate(modes):
        y = Inches(2.6 + i * 1.3)
        badge = add_rounded_rect(slide, Inches(7.2), y, Inches(1.8), Inches(0.4), color)
        add_text_box(slide, Inches(7.2), y + Inches(0.02), Inches(1.8), Inches(0.35),
                     mode, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(9.2), y, Inches(1.5), Inches(0.35),
                     threshold, font_size=13, color=BODY_COLOR, bold=True)
        add_text_box(slide, Inches(7.2), y + Inches(0.45), Inches(5.2), Inches(0.6),
                     desc, font_size=13, color=BODY_COLOR)

    add_text_box(slide, Inches(7.2), Inches(5.8), Inches(5.2), Inches(0.8),
                 "特殊策略：\n意图驱动推荐（用户主动要求）\n恢复式提问（多轮无有效偏好时触发）",
                 font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # ===== Slide 12: AI Digital Human Design =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, '系统设计——AI数字人"黑塔"')
    add_page_number(slide, 12, total_slides)

    duties = [
        ("知识传授", "围绕非遗问题解释说明\n承担导览入口作用"),
        ("推荐引导", "将提问自然引向\n内容/活动/讨论推荐"),
        ("采集偏好", "多轮对话+ASK模板\n收集类别/地区/体验偏好"),
        ("闭环触发", "行为反馈输入画像\nCRS状态持续深化"),
    ]
    for i, (duty, desc) in enumerate(duties):
        x = Inches(0.5 + i * 3.15)
        card = add_rounded_rect(slide, x, Inches(1.3), Inches(2.9), Inches(1.8), WHITE)
        add_shape(slide, x, Inches(1.3), Inches(2.9), Inches(0.4), ACCENT_COLOR)
        add_text_box(slide, x + Inches(0.1), Inches(1.33), Inches(2.7), Inches(0.35),
                     duty, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(1.8), Inches(2.6), Inches(1.1),
                     desc, font_size=13, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

    chain_card = add_rounded_rect(slide, Inches(0.5), Inches(3.4), Inches(12.2), Inches(1.6), WHITE)
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(3), Inches(0.4),
                 "AI回答链路（本地知识优先）", font_size=16, color=ACCENT_COLOR, bold=True)
    chain_steps = "接收问题+上下文 → 检索本地知识库 → 读取CRS会话状态 → 生成推荐载荷 → ASK-REC决策 → 调用豆包API生成/润色 → 返回回答+推荐卡+ASK选项+策略说明"
    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11.6), Inches(0.8),
                 chain_steps, font_size=14, color=BODY_COLOR)

    prompt_card = add_rounded_rect(slide, Inches(0.5), Inches(5.3), Inches(12.2), Inches(1.8), WHITE)
    add_text_box(slide, Inches(0.8), Inches(5.4), Inches(3), Inches(0.4),
                 "三层提示词约束", font_size=16, color=ACCENT_COLOR, bold=True)
    prompt_items = [
        "第一层：全局角色提示——确定黑塔身份、语气、服务范围和知识边界",
        "第二层：任务型提示——知识命中时润色简化，未命中时限定输出范围",
        "第三层：CRS感知提示——根据cold_start/mixed/precision动态调整语气深度",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(5.9), Inches(11.6), Inches(1), prompt_items, font_size=13, bold_prefix=False, spacing=Pt(4))

    # ===== Slide 13: Knowledge Graph Design =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统设计——知识图谱与内容治理")
    add_page_number(slide, 13, total_slides)

    kg_card = add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
                 "知识图谱设计", font_size=20, color=ACCENT_COLOR, bold=True)
    kg_items = [
        "定位：推荐和AI对话的加强器，非主排序器",
        "",
        "图谱增强方式：",
        "● 扩大推荐候选人范围——相近实体、类别关系、路径关系",
        "● 加强推荐理由解释——路径关系→自然语言推荐理由",
        "● 以用户已知实体为起点，向丰富但不脱离主题方向延伸",
        "",
        "CRS与图谱联动：",
        "● 识别用户提及的具体实体→映射到画像特征",
        "● 推荐卡附加图谱理由和路径说明",
        "● 一石二鸟：增强解释+参与兴趣建模",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.5), kg_items, font_size=13, bold_prefix=False, spacing=Pt(3))

    gov_card = add_rounded_rect(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
    add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.2), Inches(0.5),
                 "内容治理体系", font_size=20, color=ACCENT_COLOR, bold=True)
    gov_items = [
        "质量评分机制（三维度量化）：",
        "● 内容完整度：正文字数、封面图、摘要",
        "● 信息丰富度：章节归属、标签覆盖",
        "● 原创质量：来源可信度、内容哈希去重",
        "",
        "审核发布机制：",
        "● 质量分≥阈值 → 自动进入推荐池",
        "● 人工审核 → 精选标记/退回修改",
        "",
        "白名单回补机制：",
        "● 定期从达标内容中重新选拔精选",
        "● 保证推荐池的新鲜度和质量",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.2), Inches(4.5), gov_items, font_size=13, bold_prefix=False, spacing=Pt(3))

    # ===== Slide 14: Implementation - Homepage =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统实现——首页与个性化推荐")
    add_page_number(slide, 14, total_slides)

    features = [
        ("AI导览横幅", "● 数字导览中枢徽章\n● 动态标题和描述文案\n● CRS模式进度胶囊\n  （想认识你→正在了解你→已懂你）\n● 黑塔数字人组件\n  （表情随CRS模式变化）\n● 立即开始按钮"),
        ("精选推荐区", "● 第一条推荐大卡片展示\n● 今日推荐区：内容/活动/讨论\n● 推荐理由标签+新鲜度标签\n● 点击行为即时回流画像"),
        ("快速入口+TabBar", "● 非遗发展史/非遗地点快捷入口\n● 自定义TabBar：\n  主页/文化/活动/讨论/我的\n● 胶囊圆角设计\n● 选中态朱红渐变背景"),
    ]
    for i, (title, desc) in enumerate(features):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.3), Inches(3.9), Inches(5.5), WHITE)
        add_shape(slide, x, Inches(1.3), Inches(3.9), Inches(0.45), ACCENT_COLOR)
        add_text_box(slide, x + Inches(0.2), Inches(1.35), Inches(3.5), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(3.5), Inches(4.7),
                     desc, font_size=13, color=BODY_COLOR)

    # ===== Slide 15: Implementation - Content & Activity =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统实现——非遗内容与活动报名")
    add_page_number(slide, 15, total_slides)

    content_card = add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
                 "非遗内容模块", font_size=20, color=ACCENT_COLOR, bold=True)
    content_items = [
        "内容列表页：",
        "● 瀑布流双列布局",
        "● 策展精选显眼位置",
        "● 按子章节和类型打标签",
        "● 个性化内容推荐混合展示",
        "",
        "内容详情页：",
        "● 封面图展示",
        "● 三分钟看点摘要",
        "● 延伸讲解内容",
        "● AI浮窗入口（场景化问答）",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.5), content_items, font_size=14, bold_prefix=False, spacing=Pt(4))

    activity_card = add_rounded_rect(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
    add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.2), Inches(0.5),
                 "活动报名模块", font_size=20, color=ACCENT_COLOR, bold=True)
    activity_items = [
        "活动列表页：",
        "● 按月份筛选",
        '● 前两条标记"本月主推"',
        "● 个性化活动推荐混合展示",
        "",
        "活动详情页：",
        "● 封面、时间地点、组织者",
        "● 展示块（亮点/议程/提示）",
        "● 在线报名/取消报名",
        "● 报名备注+报名须知确认",
        "● AI浮窗入口",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.2), Inches(4.5), activity_items, font_size=14, bold_prefix=False, spacing=Pt(4))

    # ===== Slide 16: Implementation - Community =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统实现——社区讨论与用户中心")
    add_page_number(slide, 16, total_slides)

    community_card = add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
                 "社区讨论模块", font_size=20, color=ACCENT_COLOR, bold=True)
    community_items = [
        "讨论列表页：",
        "● 关键词搜索、标签筛选（10个标签）",
        "● 排序切换（热门/最新）、只看收藏",
        "● 热门话题TOP3展示",
        "",
        "发帖功能：",
        "● 三种模板：提问型/体验分享型/活动反馈型",
        "● 可选标签、上传封面图",
        "",
        "话题详情页：",
        "● 帖子正文+评论列表",
        "● 发表评论、点赞、收藏",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.5), community_items, font_size=14, bold_prefix=False, spacing=Pt(4))

    user_card = add_rounded_rect(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
    add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.2), Inches(0.5),
                 "用户中心模块", font_size=20, color=ACCENT_COLOR, bold=True)
    user_items = [
        "用户画像卡片：",
        "● 偏好关键词展示",
        "● 活跃场景展示",
        "● 关注地区展示",
        "",
        "统计卡片：",
        "● 发帖数、活动记录、已选偏好",
        "",
        "偏好设置页：",
        "● 非遗类别（工艺/戏曲/民俗/医药）",
        "● 场景类型（知识阅读/活动体验/论坛交流）",
        "● 地区（华东/华南/西南/华北/西北/东北）",
        "● 保存后作为显式信号纳入画像",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.2), Inches(4.5), user_items, font_size=14, bold_prefix=False, spacing=Pt(4))

    # ===== Slide 17: Implementation - AI Digital Human =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, '系统实现——AI数字人"黑塔"核心交互')
    add_page_number(slide, 17, total_slides)

    mode_details = [
        ("cold_start\n冷启动模式", "C < 28", [
            "主动提问为主",
            "ASK卡片收集偏好",
            "类目→地区→场景→程度",
            "不直接给出推荐结果",
            "数字人好奇张望表情",
        ], RGBColor(0x66, 0x99, 0xCC)),
        ("mixed\n混合模式", "28 ≤ C < 62", [
            "边推荐边追问",
            "1-2张AI推荐卡+B组追问",
            "上方推荐+下方ASK选项",
            "推荐理由贴合当前选择",
            "数字人托腮思考表情",
        ], RGBColor(0xDD, 0xAA, 0x44)),
        ("precision\n精准模式", "C ≥ 62", [
            "停止追问，直接输出",
            "高度匹配推荐结果",
            "精准匹配偏好表达",
            "推荐卡占主要区域",
            "数字人自信微笑表情",
        ], RGBColor(0x66, 0xBB, 0x66)),
    ]
    for i, (mode, threshold, items, color) in enumerate(mode_details):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.3), Inches(3.9), Inches(5.5), WHITE)
        badge = add_rounded_rect(slide, x + Inches(0.5), Inches(1.5), Inches(2.9), Inches(0.9), color)
        add_text_box(slide, x + Inches(0.5), Inches(1.5), Inches(2.9), Inches(0.9),
                     mode, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.5), Inches(2.5), Inches(2.9), Inches(0.35),
                     threshold, font_size=13, color=BODY_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + Inches(0.3), Inches(3.0), Inches(3.3), Inches(3.5),
                        [f"● {it}" for it in items], font_size=13, bold_prefix=False, spacing=Pt(6))

    # ===== Slide 18: Implementation - AI Features =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统实现——AI对话功能与信息层次")
    add_page_number(slide, 18, total_slides)

    feature_cards = [
        ("语音播报", "豆包TTS优先\nEdge-TTS降级\n双引擎架构"),
        ("延伸推荐", "回答下方展示\n相关内容/活动/讨论\n推荐卡片即时刷新"),
        ("行动清单", "浏览、报名、讨论\n下一步行动建议\n闭环触发机制"),
        ("AI浮窗", "内容页/活动页嵌入\n场景化问答\n随时发起对话"),
    ]
    for i, (title, desc) in enumerate(feature_cards):
        x = Inches(0.5 + i * 3.15)
        card = add_rounded_rect(slide, x, Inches(1.3), Inches(2.9), Inches(2.2), WHITE)
        add_shape(slide, x, Inches(1.3), Inches(2.9), Inches(0.4), ACCENT_COLOR)
        add_text_box(slide, x + Inches(0.1), Inches(1.33), Inches(2.7), Inches(0.35),
                     title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(1.8), Inches(2.6), Inches(1.5),
                     desc, font_size=13, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

    info_card = add_rounded_rect(slide, Inches(0.5), Inches(3.8), Inches(12.2), Inches(3.2), WHITE)
    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(5), Inches(0.5),
                 "AI对话页三层信息结构", font_size=18, color=ACCENT_COLOR, bold=True)

    info_layers = [
        ("第一层：任务主线层", "黑塔数字人形象 | 对话消息区 | 输入框 | 当前模式提示 | ASK提问卡片"),
        ("第二层：推荐结果层", "AI推荐卡 | 推荐理由 | 下一步建议（浏览/报名/讨论）"),
        ("第三层：解释层", "置信度详细部分 | 策略解释 | 图谱依据"),
    ]
    for i, (layer, desc) in enumerate(info_layers):
        y = Inches(4.5 + i * 0.8)
        add_text_box(slide, Inches(0.8), y, Inches(2.8), Inches(0.35),
                     layer, font_size=14, color=ACCENT_COLOR, bold=True)
        add_text_box(slide, Inches(3.8), y, Inches(8.5), Inches(0.6),
                     desc, font_size=13, color=BODY_COLOR)

    # ===== Slide 19: Implementation - Admin =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统实现——Web管理端")
    add_page_number(slide, 19, total_slides)

    admin_modules = [
        ("内容管理", "卡片网格展示\n封面/标题/状态/精选/质量分\n关键词搜索+状态筛选"),
        ("质量检查", "四项统计指标\n五档质量分布柱状图\n质量诊断弹窗+改进建议\n重建推荐池"),
        ("活动管理", "卡片形式展示\n封面/标题/地点/时间/人数\n上架下架管理"),
        ("帖子管理", "帖子列表展示\n封面/标题/作者/点赞/评论\n编辑/精选/删除"),
        ("用户管理", "表格形式展示\n昵称/OpenID/手机/角色/状态\n行为统计查看"),
        ("统计导出", "四项平台指标\n行为日志筛选\nCSV导出功能"),
        ("知识库管理", "表格形式展示\n问题/章节/关键词/状态\n新增/编辑/启禁用"),
    ]
    for i, (title, desc) in enumerate(admin_modules):
        col = i % 4
        row = i // 4
        x = Inches(0.4 + col * 3.2)
        y = Inches(1.3 + row * 2.9)
        card = add_rounded_rect(slide, x, y, Inches(3), Inches(2.5), WHITE)
        add_shape(slide, x, y, Inches(3), Inches(0.4), ACCENT_COLOR)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.03), Inches(2.8), Inches(0.35),
                     title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.7), Inches(1.8),
                     desc, font_size=12, color=BODY_COLOR)

    # ===== Slide 20: Testing =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "系统测试")
    add_page_number(slide, 20, total_slides)

    test_sections = [
        ("功能测试", [
            "黑盒测试：10条用例，通过10条",
            "  覆盖认证/内容/活动/讨论/推荐/AI/图谱/管理8模块",
            "白盒测试：7条用例，通过7条",
            "  CRS决策引擎/ASK模板/知识图谱/数据模型",
        ]),
        ("性能测试", [
            "API平均响应时间：43.16ms",
            "最大响应时间：51.33ms（均<100ms）",
            "并发测试：10用户50请求",
            "  成功率100%，吞吐225.5请求/秒",
            "AI对话平均响应：7.9秒",
            "  P50=2.6s，P90=4.3s",
        ]),
        ("兼容性测试", [
            "浏览器：Chrome/Firefox/Safari/Edge ✓",
            "微信基础库：≥2.20.0 ✓",
            "设备：iOS + Android ✓",
            "数据库：SQLite/MySQL/PostgreSQL ✓",
        ]),
    ]
    for i, (title, items) in enumerate(test_sections):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.3), Inches(3.9), Inches(5.5), WHITE)
        add_shape(slide, x, Inches(1.3), Inches(3.9), Inches(0.45), ACCENT_COLOR)
        add_text_box(slide, x + Inches(0.2), Inches(1.35), Inches(3.5), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + Inches(0.2), Inches(1.9), Inches(3.5), Inches(4.5),
                        items, font_size=13, bold_prefix=False, spacing=Pt(5))

    # ===== Slide 21: Experiment =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "实验与分析")
    add_page_number(slide, 21, total_slides)

    exp_card = add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
                 "AI对话质量评估", font_size=20, color=ACCENT_COLOR, bold=True)

    exp_items = [
        "知识库命中率：90.0%",
        "● 20个测试问题，15个预期命中全部命中",
        "● 5个冷门问题中3个未命中→回退豆包直答/联网搜索",
        "● 命中项置信度集中在0.80-0.85区间",
        "",
        "AI回答响应时间：",
        "● P50 = 2.6秒，P90 = 4.3秒",
        "● 流式输出机制下3秒内可见首个回答片段",
        "",
        "回答来源分布：",
        "● KB命中+润色 占比最大（本地知识优先链路有效）",
        "● 豆包直答+豆包组合 占比较大（大模型能力体现）",
        "● 联网搜索+兜底回答 属补充手段",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.5), exp_items, font_size=13, bold_prefix=False, spacing=Pt(3))

    eval_card = add_rounded_rect(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
    add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.2), Inches(0.5),
                 "推荐评价指标", font_size=20, color=ACCENT_COLOR, bold=True)

    metrics = [
        ("Precision@5", "推荐列表前5项中与用户偏好相关的比例\n→ 反映推荐准确性"),
        ("Diversity@5", "推荐列表覆盖的非遗类别数/列表长度\n→ 反映推荐多样性"),
        ("Coverage", "推荐结果覆盖内容占全部已发布内容的比例\n→ 反映推荐覆盖广度"),
        ("NDCG@5", "考虑位置权重的排序质量\n→ 反映推荐排序质量"),
    ]
    for i, (metric, desc) in enumerate(metrics):
        y = Inches(2.2 + i * 1.15)
        badge = add_rounded_rect(slide, Inches(7.2), y, Inches(2.2), Inches(0.35), ACCENT_COLOR)
        add_text_box(slide, Inches(7.2), y + Inches(0.02), Inches(2.2), Inches(0.3),
                     metric, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(9.6), y, Inches(3), Inches(0.9),
                     desc, font_size=12, color=BODY_COLOR)

    # ===== Slide 22: Conclusion =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_decorated_title(slide, "结论与展望")
    add_page_number(slide, 22, total_slides)

    conclusion_card = add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(7.5), Inches(5.5), WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(0.5),
                 "研究结论", font_size=20, color=ACCENT_COLOR, bold=True)
    conclusion_items = [
        "● 构建了基于CRS推荐与AI数字人的非遗文化传播微信小程序",
        '● AI数字人"黑塔"通过ASK追问逐步构建用户偏好画像',
        "● 从冷启动到精准推荐只需3-4轮交互即可收敛",
        "● 本地知识优先+大模型回退策略兼顾知识准确性与语义理解",
        "● ASK-REC决策引擎执行三阶段置信度收敛机制",
        "● 知识图谱实现实体偏好映射和推荐解释生成",
        "● 系统性能满足中小规模使用场景",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(7), Inches(4.5), conclusion_items, font_size=14, bold_prefix=False, spacing=Pt(6))

    future_card = add_rounded_rect(slide, Inches(8.5), Inches(1.3), Inches(4.3), Inches(5.5), RGBColor(0xFD, 0xF0, 0xE8))
    add_text_box(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(0.5),
                 "未来展望", font_size=20, color=ACCENT_COLOR, bold=True)
    future_items = [
        "◆ 加入更细致的知识图谱推理能力，提高个性化与智能化程度",
        "◆ 同线下非遗场馆深度对接，达成预约体验、位置推荐",
        "◆ 支持多模态交互，使AI数字人引导更具拟人性",
        "◆ 伴随大语言模型技术发展，对话方式更加自然",
    ]
    add_bullet_list(slide, Inches(8.8), Inches(2.1), Inches(3.7), Inches(4.5), future_items, font_size=13, bold_prefix=False, spacing=Pt(10))

    # ===== Slide 23: Thanks =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ACCENT_COLOR)
    add_shape(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), ACCENT_COLOR)

    add_shape(slide, Inches(3), Inches(2.2), Inches(7.3), Inches(0.02), GOLD)
    add_text_box(slide, Inches(2), Inches(2.5), Inches(9.3), Inches(1.2),
                 "感谢各位老师的指导与评审", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(3), Inches(3.8), Inches(7.3), Inches(0.02), GOLD)

    add_text_box(slide, Inches(2), Inches(4.3), Inches(9.3), Inches(0.6),
                 "恳请各位老师批评指正", font_size=22, color=GOLD, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.5),
                 "答辩人：王子轩    指导教师：徐龙琴 教授  高静 高级工程师",
                 font_size=16, color=RGBColor(0xBB, 0xAA, 0x99), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(6.0), Inches(9.3), Inches(0.5),
                 "仲恺农业工程学院  信息科学与技术学院  2026年5月9日",
                 font_size=14, color=RGBColor(0x99, 0x88, 0x77), alignment=PP_ALIGN.CENTER)

    output_path = r"d:\桌面\毕业设计\答辩演示文件夹\答辩PPT_基于CRS推荐与AI数字人的非遗文化传播系统.pptx"
    prs.save(output_path)
    print(f"PPT生成完成: {output_path}")
    print(f"文件大小: {os.path.getsize(output_path)} 字节")
    print(f"总页数: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
