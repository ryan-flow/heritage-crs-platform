from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

INK_BLACK = RGBColor(0x1A, 0x14, 0x10)
CINNABAR = RGBColor(0xC4, 0x1E, 0x3A)
GOLD_LEAF = RGBColor(0xD4, 0xA5, 0x37)
WARM_CREAM = RGBColor(0xF7, 0xF2, 0xE8)
CHARCOAL = RGBColor(0x2C, 0x24, 0x1E)
SILVER_MIST = RGBColor(0x9A, 0x8C, 0x7E)
DEEP_RED = RGBColor(0x8B, 0x1A, 0x2B)
JADE_GREEN = RGBColor(0x4A, 0x7C, 0x59)
INK_BLUE = RGBColor(0x2A, 0x3D, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GOLD = RGBColor(0xF0, 0xE4, 0xC8)
DARK_OVERLAY = RGBColor(0x0D, 0x0A, 0x08)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None, rotation=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if rotation:
        shape.rotation = rotation
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_diamond(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=INK_BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Source Han Sans SC"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height, runs_data, alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing:
        p.line_spacing = line_spacing
    for i, rd in enumerate(runs_data):
        run = p.add_run()
        run.text = rd.get("text", "")
        run.font.size = Pt(rd.get("size", 16))
        run.font.color.rgb = rd.get("color", INK_BLACK)
        run.font.bold = rd.get("bold", False)
        run.font.name = rd.get("font", "Source Han Sans SC")
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=INK_BLACK, spacing=Pt(4), bullet_color=CINNABAR):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        if isinstance(item, tuple):
            label, desc = item
            r1 = p.add_run()
            r1.text = label
            r1.font.size = Pt(font_size)
            r1.font.color.rgb = CINNABAR
            r1.font.bold = True
            r1.font.name = "Source Han Sans SC"
            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = color
            r2.font.bold = False
            r2.font.name = "Source Han Sans SC"
        else:
            r1 = p.add_run()
            r1.text = "\u25C6 "
            r1.font.size = Pt(font_size - 2)
            r1.font.color.rgb = bullet_color
            r1.font.bold = False
            r1.font.name = "Source Han Sans SC"
            r2 = p.add_run()
            r2.text = item
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = color
            r2.font.bold = False
            r2.font.name = "Source Han Sans SC"
    return txBox


def add_section_header(slide, number, title, subtitle="", top=Inches(0.35)):
    add_shape(slide, Inches(0.5), top, Inches(0.06), Inches(0.7), CINNABAR)
    add_text_box(slide, Inches(0.75), top - Inches(0.05), Inches(1), Inches(0.4),
                 number, font_size=12, color=CINNABAR, bold=True, font_name="Source Han Serif SC")
    add_text_box(slide, Inches(0.75), top + Inches(0.25), Inches(8), Inches(0.5),
                 title, font_size=26, color=INK_BLACK, bold=True, font_name="Source Han Serif SC")
    if subtitle:
        add_text_box(slide, Inches(0.75), top + Inches(0.7), Inches(10), Inches(0.35),
                     subtitle, font_size=13, color=SILVER_MIST)
    add_shape(slide, Inches(0.5), top + Inches(1.0), Inches(12.3), Inches(0.008), LIGHT_GOLD)


def add_page_number(slide, num, total):
    add_text_box(slide, Inches(12.0), Inches(7.1), Inches(1), Inches(0.3),
                 f"{num} / {total}", font_size=9, color=SILVER_MIST, alignment=PP_ALIGN.RIGHT, font_name="Source Han Serif SC")


def add_decorative_corner(slide, position="top-right"):
    if position == "top-right":
        add_shape(slide, Inches(12.5), Inches(0), Inches(0.833), Inches(0.04), CINNABAR)
        add_shape(slide, Inches(13.293), Inches(0), Inches(0.04), Inches(0.8), CINNABAR)
    elif position == "bottom-left":
        add_shape(slide, Inches(0), Inches(7.46), Inches(0.833), Inches(0.04), CINNABAR)
        add_shape(slide, Inches(0), Inches(6.7), Inches(0.04), Inches(0.8), CINNABAR)


def add_ink_wash_accent(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE0, 0xD4)
    shape.line.fill.background()
    shape.rotation = 15


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]
    total_slides = 22

    # ===== Slide 1: Cover =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, INK_BLACK)

    add_shape(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), CINNABAR)
    add_shape(slide, Inches(0.12), Inches(0), Inches(0.04), Inches(7.5), GOLD_LEAF)

    add_shape(slide, Inches(1.2), Inches(1.5), Inches(0.04), Inches(4.5), GOLD_LEAF)

    add_text_box(slide, Inches(1.8), Inches(1.5), Inches(9), Inches(0.5),
                 "毕业设计答辩", font_size=16, color=GOLD_LEAF, bold=False, font_name="Source Han Serif SC")
    add_text_box(slide, Inches(1.8), Inches(2.1), Inches(9), Inches(1.5),
                 "基于CRS推荐与AI数字人的\n非遗文化传播系统实现",
                 font_size=38, color=WHITE, bold=True, font_name="Source Han Serif SC")

    add_shape(slide, Inches(1.8), Inches(3.9), Inches(3), Inches(0.02), GOLD_LEAF)

    info_lines = [
        ("答辩人  ", "王子轩"),
        ("学    号  ", "202210274225"),
        ("专    业  ", "数据科学与大数据技术222班"),
        ("学    院  ", "信息科学与技术学院"),
        ("指导教师  ", "徐龙琴 教授  高静 高级工程师"),
    ]
    for i, (label, value) in enumerate(info_lines):
        y = Inches(4.2 + i * 0.42)
        add_rich_text_box(slide, Inches(1.8), y, Inches(8), Inches(0.4),
                          [{"text": label, "size": 14, "color": SILVER_MIST, "bold": False},
                           {"text": value, "size": 14, "color": WHITE, "bold": True}])

    add_text_box(slide, Inches(1.8), Inches(6.5), Inches(4), Inches(0.4),
                 "仲恺农业工程学院", font_size=18, color=GOLD_LEAF, bold=True, font_name="Source Han Serif SC")
    add_text_box(slide, Inches(1.8), Inches(6.9), Inches(4), Inches(0.3),
                 "2026.05.09", font_size=13, color=SILVER_MIST, font_name="Source Han Serif SC")

    add_diamond(slide, Inches(11.5), Inches(1.0), Inches(0.3), CINNABAR)
    add_diamond(slide, Inches(10.2), Inches(5.8), Inches(0.2), GOLD_LEAF)
    add_diamond(slide, Inches(12.0), Inches(4.5), Inches(0.15), RGBColor(0x55, 0x44, 0x33))

    add_ink_wash_accent(slide, Inches(9.5), Inches(2.0), Inches(3.5), Inches(3.5))

    # ===== Slide 2: TOC =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_decorative_corner(slide, "top-right")

    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(4), Inches(0.6),
                 "CONTENTS", font_size=11, color=CINNABAR, bold=True, font_name="Source Han Serif SC")
    add_text_box(slide, Inches(0.5), Inches(0.75), Inches(6), Inches(0.6),
                 "汇报提纲", font_size=30, color=INK_BLACK, bold=True, font_name="Source Han Serif SC")
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(2), Inches(0.02), CINNABAR)
    add_page_number(slide, 2, total_slides)

    toc_items = [
        ("01", "研究背景与意义", "非遗数字化传播的时代需求"),
        ("02", "国内外研究现状", "CRS、AI数字人研究进展与不足"),
        ("03", "相关技术与理论", "FastAPI、LLM、CRS核心理论"),
        ("04", "系统分析", "需求分析、可行性、数据流程"),
        ("05", "系统设计", "架构、推荐模块、AI数字人、知识图谱"),
        ("06", "系统实现", "小程序端与管理端功能展示"),
        ("07", "系统测试与实验", "功能/性能/兼容性测试与实验分析"),
        ("08", "结论与展望", "研究成果总结与未来方向"),
    ]
    for i, (num, title, desc) in enumerate(toc_items):
        y = Inches(1.6 + i * 0.7)
        add_shape(slide, Inches(0.8), y, Inches(0.5), Inches(0.5), CINNABAR if i % 2 == 0 else CHARCOAL)
        add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.4),
                     num, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")
        add_text_box(slide, Inches(1.5), y + Inches(0.02), Inches(4), Inches(0.35),
                     title, font_size=17, color=INK_BLACK, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, Inches(1.5), y + Inches(0.32), Inches(5), Inches(0.25),
                     desc, font_size=11, color=SILVER_MIST)

    add_ink_wash_accent(slide, Inches(8.5), Inches(1.0), Inches(5), Inches(5))

    # ===== Slide 3: Research Background =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "01", "研究背景", "非遗数字化传播的时代需求与现有平台不足")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 3, total_slides)

    left_card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5.3), WHITE, LIGHT_GOLD)
    add_text_box(slide, Inches(0.8), Inches(1.65), Inches(5.4), Inches(0.4),
                 "时代需求", font_size=18, color=CINNABAR, bold=True, font_name="Source Han Serif SC")
    bg_items = [
        "国家级非遗代表性项目1000+，联合国非遗名录数量世界第一",
        "短视频化：抖音非遗视频超2.2亿条，138亿次点赞",
        "资源数字化：动画、游戏、产品设计需数字化非遗素材",
        "体验化转型：非遗从博物馆附属走向独立活动产品",
        "传播重心从专业领域转向大众视野",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.4), Inches(4), bg_items, font_size=13, spacing=Pt(8))

    right_card = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(5.3), INK_BLACK)
    add_text_box(slide, Inches(7.2), Inches(1.65), Inches(5.4), Inches(0.4),
                 "现有平台五大不足", font_size=18, color=GOLD_LEAF, bold=True, font_name="Source Han Serif SC")
    gaps = [
        ("推荐能力弱  ", "热门排序为主，无个性化推荐"),
        ("交互方式单一  ", "图文展示为主，缺少语音对话交互"),
        ("场景感知缺失  ", "不能捕捉时间地点文化情境信号"),
        ("冷启动严重  ", "用户行为稀疏，新用户流失率高"),
        ("推荐不可解释  ", "不给出推荐理由，影响信任度"),
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.4), Inches(4), gaps, font_size=13, color=WHITE, spacing=Pt(10), bullet_color=GOLD_LEAF)

    # ===== Slide 4: Research Significance =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "01", "研究意义", "理论创新与实践价值")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 4, total_slides)

    theory_card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5), CINNABAR)
    add_text_box(slide, Inches(0.8), Inches(1.55), Inches(5.2), Inches(0.4),
                 "理论意义", font_size=18, color=WHITE, bold=True, font_name="Source Han Serif SC")
    theory_items = [
        "将CRS对话推荐思想引入非遗文化传播场景，探索SAUR交互模式在文化领域的适用性",
        "提出以CRS为引导、知识图谱为增强、AI数字人为载体的智能化服务框架",
        "构建三维置信度评估模型，实现冷启动到精准推荐的渐进收敛",
        "为非遗推荐领域提供可解释的推荐机制设计参考",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.2), Inches(4), theory_items, font_size=13, spacing=Pt(10))

    practice_card = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.5), JADE_GREEN)
    add_text_box(slide, Inches(7.2), Inches(1.55), Inches(5.2), Inches(0.4),
                 "实践意义", font_size=18, color=WHITE, bold=True, font_name="Source Han Serif SC")
    practice_items = [
        '开发完整的非遗文化传播微信小程序，覆盖浏览、问答、推荐、活动、社区全链路',
        'AI数字人黑塔实现多轮问答、语音播报、CRS状态显示、延伸推荐、行动清单生成',
        "本地知识库优先+大模型回退策略，兼顾知识准确性与语义理解",
        "为非遗数字化传播与个性化服务提供可落地的实践参考",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(4), practice_items, font_size=13, spacing=Pt(10), bullet_color=JADE_GREEN)

    # ===== Slide 5: Research Status 1 =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "02", "国内外研究现状", "推荐系统、CRS、大语言模型融合")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 5, total_slides)

    sections = [
        ("推荐系统", CINNABAR, [
            "协同过滤：依赖历史行为，冷启动突出",
            "内容推荐：特征匹配，多样性不足",
            "混合推荐：融合多源信号，适合稀疏场景",
        ]),
        ("对话推荐CRS", INK_BLUE, [
            "Jannach等：CRS通过多轮交互获取反馈",
            "赵梦媛等：传统推荐三大缺陷",
            "Wang等：知识增强提示学习统一框架",
        ]),
        ("LLM+推荐", JADE_GREEN, [
            "谢广明等：判别式与生成式范式",
            "Wu等：LLM for Recommendation综述",
            "趋势：走向LLM驱动的交互式推荐",
        ]),
    ]
    for i, (title, accent, items) in enumerate(sections):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.5), Inches(3.9), Inches(5.2), WHITE, LIGHT_GOLD)
        add_shape(slide, x, Inches(1.5), Inches(3.9), Inches(0.06), accent)
        add_text_box(slide, x + Inches(0.3), Inches(1.75), Inches(3.3), Inches(0.4),
                     title, font_size=18, color=accent, bold=True, font_name="Source Han Serif SC")
        add_bullet_list(slide, x + Inches(0.3), Inches(2.4), Inches(3.3), Inches(3.5),
                        items, font_size=13, spacing=Pt(12), bullet_color=accent)

    # ===== Slide 6: Research Status 2 =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "02", "AI数字人与研究不足", "现有研究空白与本课题定位")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 6, total_slides)

    top_left = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.3), WHITE, LIGHT_GOLD)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.35),
                 "AI数字人与文化导览", font_size=16, color=CINNABAR, bold=True, font_name="Source Han Serif SC")
    ai_items = [
        '康宁等辽小博：知识图谱+大模型融合微调的AI智慧导览系统',
        "Carreiro等InHeritage：游戏化+AR的文化遗产传播应用",
        "趋势：从静态展示走向AI驱动的智能导览",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(1.5), ai_items, font_size=12, spacing=Pt(6))

    top_right = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(2.3), WHITE, LIGHT_GOLD)
    add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5.2), Inches(0.35),
                 "知识图谱增强推荐", font_size=16, color=INK_BLUE, bold=True, font_name="Source Han Serif SC")
    kg_items = [
        "Guo等：知识图谱通过实体关系扩展推荐候选集",
        "王敏等：KG+LLM增强推荐系统研究",
        "汪天雄等：非遗数字化发展中知识图谱视角的前景",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.2), Inches(1.5), kg_items, font_size=12, spacing=Pt(6), bullet_color=INK_BLUE)

    gap_card = add_rounded_rect(slide, Inches(0.5), Inches(4.2), Inches(12.2), Inches(2.8), INK_BLACK)
    add_text_box(slide, Inches(0.8), Inches(4.35), Inches(11.6), Inches(0.4),
                 "现有研究不足与本课题定位", font_size=18, color=GOLD_LEAF, bold=True, font_name="Source Han Serif SC")
    gap_items = [
        "1.  现有非遗平台缺少CRS对话推荐能力，无法在交互中主动发现用户偏好",
        "2.  AI导览系统多侧重知识问答，缺少推荐引导与偏好采集的闭环设计",
        "3.  知识图谱在非遗推荐中的应用尚处起步阶段，推荐解释性不足",
        "4.  缺少面向非遗场景的完整系统实践，CRS+AI数字人+知识图谱三者协同方案尚未见报道",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(4.9), Inches(8), Inches(1.8), gap_items, font_size=13, color=WHITE, spacing=Pt(6), bullet_color=GOLD_LEAF)

    add_shape(slide, Inches(9.5), Inches(5.0), Inches(3), Inches(1.5), CINNABAR)
    add_text_box(slide, Inches(9.5), Inches(5.1), Inches(3), Inches(1.3),
                 "本课题定位\nCRS为引导\n图谱为增强\n数字人为载体",
                 font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")

    # ===== Slide 7: Technologies =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "03", "相关技术与理论", "系统技术栈全景")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 7, total_slides)

    techs = [
        ("前端", "微信小程序\nWXML/WXSS/JS", "原生开发\n微信能力直调", CINNABAR),
        ("后端", "FastAPI\nSQLAlchemy", "异步高性能\nORM多库切换", INK_BLUE),
        ("数据", "SQLite\n本地存储", "轻量关系型\n小规模部署", JADE_GREEN),
        ("AI", "豆包大模型\n本地知识库", "三级回退策略\nTTS双引擎", GOLD_LEAF),
        ("推荐", "CRS对话推荐\n知识图谱", "ASK-REC状态机\n规则加权混合", DEEP_RED),
    ]
    for i, (cat, name, desc, accent) in enumerate(techs):
        x = Inches(0.4 + i * 2.55)
        card = add_rounded_rect(slide, x, Inches(1.5), Inches(2.35), Inches(5.3), WHITE, LIGHT_GOLD)
        add_shape(slide, x, Inches(1.5), Inches(2.35), Inches(0.06), accent)
        add_text_box(slide, x + Inches(0.15), Inches(1.75), Inches(2.05), Inches(0.35),
                     cat, font_size=13, color=accent, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, x + Inches(0.15), Inches(2.2), Inches(2.05), Inches(1.2),
                     name, font_size=16, color=INK_BLACK, bold=True, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")
        add_shape(slide, x + Inches(0.3), Inches(3.5), Inches(1.75), Inches(0.01), LIGHT_GOLD)
        add_text_box(slide, x + Inches(0.15), Inches(3.7), Inches(2.05), Inches(2.5),
                     desc, font_size=12, color=SILVER_MIST, alignment=PP_ALIGN.CENTER)

    # ===== Slide 8: System Analysis =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "04", "系统分析", "三类角色与六大功能模块")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 8, total_slides)

    roles = [
        ("普通用户", "浏览内容、AI对话\n个性化推荐、活动报名\n社区讨论、偏好管理", CINNABAR),
        ("内容管理员", "内容审核发布\n质量评定、数据收集\n基本用户管理", INK_BLUE),
        ("系统管理员", "系统运维配置\n运行监控\n权限控制", JADE_GREEN),
    ]
    for i, (role, desc, accent) in enumerate(roles):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.5), Inches(3.9), Inches(1.6), WHITE, LIGHT_GOLD)
        add_shape(slide, x, Inches(1.5), Inches(3.9), Inches(0.06), accent)
        add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.35),
                     role, font_size=16, color=accent, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, x + Inches(0.2), Inches(2.1), Inches(3.5), Inches(0.9),
                     desc, font_size=12, color=INK_BLACK)

    modules = [
        ("非遗内容展示", "浏览/详情/分类/搜索"),
        ("个性化推荐", "首页推荐/冷启动/动态调整"),
        ("AI数字人对话", "知识问答/多轮会话/ASK追问"),
        ("活动报名", "活动浏览/在线报名/记录管理"),
        ("社区互动", "帖子/评论/点赞/话题标签"),
        ("用户中心", "微信登录/偏好设置/收藏管理"),
    ]
    for i, (mod, desc) in enumerate(modules):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(3.5 + row * 1.8)
        card = add_rounded_rect(slide, x, y, Inches(3.9), Inches(1.5), WHITE, LIGHT_GOLD)
        add_diamond(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.2), CINNABAR)
        add_text_box(slide, x + Inches(0.55), y + Inches(0.15), Inches(3.2), Inches(0.35),
                     mod, font_size=14, color=INK_BLACK, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, x + Inches(0.55), y + Inches(0.5), Inches(3.2), Inches(0.8),
                     desc, font_size=12, color=SILVER_MIST)

    # ===== Slide 9: Feasibility =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "04", "可行性分析", "技术/经济/操作/社会法律四维可行性")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 9, total_slides)

    feasibilities = [
        ("技术", "微信小程序成熟框架\nFastAPI+SQLAlchemy ORM\n知识库-搜索-大模型三级回退\n规则加权混合推荐", CINNABAR),
        ("经济", "SQLite零成本\n云服务器月费百元内\n豆包大模型免费额度\n小程序注册发布平价", INK_BLUE),
        ("操作", "小程序免安装扫码即用\nAI数字人零学习成本\n管理端Web界面直观操作", JADE_GREEN),
        ("社会法律", "本地知识库降低AI幻觉\n符合个人信息保护法\n传播非遗有积极社会意义\n符合非物质文化遗产法", GOLD_LEAF),
    ]
    for i, (title, desc, accent) in enumerate(feasibilities):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(1.5 + row * 2.8)
        card = add_rounded_rect(slide, x, y, Inches(6), Inches(2.4), WHITE, LIGHT_GOLD)
        add_shape(slide, x, y, Inches(0.06), Inches(2.4), accent)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.4), Inches(0.4),
                     title, font_size=18, color=accent, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, x + Inches(0.3), y + Inches(0.65), Inches(5.4), Inches(1.6),
                     desc, font_size=13, color=INK_BLACK)

    # ===== Slide 10: Architecture =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "05", "系统整体架构", "前后端分离四层架构")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 10, total_slides)

    layers = [
        ("前端交互层", "微信小程序用户端  |  Web管理端", RGBColor(0xF7, 0xF2, 0xE8), INK_BLACK, CINNABAR),
        ("接口服务层", "AI问答 | CRS状态 | 用户画像 | 内容活动 | 讨论区 | 登录", RGBColor(0xE8, 0xE0, 0xD4), INK_BLACK, CINNABAR),
        ("核心能力层", "AI服务 | 推荐服务 | 本地知识检索 | 知识图谱增强 | 联网补充 | TTS语音", CINNABAR, WHITE, GOLD_LEAF),
        ("数据与外部服务层", "SQLite | 本地知识 | 推荐日志 | CRS会话 | 豆包API | 联网搜索 | 微信登录", INK_BLACK, WHITE, GOLD_LEAF),
    ]
    for i, (name, content, bg, text_c, accent_c) in enumerate(layers):
        y = Inches(1.5 + i * 1.4)
        card = add_rounded_rect(slide, Inches(1.8), y, Inches(10.5), Inches(1.15), bg, accent_c)
        add_text_box(slide, Inches(2.1), y + Inches(0.1), Inches(2.5), Inches(0.35),
                     name, font_size=16, color=text_c, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, Inches(2.1), y + Inches(0.5), Inches(10), Inches(0.5),
                     content, font_size=12, color=text_c)

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(1.1), Inches(0.4),
                 "设计原则", font_size=12, color=CINNABAR, bold=True, font_name="Source Han Serif SC")
    add_text_box(slide, Inches(0.5), Inches(1.9), Inches(1.1), Inches(3.5),
                 "前后端分离\n主线趋同\n轻量优先\n可扩展", font_size=11, color=SILVER_MIST)

    # ===== Slide 11: CRS Design =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "05", "CRS对话推荐机制", "三维置信度评估 + ASK-REC状态机")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 11, total_slides)

    left_card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_text_box(slide, Inches(0.8), Inches(1.65), Inches(5.2), Inches(0.4),
                 "三维置信度评估模型", font_size=18, color=CINNABAR, bold=True, font_name="Source Han Serif SC")
    add_shape(slide, Inches(0.8), Inches(2.15), Inches(4.5), Inches(0.6), INK_BLACK)
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(4.5), Inches(0.5),
                 "C = 0.40\u00d7Se + 0.35\u00d7Si + 0.25\u00d7Sd",
                 font_size=18, color=GOLD_LEAF, bold=True, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")
    conf_items = [
        ("Se 显式偏好  ", "用户偏好设置、ASK选项回答"),
        ("Si 隐式行为  ", "曝光、点击、浏览、活动报名、讨论互动"),
        ("Sd 对话语义  ", "AI提问中的非遗关键词、地区特征、场景偏好"),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(3.0), Inches(5.2), Inches(3), conf_items, font_size=13, spacing=Pt(12))

    right_card = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_text_box(slide, Inches(7.2), Inches(1.65), Inches(5.2), Inches(0.4),
                 "ASK-REC 状态机", font_size=18, color=CINNABAR, bold=True, font_name="Source Han Serif SC")

    modes = [
        ("cold_start", "C < 28", "优先ASK采集偏好", RGBColor(0x66, 0x99, 0xCC)),
        ("mixed", "28 \u2264 C < 62", "边推荐边追问", RGBColor(0xDD, 0xAA, 0x44)),
        ("precision", "C \u2265 62", "精准推荐为主", RGBColor(0x66, 0xBB, 0x66)),
    ]
    for i, (mode, threshold, desc, color) in enumerate(modes):
        y = Inches(2.3 + i * 1.3)
        badge = add_rounded_rect(slide, Inches(7.2), y, Inches(2.2), Inches(0.4), color)
        add_text_box(slide, Inches(7.2), y + Inches(0.03), Inches(2.2), Inches(0.35),
                     mode, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")
        add_text_box(slide, Inches(9.6), y, Inches(1.5), Inches(0.35),
                     threshold, font_size=12, color=INK_BLACK, bold=True)
        add_text_box(slide, Inches(7.2), y + Inches(0.5), Inches(5.2), Inches(0.6),
                     desc, font_size=13, color=SILVER_MIST)

    add_text_box(slide, Inches(7.2), Inches(5.8), Inches(5.2), Inches(0.8),
                 "特殊策略：意图驱动推荐 | 恢复式提问",
                 font_size=12, color=SILVER_MIST, font_name="Source Han Serif SC")

    # ===== Slide 12: AI Digital Human =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "05", 'AI数字人"黑塔"', "四大职责 + 回答链路 + 提示词约束")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 12, total_slides)

    duties = [
        ("知识传授", "非遗问题解释\n导览入口", CINNABAR),
        ("推荐引导", "提问引向推荐\n不停留在问答", INK_BLUE),
        ("采集偏好", "ASK模板收集\n类别/地区/体验", JADE_GREEN),
        ("闭环触发", "行为反馈画像\nCRS状态深化", GOLD_LEAF),
    ]
    for i, (duty, desc, accent) in enumerate(duties):
        x = Inches(0.5 + i * 3.15)
        card = add_rounded_rect(slide, x, Inches(1.5), Inches(2.9), Inches(1.8), WHITE, LIGHT_GOLD)
        add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(0.06), accent)
        add_text_box(slide, x + Inches(0.15), Inches(1.7), Inches(2.6), Inches(0.35),
                     duty, font_size=15, color=accent, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, x + Inches(0.15), Inches(2.1), Inches(2.6), Inches(1),
                     desc, font_size=12, color=INK_BLACK, alignment=PP_ALIGN.CENTER)

    chain_card = add_rounded_rect(slide, Inches(0.5), Inches(3.6), Inches(12.2), Inches(1.4), INK_BLACK)
    add_text_box(slide, Inches(0.8), Inches(3.7), Inches(4), Inches(0.35),
                 "AI回答链路（本地知识优先）", font_size=15, color=GOLD_LEAF, bold=True, font_name="Source Han Serif SC")
    add_text_box(slide, Inches(0.8), Inches(4.15), Inches(11.6), Inches(0.6),
                 "接收问题 \u2192 检索本地知识库 \u2192 读取CRS会话 \u2192 生成推荐载荷 \u2192 ASK-REC决策 \u2192 调用豆包API \u2192 返回回答+推荐卡+ASK选项",
                 font_size=13, color=WHITE)

    prompt_card = add_rounded_rect(slide, Inches(0.5), Inches(5.3), Inches(12.2), Inches(1.7), WHITE, LIGHT_GOLD)
    add_text_box(slide, Inches(0.8), Inches(5.4), Inches(4), Inches(0.35),
                 "三层提示词约束", font_size=15, color=CINNABAR, bold=True, font_name="Source Han Serif SC")
    prompt_items = [
        "第一层  全局角色提示：确定黑塔身份、语气、服务范围和知识边界",
        "第二层  任务型提示：知识命中时润色简化，未命中时限定输出范围",
        "第三层  CRS感知提示：根据cold_start/mixed/precision动态调整语气深度",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(5.85), Inches(11.6), Inches(1), prompt_items, font_size=12, spacing=Pt(5), bullet_color=CINNABAR)

    # ===== Slide 13: KG & Governance =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "05", "知识图谱与内容治理", "图谱增强推荐 + 三维度质量评分")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 13, total_slides)

    kg_card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.06), INK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.2), Inches(0.4),
                 "知识图谱设计", font_size=18, color=INK_BLUE, bold=True, font_name="Source Han Serif SC")
    kg_items = [
        "定位：推荐和AI对话的加强器，非主排序器",
        "扩大推荐候选：相近实体、类别关系、路径关系",
        "加强推荐解释：路径关系转为自然语言推荐理由",
        "以用户已知实体为起点，向丰富方向延伸",
        "CRS联动：识别实体\u2192映射画像\u2192附加图谱理由",
        "一石二鸟：增强解释 + 参与兴趣建模",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(4), kg_items, font_size=13, spacing=Pt(8), bullet_color=INK_BLUE)

    gov_card = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.06), JADE_GREEN)
    add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.2), Inches(0.4),
                 "内容治理体系", font_size=18, color=JADE_GREEN, bold=True, font_name="Source Han Serif SC")
    gov_items = [
        "质量评分（三维度量化）：",
        "  内容完整度：正文字数、封面图、摘要",
        "  信息丰富度：章节归属、标签覆盖",
        "  原创质量：来源可信度、内容哈希去重",
        "审核发布：质量分\u2265阈值自动进推荐池",
        "白名单回补：定期从达标内容重新选拔精选",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(4), gov_items, font_size=13, spacing=Pt(8), bullet_color=JADE_GREEN)

    # ===== Slide 14: Homepage =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "06", "首页与个性化推荐", "AI导览横幅 + 精选推荐 + 快速入口")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 14, total_slides)

    features = [
        ("AI导览横幅", "数字导览中枢徽章\nCRS模式进度胶囊\n（想认识你\u2192正在了解你\u2192已懂你）\n黑塔数字人组件\n（表情随CRS模式变化）\n立即开始按钮", CINNABAR),
        ("精选推荐区", "第一条推荐大卡片展示\n今日推荐：内容/活动/讨论\n推荐理由标签+新鲜度标签\n点击行为即时回流画像", INK_BLUE),
        ("快速入口+TabBar", "非遗发展史/非遗地点快捷入口\n自定义TabBar\n主页/文化/活动/讨论/我的\n胶囊圆角设计\n选中态朱红渐变背景", JADE_GREEN),
    ]
    for i, (title, desc, accent) in enumerate(features):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.5), Inches(3.9), Inches(5.3), WHITE, LIGHT_GOLD)
        add_shape(slide, x, Inches(1.5), Inches(3.9), Inches(0.06), accent)
        add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.35),
                     title, font_size=16, color=accent, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, x + Inches(0.2), Inches(2.2), Inches(3.5), Inches(4.2),
                     desc, font_size=13, color=INK_BLACK)

    # ===== Slide 15: Content & Activity =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "06", "非遗内容与活动报名", "瀑布流布局 + 在线报名")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 15, total_slides)

    content_card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.06), CINNABAR)
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.2), Inches(0.4),
                 "非遗内容模块", font_size=18, color=CINNABAR, bold=True, font_name="Source Han Serif SC")
    content_items = [
        "内容列表：瀑布流双列布局",
        "策展精选显眼位置",
        "按子章节和类型打标签",
        "个性化内容推荐混合展示",
        "内容详情：封面图+三分钟看点+延伸讲解",
        "AI浮窗入口（场景化问答）",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(4), content_items, font_size=13, spacing=Pt(8))

    activity_card = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.06), JADE_GREEN)
    add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.2), Inches(0.4),
                 "活动报名模块", font_size=18, color=JADE_GREEN, bold=True, font_name="Source Han Serif SC")
    activity_items = [
        "活动列表：按月份筛选",
        '前两条标记本月主推',
        "个性化活动推荐混合展示",
        "活动详情：封面/时间地点/组织者",
        "展示块（亮点/议程/提示）",
        "在线报名/取消报名+AI浮窗",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(4), activity_items, font_size=13, spacing=Pt(8), bullet_color=JADE_GREEN)

    # ===== Slide 16: Community & User =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "06", "社区讨论与用户中心", "话题互动 + 偏好管理")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 16, total_slides)

    community_card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.06), INK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.2), Inches(0.4),
                 "社区讨论模块", font_size=18, color=INK_BLUE, bold=True, font_name="Source Han Serif SC")
    community_items = [
        "讨论列表：关键词搜索+标签筛选（10个标签）",
        "排序切换（热门/最新）、只看收藏",
        "热门话题TOP3展示",
        "发帖三种模板：提问型/体验分享型/活动反馈型",
        "话题详情：帖子正文+评论列表",
        "发表评论、点赞、收藏",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(4), community_items, font_size=13, spacing=Pt(8), bullet_color=INK_BLUE)

    user_card = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.06), GOLD_LEAF)
    add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.2), Inches(0.4),
                 "用户中心模块", font_size=18, color=GOLD_LEAF, bold=True, font_name="Source Han Serif SC")
    user_items = [
        "用户画像卡片：偏好关键词/活跃场景/关注地区",
        "统计卡片：发帖数/活动记录/已选偏好",
        "偏好设置：非遗类别（工艺/戏曲/民俗/医药）",
        "场景类型（知识阅读/活动体验/论坛交流）",
        "地区（华东/华南/西南/华北/西北/东北）",
        "保存后作为显式信号纳入画像",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(4), user_items, font_size=13, spacing=Pt(8), bullet_color=GOLD_LEAF)

    # ===== Slide 17: AI Modes =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "06", 'AI数字人"黑塔"核心交互', "cold_start \u2192 mixed \u2192 precision 三阶段渐进")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 17, total_slides)

    mode_details = [
        ("cold_start\n冷启动", "C < 28", [
            "主动提问为主",
            "ASK卡片收集偏好",
            "类目\u2192地区\u2192场景\u2192程度",
            "不直接给出推荐结果",
            "数字人好奇张望表情",
        ], RGBColor(0x66, 0x99, 0xCC)),
        ("mixed\n混合模式", "28 \u2264 C < 62", [
            "边推荐边追问",
            "1-2张AI推荐卡+B组追问",
            "上方推荐+下方ASK选项",
            "推荐理由贴合当前选择",
            "数字人托腮思考表情",
        ], RGBColor(0xDD, 0xAA, 0x44)),
        ("precision\n精准模式", "C \u2265 62", [
            "停止追问直接输出",
            "高度匹配推荐结果",
            "精准匹配偏好表达",
            "推荐卡占主要区域",
            "数字人自信微笑表情",
        ], RGBColor(0x66, 0xBB, 0x66)),
    ]
    for i, (mode, threshold, items, color) in enumerate(mode_details):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.5), Inches(3.9), Inches(5.3), WHITE, LIGHT_GOLD)
        badge = add_rounded_rect(slide, x + Inches(0.5), Inches(1.7), Inches(2.9), Inches(0.8), color)
        add_text_box(slide, x + Inches(0.5), Inches(1.7), Inches(2.9), Inches(0.8),
                     mode, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")
        add_text_box(slide, x + Inches(0.5), Inches(2.6), Inches(2.9), Inches(0.3),
                     threshold, font_size=12, color=INK_BLACK, bold=True, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + Inches(0.3), Inches(3.1), Inches(3.3), Inches(3.5),
                        items, font_size=12, spacing=Pt(6), bullet_color=color)

    # ===== Slide 18: AI Features =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "06", "AI对话功能与信息层次", "语音播报 + 延伸推荐 + 行动清单 + 三层信息结构")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 18, total_slides)

    feature_cards = [
        ("语音播报", "豆包TTS优先\nEdge-TTS降级\n双引擎架构", CINNABAR),
        ("延伸推荐", "回答下方展示\n内容/活动/讨论\n推荐卡片即时刷新", INK_BLUE),
        ("行动清单", "浏览/报名/讨论\n下一步行动建议\n闭环触发机制", JADE_GREEN),
        ("AI浮窗", "内容页/活动页嵌入\n场景化问答\n随时发起对话", GOLD_LEAF),
    ]
    for i, (title, desc, accent) in enumerate(feature_cards):
        x = Inches(0.5 + i * 3.15)
        card = add_rounded_rect(slide, x, Inches(1.5), Inches(2.9), Inches(2.2), WHITE, LIGHT_GOLD)
        add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(0.06), accent)
        add_text_box(slide, x + Inches(0.15), Inches(1.7), Inches(2.6), Inches(0.35),
                     title, font_size=14, color=accent, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, x + Inches(0.15), Inches(2.1), Inches(2.6), Inches(1.4),
                     desc, font_size=12, color=INK_BLACK, alignment=PP_ALIGN.CENTER)

    info_card = add_rounded_rect(slide, Inches(0.5), Inches(4.0), Inches(12.2), Inches(3), INK_BLACK)
    add_text_box(slide, Inches(0.8), Inches(4.15), Inches(5), Inches(0.4),
                 "AI对话页三层信息结构", font_size=16, color=GOLD_LEAF, bold=True, font_name="Source Han Serif SC")
    info_layers = [
        ("第一层  任务主线层", "黑塔数字人形象 | 对话消息区 | 输入框 | 当前模式提示 | ASK提问卡片"),
        ("第二层  推荐结果层", "AI推荐卡 | 推荐理由 | 下一步建议（浏览/报名/讨论）"),
        ("第三层  解释层", "置信度详细部分 | 策略解释 | 图谱依据"),
    ]
    for i, (layer, desc) in enumerate(info_layers):
        y = Inches(4.7 + i * 0.7)
        add_rich_text_box(slide, Inches(0.8), y, Inches(11.6), Inches(0.6),
                          [{"text": layer, "size": 13, "color": GOLD_LEAF, "bold": True, "font": "Source Han Serif SC"},
                           {"text": "  " + desc, "size": 12, "color": WHITE, "bold": False}])

    # ===== Slide 19: Admin =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "06", "Web管理端", "七大管理模块")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 19, total_slides)

    admin_modules = [
        ("内容管理", "卡片网格\n封面/标题/状态/质量分", CINNABAR),
        ("质量检查", "四项统计指标\n五档质量分布柱状图", INK_BLUE),
        ("活动管理", "卡片展示\n上架下架管理", JADE_GREEN),
        ("帖子管理", "帖子列表\n编辑/精选/删除", GOLD_LEAF),
        ("用户管理", "表格展示\n行为统计查看", DEEP_RED),
        ("统计导出", "平台指标\nCSV导出功能", RGBColor(0x66, 0x99, 0xCC)),
        ("知识库管理", "表格展示\n新增/编辑/启禁用", RGBColor(0x99, 0x66, 0xCC)),
    ]
    for i, (title, desc, accent) in enumerate(admin_modules):
        col = i % 4
        row = i // 4
        x = Inches(0.4 + col * 3.2)
        y = Inches(1.5 + row * 2.8)
        card = add_rounded_rect(slide, x, y, Inches(3), Inches(2.4), WHITE, LIGHT_GOLD)
        add_shape(slide, x, y, Inches(3), Inches(0.06), accent)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.7), Inches(0.35),
                     title, font_size=14, color=accent, bold=True, font_name="Source Han Serif SC")
        add_text_box(slide, x + Inches(0.15), y + Inches(0.65), Inches(2.7), Inches(1.5),
                     desc, font_size=12, color=INK_BLACK)

    # ===== Slide 20: Testing =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "07", "系统测试", "功能测试 + 性能测试 + 兼容性测试")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 20, total_slides)

    test_sections = [
        ("功能测试", CINNABAR, [
            "黑盒测试：10条用例，通过10条",
            "覆盖8模块：认证/内容/活动/讨论/推荐/AI/图谱/管理",
            "白盒测试：7条用例，通过7条",
            "CRS决策引擎/ASK模板/知识图谱/数据模型",
        ]),
        ("性能测试", INK_BLUE, [
            "API平均响应时间：43.16ms",
            "最大响应时间：51.33ms（均<100ms）",
            "并发10用户50请求：成功率100%",
            "吞吐量225.5请求/秒",
            "AI对话平均响应：7.9秒",
        ]),
        ("兼容性测试", JADE_GREEN, [
            "浏览器：Chrome/Firefox/Safari/Edge",
            "微信基础库：\u22652.20.0",
            "设备：iOS + Android",
            "数据库：SQLite/MySQL/PostgreSQL",
        ]),
    ]
    for i, (title, accent, items) in enumerate(test_sections):
        x = Inches(0.5 + i * 4.2)
        card = add_rounded_rect(slide, x, Inches(1.5), Inches(3.9), Inches(5.3), WHITE, LIGHT_GOLD)
        add_shape(slide, x, Inches(1.5), Inches(3.9), Inches(0.06), accent)
        add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.35),
                     title, font_size=16, color=accent, bold=True, font_name="Source Han Serif SC")
        add_bullet_list(slide, x + Inches(0.2), Inches(2.2), Inches(3.5), Inches(4),
                        items, font_size=12, spacing=Pt(8), bullet_color=accent)

    # ===== Slide 21: Experiment =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "07", "实验与分析", "AI对话质量评估 + 推荐评价指标")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 21, total_slides)

    exp_card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.06), CINNABAR)
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.2), Inches(0.4),
                 "AI对话质量评估", font_size=18, color=CINNABAR, bold=True, font_name="Source Han Serif SC")

    add_shape(slide, Inches(0.8), Inches(2.2), Inches(2.5), Inches(0.5), CINNABAR)
    add_text_box(slide, Inches(0.8), Inches(2.22), Inches(2.5), Inches(0.45),
                 "知识库命中率 90.0%", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")

    exp_items = [
        "20个测试问题，15个预期命中全部命中",
        "5个冷门问题中3个未命中\u2192回退豆包/联网",
        "命中项置信度集中在0.80-0.85区间",
        "P50=2.6秒，P90=4.3秒",
        "流式输出3秒内可见首个回答片段",
        "KB命中+润色占比最大（链路有效）",
        "豆包直答+组合占比较大（大模型体现）",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.9), Inches(5.2), Inches(3.5), exp_items, font_size=12, spacing=Pt(6))

    eval_card = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.06), INK_BLUE)
    add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.2), Inches(0.4),
                 "推荐评价指标", font_size=18, color=INK_BLUE, bold=True, font_name="Source Han Serif SC")

    metrics = [
        ("Precision@5", "推荐列表前5项中与用户偏好相关的比例\n\u2192 反映推荐准确性"),
        ("Diversity@5", "推荐列表覆盖的非遗类别数/列表长度\n\u2192 反映推荐多样性"),
        ("Coverage", "推荐结果覆盖内容占全部已发布内容的比例\n\u2192 反映推荐覆盖广度"),
        ("NDCG@5", "考虑位置权重的排序质量\n\u2192 反映推荐排序质量"),
    ]
    for i, (metric, desc) in enumerate(metrics):
        y = Inches(2.3 + i * 1.15)
        badge = add_rounded_rect(slide, Inches(7.2), y, Inches(2.2), Inches(0.35), INK_BLUE)
        add_text_box(slide, Inches(7.2), y + Inches(0.02), Inches(2.2), Inches(0.3),
                     metric, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")
        add_text_box(slide, Inches(9.6), y, Inches(3), Inches(0.9),
                     desc, font_size=11, color=INK_BLACK)

    # ===== Slide 22: Conclusion =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WARM_CREAM)
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CINNABAR)
    add_section_header(slide, "08", "结论与展望", "研究成果总结与未来方向")
    add_decorative_corner(slide, "top-right")
    add_page_number(slide, 22, total_slides)

    conclusion_card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.3), WHITE, LIGHT_GOLD)
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(7.5), Inches(0.06), CINNABAR)
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
                 "研究结论", font_size=18, color=CINNABAR, bold=True, font_name="Source Han Serif SC")
    conclusion_items = [
        "构建了基于CRS推荐与AI数字人的非遗文化传播微信小程序",
        'AI数字人黑塔通过ASK追问逐步构建用户偏好画像',
        "从冷启动到精准推荐只需3-4轮交互即可收敛",
        "本地知识优先+大模型回退策略兼顾知识准确性与语义理解",
        "ASK-REC决策引擎执行三阶段置信度收敛机制",
        "知识图谱实现实体偏好映射和推荐解释生成",
        "系统性能满足中小规模使用场景",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(7), Inches(4), conclusion_items, font_size=13, spacing=Pt(8))

    future_card = add_rounded_rect(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.3), INK_BLACK)
    add_text_box(slide, Inches(8.8), Inches(1.7), Inches(3.7), Inches(0.4),
                 "未来展望", font_size=18, color=GOLD_LEAF, bold=True, font_name="Source Han Serif SC")
    future_items = [
        "加入更细致的知识图谱推理能力",
        "同线下非遗场馆深度对接",
        "支持多模态交互",
        "大语言模型技术发展推动对话更自然",
    ]
    add_bullet_list(slide, Inches(8.8), Inches(2.3), Inches(3.7), Inches(4), future_items, font_size=13, color=WHITE, spacing=Pt(14), bullet_color=GOLD_LEAF)

    # ===== Slide 23: Thanks =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, INK_BLACK)

    add_shape(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), CINNABAR)
    add_shape(slide, Inches(0.12), Inches(0), Inches(0.04), Inches(7.5), GOLD_LEAF)

    add_shape(slide, Inches(2), Inches(2.3), Inches(9.3), Inches(0.02), GOLD_LEAF)
    add_text_box(slide, Inches(2), Inches(2.6), Inches(9.3), Inches(1),
                 "感谢各位老师的指导与评审", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")
    add_shape(slide, Inches(2), Inches(3.7), Inches(9.3), Inches(0.02), GOLD_LEAF)

    add_text_box(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.5),
                 "恳请各位老师批评指正", font_size=20, color=GOLD_LEAF, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")

    add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.4),
                 "答辩人：王子轩    指导教师：徐龙琴 教授  高静 高级工程师",
                 font_size=14, color=SILVER_MIST, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")
    add_text_box(slide, Inches(2), Inches(5.9), Inches(9.3), Inches(0.4),
                 "仲恺农业工程学院  信息科学与技术学院  2026.05.09",
                 font_size=12, color=SILVER_MIST, alignment=PP_ALIGN.CENTER, font_name="Source Han Serif SC")

    add_diamond(slide, Inches(11.5), Inches(1.0), Inches(0.3), CINNABAR)
    add_diamond(slide, Inches(1.5), Inches(6.0), Inches(0.2), GOLD_LEAF)

    output_path = r"d:\桌面\毕业设计\答辩演示文件夹\答辩PPT_v2_基于CRS推荐与AI数字人的非遗文化传播系统.pptx"
    prs.save(output_path)
    print(f"PPT生成完成: {output_path}")
    print(f"文件大小: {os.path.getsize(output_path)} 字节")
    print(f"总页数: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
