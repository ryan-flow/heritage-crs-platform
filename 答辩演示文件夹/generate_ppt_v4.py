from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

BG_WARM = RGBColor(0xF5, 0xF0, 0xE8)
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)
BG_DARK = RGBColor(0x2C, 0x1A, 0x12)
CINNABAR = RGBColor(0xC4, 0x1E, 0x3A)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
DARK_RED = RGBColor(0x8B, 0x22, 0x22)
BROWN = RGBColor(0x5C, 0x3A, 0x21)
BODY = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x88, 0x77, 0x66)
LIGHT_TAN = RGBColor(0xE8, 0xD5, 0xC4)
LIGHT_GOLD = RGBColor(0xF0, 0xE4, 0xC8)
TAG_BLUE = RGBColor(0x3B, 0x82, 0xF6)
TAG_GREEN = RGBColor(0x22, 0xC5, 0x5E)
TAG_AMBER = RGBColor(0xD4, 0x8B, 0x2C)
TAG_PURPLE = RGBColor(0x7C, 0x5C, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SERIF = "Source Han Serif SC"
SANS = "Source Han Sans SC"
MONO = "Consolas"
T = 22


def set_bg(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c


def rect(s, l, t, w, h, fill, line=None, radius=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def txt(s, l, t, w, h, text, sz=16, c=BODY, bold=False, align=PP_ALIGN.LEFT, fn=SANS):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = bold
    p.font.name = fn
    p.alignment = align
    return tb


def rich(s, l, t, w, h, runs, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for rd in runs:
        r = p.add_run()
        r.text = rd.get("t", "")
        r.font.size = Pt(rd.get("s", 14))
        r.font.color.rgb = rd.get("c", BODY)
        r.font.bold = rd.get("b", False)
        r.font.name = rd.get("f", SANS)
    return tb


def bullets(s, l, t, w, h, items, sz=13, c=BODY, sp=Pt(5), bc=CINNABAR):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = sp
        if isinstance(item, tuple):
            label, desc = item
            r1 = p.add_run()
            r1.text = label
            r1.font.size = Pt(sz)
            r1.font.color.rgb = bc
            r1.font.bold = True
            r1.font.name = SANS
            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(sz)
            r2.font.color.rgb = c
            r2.font.name = SANS
        else:
            r1 = p.add_run()
            r1.text = "\u25B8 "
            r1.font.size = Pt(sz - 1)
            r1.font.color.rgb = bc
            r1.font.name = SANS
            r2 = p.add_run()
            r2.text = item
            r2.font.size = Pt(sz)
            r2.font.color.rgb = c
            r2.font.name = SANS
    return tb


def card(s, l, t, w, h, accent=CINNABAR):
    rect(s, l, t, w, h, BG_CARD, LIGHT_TAN, radius=True)
    rect(s, l, t + Inches(0.02), w, Inches(0.035), accent)


def sec_num(s, num, top=Inches(0.5)):
    rect(s, Inches(0.7), top, Inches(0.48), Inches(0.48), CINNABAR, radius=True)
    txt(s, Inches(0.7), top + Inches(0.05), Inches(0.48), Inches(0.38),
        num, sz=14, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=SERIF)


def pnum(s, n):
    txt(s, Inches(12.2), Inches(7.1), Inches(0.9), Inches(0.25),
        f"{n}/{T}", sz=9, c=MUTED, align=PP_ALIGN.RIGHT, fn=MONO)


def divider(s):
    rect(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.01), LIGHT_TAN)


def create():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    bl = prs.slide_layouts[6]

    # === 1 COVER ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_DARK)
    rect(s, Inches(0), Inches(0), Inches(0.1), Inches(7.5), CINNABAR)
    rect(s, Inches(0.1), Inches(0), Inches(0.03), Inches(7.5), GOLD)
    rect(s, Inches(1.2), Inches(1.5), Inches(0.04), Inches(4.5), GOLD)
    txt(s, Inches(1.8), Inches(1.5), Inches(9), Inches(0.4),
        "\u6bd5\u4e1a\u8bbe\u8ba1\u7b54\u8fa9", sz=15, c=GOLD, fn=SERIF)
    txt(s, Inches(1.8), Inches(2.0), Inches(9), Inches(1.5),
        "\u57fa\u4e8eCRS\u63a8\u8350\u4e0eAI\u6570\u5b57\u4eba\u7684\n\u975e\u9057\u6587\u5316\u4f20\u64ad\u7cfb\u7edf\u5b9e\u73b0",
        sz=38, c=WHITE, bold=True, fn=SERIF)
    rect(s, Inches(1.8), Inches(3.7), Inches(3), Inches(0.02), GOLD)
    info = [("\u7b54\u8fa9\u4eba ", "\u738b\u5b50\u8f69"), ("\u5b66    \u53f7 ", "202210274225"),
            ("\u4e13    \u4e1a ", "\u6570\u636e\u79d1\u5b66\u4e0e\u5927\u6570\u636e\u6280\u672f222\u73ed"),
            ("\u5b66    \u9662 ", "\u4fe1\u606f\u79d1\u5b66\u4e0e\u6280\u672f\u5b66\u9662"),
            ("\u6307\u5bfc\u6559\u5e08 ", "\u5f90\u9f99\u7434 \u6559\u6388  \u9ad8\u9759 \u9ad8\u7ea7\u5de5\u7a0b\u5e08")]
    for i, (lb, vl) in enumerate(info):
        rich(s, Inches(1.8), Inches(4.0 + i * 0.42), Inches(8), Inches(0.4),
             [{"t": lb, "s": 14, "c": RGBColor(0xBB, 0xAA, 0x99)}, {"t": vl, "s": 14, "c": WHITE, "b": True}])
    txt(s, Inches(1.8), Inches(6.3), Inches(4), Inches(0.35),
        "\u4ef2\u607a\u519c\u4e1a\u5de5\u7a0b\u5b66\u9662", sz=18, c=GOLD, bold=True, fn=SERIF)
    txt(s, Inches(1.8), Inches(6.7), Inches(4), Inches(0.25),
        "2026.05.09", sz=13, c=MUTED, fn=MONO)

    # === 2 TOC ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    txt(s, Inches(0.7), Inches(0.4), Inches(3), Inches(0.3),
        "CONTENTS", sz=10, c=CINNABAR, fn=MONO)
    txt(s, Inches(0.7), Inches(0.7), Inches(6), Inches(0.5),
        "\u6c47\u62a5\u63d0\u7eb2", sz=28, c=DARK_RED, bold=True, fn=SERIF)
    rect(s, Inches(0.7), Inches(1.2), Inches(1.5), Inches(0.02), CINNABAR)
    pnum(s, 2)
    toc = [("01", "\u7814\u7a76\u80cc\u666f\u4e0e\u610f\u4e49", "\u975e\u9057\u6570\u5b57\u5316\u4f20\u64ad\u7684\u65f6\u4ee3\u9700\u6c42\u4e0e\u8bfe\u9898\u5b9a\u4f4d"),
           ("02", "\u56fd\u5185\u5916\u7814\u7a76\u73b0\u72b6", "CRS\u3001AI\u6570\u5b57\u4eba\u7814\u7a76\u8fdb\u5c55\u4e0e\u4e0d\u8db3"),
           ("03", "\u76f8\u5173\u6280\u672f\u4e0e\u7406\u8bba", "FastAPI\u3001LLM\u3001CRS\u6838\u5fc3\u7406\u8bba"),
           ("04", "\u7cfb\u7edf\u5206\u6790", "\u9700\u6c42\u5206\u6790\u3001\u53ef\u884c\u6027\u3001\u6570\u636e\u6d41\u7a0b"),
           ("05", "\u7cfb\u7edf\u8bbe\u8ba1", "\u67b6\u6784\u3001\u63a8\u8350\u6a21\u5757\u3001AI\u6570\u5b57\u4eba\u3001\u77e5\u8bc6\u56fe\u8c31"),
           ("06", "\u7cfb\u7edf\u5b9e\u73b0", "\u5c0f\u7a0b\u5e8f\u7aef\u4e0e\u7ba1\u7406\u7aef\u529f\u80fd\u5c55\u793a"),
           ("07", "\u7cfb\u7edf\u6d4b\u8bd5\u4e0e\u5b9e\u9a8c", "\u529f\u80fd/\u6027\u80fd/\u517c\u5bb9\u6027\u6d4b\u8bd5\u4e0e\u5b9e\u9a8c\u5206\u6790"),
           ("08", "\u7ed3\u8bba\u4e0e\u5c55\u671b", "\u7814\u7a76\u6210\u679c\u603b\u7ed3\u4e0e\u672a\u6765\u65b9\u5411")]
    for i, (num, title, desc) in enumerate(toc):
        y = Inches(1.5 + i * 0.7)
        rect(s, Inches(0.7), y, Inches(0.45), Inches(0.45), CINNABAR if i % 2 == 0 else BROWN, radius=True)
        txt(s, Inches(0.7), y + Inches(0.05), Inches(0.45), Inches(0.35),
            num, sz=13, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=SERIF)
        txt(s, Inches(1.35), y + Inches(0.02), Inches(4), Inches(0.3),
            title, sz=16, c=DARK_RED, bold=True, fn=SERIF)
        txt(s, Inches(1.35), y + Inches(0.32), Inches(6), Inches(0.25),
            desc, sz=11, c=MUTED)

    # === 3 BACKGROUND ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "01")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u7814\u7a76\u80cc\u666f", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    txt(s, Inches(1.4), Inches(0.95), Inches(10), Inches(0.3),
        "\u975e\u9057\u6570\u5b57\u5316\u4f20\u64ad\u7684\u65f6\u4ee3\u9700\u6c42\u4e0e\u73b0\u6709\u5e73\u53f0\u4e0d\u8db3", sz=12, c=MUTED)
    divider(s)
    pnum(s, 3)

    card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4), TAG_GREEN)
    txt(s, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.35),
        "\u65f6\u4ee3\u9700\u6c42", sz=17, c=TAG_GREEN, bold=True, fn=SERIF)
    bullets(s, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4), [
        "\u56fd\u5bb6\u7ea7\u975e\u9057\u4ee3\u8868\u6027\u9879\u76ee1000+\uff0c\u8054\u5408\u56fd\u975e\u9057\u540d\u5f55\u6570\u91cf\u4e16\u754c\u7b2c\u4e00",
        "\u77ed\u89c6\u9891\u5316\uff1a\u6296\u97f3\u975e\u9057\u89c6\u9891\u8d852.2\u4ebf\u6761\uff0c138\u4ebf\u6b21\u70b9\u8d5e",
        "\u8d44\u6e90\u6570\u5b57\u5316\uff1a\u52a8\u753b\u3001\u6e38\u620f\u3001\u4ea7\u54c1\u8bbe\u8ba1\u9700\u6570\u5b57\u5316\u975e\u9057\u7d20\u6750",
        "\u4f53\u9a8c\u5316\u8f6c\u578b\uff1a\u975e\u9057\u4ece\u535a\u7269\u9986\u9644\u5c5e\u8d70\u5411\u72ec\u7acb\u6d3b\u52a8\u4ea7\u54c1",
        "\u4f20\u64ad\u91cd\u5fc3\u4ece\u4e13\u4e1a\u9886\u57df\u8f6c\u5411\u5927\u4f17\u89c6\u91ce",
    ], sz=13, sp=Pt(8), bc=TAG_GREEN)

    card(s, Inches(6.7), Inches(1.4), Inches(6), Inches(5.4), CINNABAR)
    txt(s, Inches(7.0), Inches(1.6), Inches(5.4), Inches(0.35),
        "\u73b0\u6709\u5e73\u53f0\u4e94\u5927\u4e0d\u8db3", sz=17, c=CINNABAR, bold=True, fn=SERIF)
    bullets(s, Inches(7.0), Inches(2.1), Inches(5.4), Inches(4), [
        ("\u63a8\u8350\u80fd\u529b\u5f31  ", "\u70ed\u95e8\u6392\u5e8f\u4e3a\u4e3b\uff0c\u65e0\u4e2a\u6027\u5316\u63a8\u8350"),
        ("\u4ea4\u4e92\u65b9\u5f0f\u5355\u4e00  ", "\u56fe\u6587\u5c55\u793a\u4e3a\u4e3b\uff0c\u7f3a\u5c11\u8bed\u97f3\u5bf9\u8bdd\u4ea4\u4e92"),
        ("\u573a\u666f\u611f\u77e5\u7f3a\u5931  ", "\u4e0d\u80fd\u6355\u6349\u65f6\u95f4\u5730\u70b9\u6587\u5316\u60c5\u5883\u4fe1\u53f7"),
        ("\u51b7\u542f\u52a8\u4e25\u91cd  ", "\u7528\u6237\u884c\u4e3a\u7a00\u758f\uff0c\u65b0\u7528\u6237\u6d41\u5931\u7387\u9ad8"),
        ("\u63a8\u8350\u4e0d\u53ef\u89e3\u91ca  ", "\u4e0d\u7ed9\u51fa\u63a8\u8350\u7406\u7531\uff0c\u5f71\u54cd\u4fe1\u4efb\u5ea6"),
    ], sz=13, sp=Pt(8), bc=CINNABAR)

    # === 4 SIGNIFICANCE ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "01")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u7814\u7a76\u610f\u4e49", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 4)

    card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4), TAG_BLUE)
    txt(s, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.35),
        "\u7406\u8bba\u610f\u4e49", sz=17, c=TAG_BLUE, bold=True, fn=SERIF)
    bullets(s, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4), [
        "\u5c06CRS\u5bf9\u8bdd\u63a8\u8350\u601d\u60f3\u5f15\u5165\u975e\u9057\u6587\u5316\u4f20\u64ad\u573a\u666f\uff0c\u63a2\u7d22SAUR\u4ea4\u4e92\u6a21\u5f0f",
        "\u63d0\u51fa\u4ee5CRS\u4e3a\u5f15\u5bfc\u3001\u77e5\u8bc6\u56fe\u8c31\u4e3a\u589e\u5f3a\u3001AI\u6570\u5b57\u4eba\u4e3a\u8f7d\u4f53\u7684\u667a\u80fd\u5316\u670d\u52a1\u6846\u67b6",
        "\u6784\u5efa\u4e09\u7ef4\u7f6e\u4fe1\u5ea6\u8bc4\u4f30\u6a21\u578b\uff0c\u5b9e\u73b0\u51b7\u542f\u52a8\u5230\u7cbe\u51c6\u63a8\u8350\u7684\u6e10\u8fdb\u6536\u655b",
        "\u4e3a\u975e\u9057\u63a8\u8350\u9886\u57df\u63d0\u4f9b\u53ef\u89e3\u91ca\u7684\u63a8\u8350\u673a\u5236\u8bbe\u8ba1\u53c2\u8003",
    ], sz=13, sp=Pt(10), bc=TAG_BLUE)

    card(s, Inches(6.7), Inches(1.4), Inches(6), Inches(5.4), TAG_GREEN)
    txt(s, Inches(7.0), Inches(1.6), Inches(5.4), Inches(0.35),
        "\u5b9e\u8df5\u610f\u4e49", sz=17, c=TAG_GREEN, bold=True, fn=SERIF)
    bullets(s, Inches(7.0), Inches(2.1), Inches(5.4), Inches(4), [
        "\u5f00\u53d1\u5b8c\u6574\u7684\u975e\u9057\u6587\u5316\u4f20\u64ad\u5fae\u4fe1\u5c0f\u7a0b\u5e8f\uff0c\u8986\u76d6\u6d4f\u89c8\u3001\u95ee\u7b54\u3001\u63a8\u8350\u3001\u6d3b\u52a8\u3001\u793e\u533a\u5168\u94fe\u8def",
        "AI\u6570\u5b57\u4eba\u9ed1\u5854\u5b9e\u73b0\u591a\u8f6e\u95ee\u7b54\u3001\u8bed\u97f3\u64ad\u62a5\u3001CRS\u72b6\u6001\u663e\u793a\u3001\u5ef6\u4f38\u63a8\u8350\u3001\u884c\u52a8\u6e05\u5355\u751f\u6210",
        "\u672c\u5730\u77e5\u8bc6\u5e93\u4f18\u5148+\u5927\u6a21\u578b\u56de\u9000\u7b56\u7565\uff0c\u517c\u987e\u77e5\u8bc6\u51c6\u786e\u6027\u4e0e\u8bed\u4e49\u7406\u89e3",
        "\u4e3a\u975e\u9057\u6570\u5b57\u5316\u4f20\u64ad\u4e0e\u4e2a\u6027\u5316\u670d\u52a1\u63d0\u4f9b\u53ef\u843d\u5730\u7684\u5b9e\u8df5\u53c2\u8003",
    ], sz=13, sp=Pt(10), bc=TAG_GREEN)

    # === 5 STATUS 1 ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "02")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u56fd\u5185\u5916\u7814\u7a76\u73b0\u72b6", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    txt(s, Inches(1.4), Inches(0.95), Inches(10), Inches(0.3),
        "\u63a8\u8350\u7cfb\u7edf\u3001CRS\u3001\u5927\u8bed\u8a00\u6a21\u578b\u878d\u5408", sz=12, c=MUTED)
    divider(s)
    pnum(s, 5)

    for i, (title, accent, items) in enumerate([
        ("\u63a8\u8350\u7cfb\u7edf", CINNABAR, [
            "\u534f\u540c\u8fc7\u6ee4\uff1a\u4f9d\u8d56\u5386\u53f2\u884c\u4e3a\uff0c\u51b7\u542f\u52a8\u7a81\u51fa",
            "\u5185\u5bb9\u63a8\u8350\uff1a\u7279\u5f81\u5339\u914d\uff0c\u591a\u6837\u6027\u4e0d\u8db3",
            "\u6df7\u5408\u63a8\u8350\uff1a\u878d\u5408\u591a\u6e90\u4fe1\u53f7\uff0c\u9002\u5408\u7a00\u758f\u573a\u666f"]),
        ("\u5bf9\u8bdd\u63a8\u8350CRS", TAG_BLUE, [
            "Jannach\u7b49\uff1aCRS\u901a\u8fc7\u591a\u8f6e\u4ea4\u4e92\u83b7\u53d6\u53cd\u9988",
            "\u8d75\u68a6\u5a9b\u7b49\uff1a\u4f20\u7edf\u63a8\u8350\u4e09\u5927\u7f3a\u9677",
            "Wang\u7b49\uff1a\u77e5\u8bc6\u589e\u5f3a\u63d0\u793a\u5b66\u4e60\u7edf\u4e00\u6846\u67b6"]),
        ("LLM+\u63a8\u8350", TAG_GREEN, [
            "\u8c22\u5e7f\u660e\u7b49\uff1a\u5224\u522b\u5f0f\u4e0e\u751f\u6210\u5f0f\u8303\u5f0f",
            "Wu\u7b49\uff1aLLM for Recommendation\u7efc\u8ff0",
            "\u8d8b\u52bf\uff1a\u8d70\u5411LLM\u9a71\u52a8\u7684\u4ea4\u4e92\u5f0f\u63a8\u8350"]),
    ]):
        x = Inches(0.5 + i * 4.2)
        card(s, x, Inches(1.4), Inches(3.9), Inches(5.2), accent)
        txt(s, x + Inches(0.3), Inches(1.65), Inches(3.3), Inches(0.35),
            title, sz=16, c=accent, bold=True, fn=SERIF)
        bullets(s, x + Inches(0.3), Inches(2.2), Inches(3.3), Inches(3.5), items, sz=12, sp=Pt(12), bc=accent)

    # === 6 STATUS 2 ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "02")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "AI\u6570\u5b57\u4eba\u4e0e\u7814\u7a76\u4e0d\u8db3", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 6)

    card(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.2), TAG_PURPLE)
    txt(s, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.3),
        "AI\u6570\u5b57\u4eba\u4e0e\u6587\u5316\u5bfc\u89c8", sz=15, c=TAG_PURPLE, bold=True, fn=SERIF)
    bullets(s, Inches(0.8), Inches(1.95), Inches(5.2), Inches(1.3), [
        "\u8fbd\u5c0f\u535a\uff1a\u77e5\u8bc6\u56fe\u8c31+\u5927\u6a21\u578b\u878d\u5408\u5fae\u8c03\u7684AI\u667a\u6167\u5bfc\u89c8\u7cfb\u7edf",
        "InHeritage\uff1a\u6e38\u620f\u5316+AR\u7684\u6587\u5316\u9057\u4ea7\u4f20\u64ad\u5e94\u7528",
        "\u8d8b\u52bf\uff1a\u4ece\u9759\u6001\u5c55\u793a\u8d70\u5411AI\u9a71\u52a8\u7684\u667a\u80fd\u5bfc\u89c8"], sz=12, sp=Pt(5), bc=TAG_PURPLE)

    card(s, Inches(6.7), Inches(1.3), Inches(6), Inches(2.2), TAG_AMBER)
    txt(s, Inches(7.0), Inches(1.5), Inches(5.4), Inches(0.3),
        "\u77e5\u8bc6\u56fe\u8c31\u589e\u5f3a\u63a8\u8350", sz=15, c=TAG_AMBER, bold=True, fn=SERIF)
    bullets(s, Inches(7.0), Inches(1.95), Inches(5.4), Inches(1.3), [
        "Guo\u7b49\uff1a\u77e5\u8bc6\u56fe\u8c31\u901a\u8fc7\u5b9e\u4f53\u5173\u7cfb\u6269\u5c55\u63a8\u8350\u5019\u9009\u96c6",
        "\u738b\u654f\u7b49\uff1aKG+LLM\u589e\u5f3a\u63a8\u8350\u7cfb\u7edf\u7814\u7a76",
        "\u6c6a\u5929\u96c4\u7b49\uff1a\u975e\u9057\u6570\u5b57\u5316\u53d1\u5c55\u4e2d\u77e5\u8bc6\u56fe\u8c31\u89c6\u89d2\u7684\u524d\u666f"], sz=12, sp=Pt(5), bc=TAG_AMBER)

    rect(s, Inches(0.5), Inches(3.8), Inches(12.2), Inches(3.2), LIGHT_GOLD, LIGHT_TAN, radius=True)
    rect(s, Inches(0.5), Inches(3.82), Inches(12.2), Inches(0.035), CINNABAR)
    txt(s, Inches(0.8), Inches(4.0), Inches(5), Inches(0.35),
        "\u73b0\u6709\u7814\u7a76\u4e0d\u8db3\u4e0e\u672c\u8bfe\u9898\u5b9a\u4f4d", sz=16, c=CINNABAR, bold=True, fn=SERIF)
    bullets(s, Inches(0.8), Inches(4.5), Inches(8), Inches(2), [
        "1.  \u73b0\u6709\u975e\u9057\u5e73\u53f0\u7f3a\u5c11CRS\u5bf9\u8bdd\u63a8\u8350\u80fd\u529b\uff0c\u65e0\u6cd5\u5728\u4ea4\u4e92\u4e2d\u4e3b\u52a8\u53d1\u73b0\u7528\u6237\u504f\u597d",
        "2.  AI\u5bfc\u89c8\u7cfb\u7edf\u591a\u4fa7\u91cd\u77e5\u8bc6\u95ee\u7b54\uff0c\u7f3a\u5c11\u63a8\u8350\u5f15\u5bfc\u4e0e\u504f\u597d\u91c7\u96c6\u7684\u95ed\u73af\u8bbe\u8ba1",
        "3.  \u77e5\u8bc6\u56fe\u8c31\u5728\u975e\u9057\u63a8\u8350\u4e2d\u7684\u5e94\u7528\u5c1a\u5904\u8d77\u6b65\u9636\u6bb5\uff0c\u63a8\u8350\u89e3\u91ca\u6027\u4e0d\u8db3",
        "4.  CRS+AI\u6570\u5b57\u4eba+\u77e5\u8bc6\u56fe\u8c31\u534f\u540c\u65b9\u6848\u5c1a\u672a\u89c1\u62a5\u9053"], sz=12, sp=Pt(6), bc=CINNABAR)

    rect(s, Inches(9.5), Inches(4.1), Inches(3), Inches(2.4), CINNABAR, radius=True)
    txt(s, Inches(9.5), Inches(4.3), Inches(3), Inches(2),
         "\u672c\u8bfe\u9898\u5b9a\u4f4d\n\nCRS\u4e3a\u5f15\u5bfc\n\u56fe\u8c31\u4e3a\u589e\u5f3a\n\u6570\u5b57\u4eba\u4e3a\u8f7d\u4f53",
         sz=13, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=SERIF)

    # === 7 TECH ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "03")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u76f8\u5173\u6280\u672f\u4e0e\u7406\u8bba", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 7)

    for i, (cat, name, desc, accent) in enumerate([
        ("\u524d\u7aef", "\u5fae\u4fe1\u5c0f\u7a0b\u5e8f\nWXML/WXSS/JS", "\u539f\u751f\u5f00\u53d1\n\u5fae\u4fe1\u80fd\u529b\u76f4\u8c03", CINNABAR),
        ("\u540e\u7aef", "FastAPI\nSQLAlchemy", "\u5f02\u6b65\u9ad8\u6027\u80fd\nORM\u591a\u5e93\u5207\u6362", TAG_BLUE),
        ("\u6570\u636e", "SQLite\n\u672c\u5730\u5b58\u50a8", "\u8f7b\u91cf\u5173\u7cfb\u578b\n\u5c0f\u89c4\u6a21\u90e8\u7f72", TAG_GREEN),
        ("AI", "\u8c46\u5305\u5927\u6a21\u578b\n\u672c\u5730\u77e5\u8bc6\u5e93", "\u4e09\u7ea7\u56de\u9000\u7b56\u7565\nTTS\u53cc\u5f15\u64ce", GOLD),
        ("\u63a8\u8350", "CRS\u5bf9\u8bdd\u63a8\u8350\n\u77e5\u8bc6\u56fe\u8c31", "ASK-REC\u72b6\u6001\u673a\n\u89c4\u5219\u52a0\u6743\u6df7\u5408", DARK_RED),
    ]):
        x = Inches(0.4 + i * 2.55)
        card(s, x, Inches(1.4), Inches(2.35), Inches(5.3), accent)
        txt(s, x + Inches(0.15), Inches(1.65), Inches(2.05), Inches(0.3),
            cat, sz=12, c=accent, bold=True, fn=MONO)
        txt(s, x + Inches(0.15), Inches(2.05), Inches(2.05), Inches(1),
            name, sz=15, c=DARK_RED, bold=True, align=PP_ALIGN.CENTER, fn=SERIF)
        rect(s, x + Inches(0.3), Inches(3.2), Inches(1.75), Inches(0.01), LIGHT_TAN)
        txt(s, x + Inches(0.15), Inches(3.4), Inches(2.05), Inches(2.5),
            desc, sz=11, c=MUTED, align=PP_ALIGN.CENTER)

    # === 8 ANALYSIS ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "04")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u7cfb\u7edf\u5206\u6790", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    txt(s, Inches(1.4), Inches(0.95), Inches(10), Inches(0.3),
        "\u4e09\u7c7b\u89d2\u8272\u4e0e\u516d\u5927\u529f\u80fd\u6a21\u5757", sz=12, c=MUTED)
    divider(s)
    pnum(s, 8)

    for i, (role, desc, accent) in enumerate([
        ("\u666e\u901a\u7528\u6237", "\u6d4f\u89c8\u5185\u5bb9\u3001AI\u5bf9\u8bdd\n\u4e2a\u6027\u5316\u63a8\u8350\u3001\u6d3b\u52a8\u62a5\u540d\n\u793e\u533a\u8ba8\u8bba\u3001\u504f\u597d\u7ba1\u7406", CINNABAR),
        ("\u5185\u5bb9\u7ba1\u7406\u5458", "\u5185\u5bb9\u5ba1\u6838\u53d1\u5e03\n\u8d28\u91cf\u8bc4\u5b9a\u3001\u6570\u636e\u6536\u96c6\n\u57fa\u672c\u7528\u6237\u7ba1\u7406", TAG_BLUE),
        ("\u7cfb\u7edf\u7ba1\u7406\u5458", "\u7cfb\u7edf\u8fd0\u7ef4\u914d\u7f6e\n\u8fd0\u884c\u76d1\u63a7\n\u6743\u9650\u63a7\u5236", TAG_GREEN)]):
        x = Inches(0.5 + i * 4.2)
        card(s, x, Inches(1.4), Inches(3.9), Inches(1.6), accent)
        txt(s, x + Inches(0.2), Inches(1.6), Inches(3.5), Inches(0.3),
            role, sz=15, c=accent, bold=True, fn=SERIF)
        txt(s, x + Inches(0.2), Inches(1.95), Inches(3.5), Inches(0.9),
            desc, sz=11, c=BODY)

    mods = [("\u975e\u9057\u5185\u5bb9\u5c55\u793a", CINNABAR), ("\u4e2a\u6027\u5316\u63a8\u8350", TAG_BLUE),
            ("AI\u6570\u5b57\u4eba\u5bf9\u8bdd", TAG_GREEN), ("\u6d3b\u52a8\u62a5\u540d", GOLD),
            ("\u793e\u533a\u4e92\u52a8", TAG_PURPLE), ("\u7528\u6237\u4e2d\u5fc3", TAG_AMBER)]
    for i, (mod, accent) in enumerate(mods):
        col, row = i % 3, i // 3
        x, y = Inches(0.5 + col * 4.2), Inches(3.4 + row * 1.8)
        rect(s, x, y, Inches(3.9), Inches(1.5), BG_CARD, LIGHT_TAN, radius=True)
        rect(s, x, y + Inches(0.02), Inches(0.04), Inches(1.46), accent)
        txt(s, x + Inches(0.25), Inches(3.55 + row * 1.8), Inches(3.4), Inches(0.3),
            mod, sz=14, c=accent, bold=True, fn=SERIF)

    # === 9 FEASIBILITY ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "04")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u53ef\u884c\u6027\u5206\u6790", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 9)

    for i, (title, desc, accent) in enumerate([
        ("\u6280\u672f", "\u2713 \u5fae\u4fe1\u5c0f\u7a0b\u5e8f\u6210\u719f\u6846\u67b6\n\u2713 FastAPI+SQLAlchemy ORM\n\u2713 \u77e5\u8bc6\u5e93-\u641c\u7d22-\u5927\u6a21\u578b\u4e09\u7ea7\u56de\u9000\n\u2713 \u89c4\u5219\u52a0\u6743\u6df7\u5408\u63a8\u8350", CINNABAR),
        ("\u7ecf\u6d4e", "\u2713 SQLite\u96f6\u6210\u672c\n\u2713 \u4e91\u670d\u52a1\u5668\u6708\u8d39\u767e\u5143\u5185\n\u2713 \u8c46\u5305\u5927\u6a21\u578b\u514d\u8d39\u989d\u5ea6\n\u2713 \u5c0f\u7a0b\u5e8f\u6ce8\u518c\u53d1\u5e03\u5e73\u4ef7", TAG_BLUE),
        ("\u64cd\u4f5c", "\u2713 \u5c0f\u7a0b\u5e8f\u514d\u5b89\u88c5\u626b\u7801\u5373\u7528\n\u2713 AI\u6570\u5b57\u4eba\u96f6\u5b66\u4e60\u6210\u672c\n\u2713 \u7ba1\u7406\u7aefWeb\u754c\u9762\u76f4\u89c2\u64cd\u4f5c", TAG_GREEN),
        ("\u793e\u4f1a\u6cd5\u5f8b", "\u2713 \u672c\u5730\u77e5\u8bc6\u5e93\u964d\u4f4eAI\u5e7b\u89c9\n\u2713 \u7b26\u5408\u4e2a\u4eba\u4fe1\u606f\u4fdd\u62a4\u6cd5\n\u2713 \u4f20\u64ad\u975e\u9057\u6709\u79ef\u6781\u793e\u4f1a\u610f\u4e49\n\u2713 \u7b26\u5408\u975e\u7269\u8d28\u6587\u5316\u9057\u4ea7\u6cd5", GOLD)]):
        col, row = i % 2, i // 2
        x, y = Inches(0.5 + col * 6.3), Inches(1.3 + row * 2.9)
        card(s, x, y, Inches(6), Inches(2.5), accent)
        txt(s, x + Inches(0.3), y + Inches(0.2), Inches(5.4), Inches(0.35),
            title, sz=17, c=accent, bold=True, fn=SERIF)
        txt(s, x + Inches(0.3), y + Inches(0.65), Inches(5.4), Inches(1.6),
            desc, sz=12, c=BODY)

    # === 10 ARCH ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "05")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u7cfb\u7edf\u6574\u4f53\u67b6\u6784", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 10)

    for i, (name, content, bg, tc, ac) in enumerate([
        ("\u524d\u7aef\u4ea4\u4e92\u5c42", "\u5fae\u4fe1\u5c0f\u7a0b\u5e8f\u7528\u6237\u7aef  |  Web\u7ba1\u7406\u7aef", BG_CARD, BODY, CINNABAR),
        ("\u63a5\u53e3\u670d\u52a1\u5c42", "AI\u95ee\u7b54 | CRS\u72b6\u6001 | \u7528\u6237\u753b\u50cf | \u5185\u5bb9\u6d3b\u52a8 | \u8ba8\u8bba\u533a | \u767b\u5f55", LIGHT_GOLD, BODY, TAG_BLUE),
        ("\u6838\u5fc3\u80fd\u529b\u5c42", "AI\u670d\u52a1 | \u63a8\u8350\u670d\u52a1 | \u672c\u5730\u77e5\u8bc6\u68c0\u7d22 | \u77e5\u8bc6\u56fe\u8c31\u589e\u5f3a | \u8054\u7f51\u8865\u5145 | TTS\u8bed\u97f3", CINNABAR, WHITE, GOLD),
        ("\u6570\u636e\u4e0e\u5916\u90e8\u670d\u52a1\u5c42", "SQLite | \u672c\u5730\u77e5\u8bc6 | \u63a8\u8350\u65e5\u5fd7 | CRS\u4f1a\u8bdd | \u8c46\u5305API | \u8054\u7f51\u641c\u7d22 | \u5fae\u4fe1\u767b\u5f55", BG_DARK, RGBColor(0xBB, 0xAA, 0x99), GOLD)]):
        y = Inches(1.3 + i * 1.45)
        rect(s, Inches(1.8), y, Inches(10.5), Inches(1.2), bg, LIGHT_TAN, radius=True)
        rect(s, Inches(1.8), y + Inches(0.02), Inches(10.5), Inches(0.035), ac)
        txt(s, Inches(2.1), y + Inches(0.12), Inches(2.5), Inches(0.3),
            name, sz=15, c=ac, bold=True, fn=SERIF)
        txt(s, Inches(2.1), y + Inches(0.5), Inches(10), Inches(0.5),
            content, sz=11, c=tc)

    txt(s, Inches(0.5), Inches(1.5), Inches(1.1), Inches(0.3),
        "\u8bbe\u8ba1\u539f\u5219", sz=11, c=CINNABAR, bold=True, fn=MONO)
    txt(s, Inches(0.5), Inches(1.85), Inches(1.1), Inches(3),
        "\u524d\u540e\u7aef\u5206\u79bb\n\u4e3b\u7ebf\u8d8b\u540c\n\u8f7b\u91cf\u4f18\u5148\n\u53ef\u6269\u5c55", sz=10, c=MUTED)

    # === 11 CRS ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "05")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "CRS\u5bf9\u8bdd\u63a8\u8350\u673a\u5236", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    txt(s, Inches(1.4), Inches(0.95), Inches(10), Inches(0.3),
        "\u4e09\u7ef4\u7f6e\u4fe1\u5ea6\u8bc4\u4f30 + ASK-REC\u72b6\u6001\u673a", sz=12, c=MUTED)
    divider(s)
    pnum(s, 11)

    card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4), CINNABAR)
    txt(s, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.35),
        "\u4e09\u7ef4\u7f6e\u4fe1\u5ea6\u8bc4\u4f30\u6a21\u578b", sz=17, c=CINNABAR, bold=True, fn=SERIF)
    rect(s, Inches(0.8), Inches(2.1), Inches(4.5), Inches(0.5), BG_DARK, LIGHT_TAN, radius=True)
    txt(s, Inches(0.8), Inches(2.15), Inches(4.5), Inches(0.4),
         "C = 0.40\u00d7Se + 0.35\u00d7Si + 0.25\u00d7Sd",
         sz=16, c=GOLD, bold=True, align=PP_ALIGN.CENTER, fn=MONO)
    bullets(s, Inches(0.8), Inches(2.9), Inches(5.2), Inches(3), [
        ("Se \u663e\u5f0f\u504f\u597d  ", "\u7528\u6237\u504f\u597d\u8bbe\u7f6e\u3001ASK\u9009\u9879\u56de\u7b54"),
        ("Si \u9690\u5f0f\u884c\u4e3a  ", "\u66dd\u5149\u3001\u70b9\u51fb\u3001\u6d4f\u89c8\u3001\u6d3b\u52a8\u62a5\u540d\u3001\u8ba8\u8bba\u4e92\u52a8"),
        ("Sd \u5bf9\u8bdd\u8bed\u4e49  ", "AI\u63d0\u95ee\u4e2d\u7684\u975e\u9057\u5173\u952e\u8bcd\u3001\u5730\u533a\u7279\u5f81\u3001\u573a\u666f\u504f\u597d")],
        sz=13, sp=Pt(12), bc=CINNABAR)

    card(s, Inches(6.7), Inches(1.4), Inches(6), Inches(5.4), TAG_BLUE)
    txt(s, Inches(7.0), Inches(1.6), Inches(5.4), Inches(0.35),
        "ASK-REC \u72b6\u6001\u673a", sz=17, c=TAG_BLUE, bold=True, fn=SERIF)
    for i, (mode, th, desc, color) in enumerate([
        ("cold_start", "C < 28", "\u4f18\u5148ASK\u91c7\u96c6\u504f\u597d", TAG_BLUE),
        ("mixed", "28 \u2264 C < 62", "\u8fb9\u63a8\u8350\u8fb9\u8ffd\u95ee", TAG_AMBER),
        ("precision", "C \u2265 62", "\u7cbe\u51c6\u63a8\u8350\u4e3a\u4e3b", TAG_GREEN)]):
        y = Inches(2.2 + i * 1.3)
        rect(s, Inches(7.0), y, Inches(2.2), Inches(0.35), color, radius=True)
        txt(s, Inches(7.0), y + Inches(0.02), Inches(2.2), Inches(0.3),
            mode, sz=12, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=MONO)
        txt(s, Inches(9.4), y, Inches(1.5), Inches(0.3),
            th, sz=11, c=DARK_RED, bold=True, fn=MONO)
        txt(s, Inches(7.0), y + Inches(0.45), Inches(5.2), Inches(0.6),
            desc, sz=12, c=BODY)
    txt(s, Inches(7.0), Inches(5.8), Inches(5.2), Inches(0.6),
        "\u7279\u6b8a\u7b56\u7565\uff1a\u610f\u56fe\u9a71\u52a8\u63a8\u8350 | \u6062\u590d\u5f0f\u63d0\u95ee", sz=11, c=MUTED, fn=MONO)

    # === 12 AI ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "05")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        'AI\u6570\u5b57\u4eba\u201c\u9ed1\u5854\u201d', sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 12)

    for i, (duty, desc, accent) in enumerate([
        ("\u77e5\u8bc6\u4f20\u6388", "\u975e\u9057\u95ee\u9898\u89e3\u91ca\n\u5bfc\u89c8\u5165\u53e3", CINNABAR),
        ("\u63a8\u8350\u5f15\u5bfc", "\u63d0\u95ee\u5f15\u5411\u63a8\u8350\n\u4e0d\u505c\u7559\u5728\u95ee\u7b54", TAG_BLUE),
        ("\u91c7\u96c6\u504f\u597d", "ASK\u6a21\u677f\u6536\u96c6\n\u7c7b\u522b/\u5730\u533a/\u4f53\u9a8c", TAG_GREEN),
        ("\u95ed\u73af\u89e6\u53d1", "\u884c\u4e3a\u53cd\u9988\u753b\u50cf\nCRS\u72b6\u6001\u6df1\u5316", GOLD)]):
        x = Inches(0.5 + i * 3.15)
        card(s, x, Inches(1.3), Inches(2.9), Inches(1.8), accent)
        txt(s, x + Inches(0.15), Inches(1.55), Inches(2.6), Inches(0.3),
            duty, sz=14, c=accent, bold=True, fn=SERIF)
        txt(s, x + Inches(0.15), Inches(1.9), Inches(2.6), Inches(1),
            desc, sz=11, c=BODY, align=PP_ALIGN.CENTER)

    rect(s, Inches(0.5), Inches(3.4), Inches(12.2), Inches(1.4), BG_CARD, LIGHT_TAN, radius=True)
    rect(s, Inches(0.5), Inches(3.42), Inches(12.2), Inches(0.035), GOLD)
    txt(s, Inches(0.8), Inches(3.6), Inches(5), Inches(0.3),
        "AI\u56de\u7b54\u94fe\u8def\uff08\u672c\u5730\u77e5\u8bc6\u4f18\u5148\uff09", sz=14, c=GOLD, bold=True, fn=SERIF)
    txt(s, Inches(0.8), Inches(4.0), Inches(11.6), Inches(0.5),
         "\u63a5\u6536\u95ee\u9898 \u2192 \u68c0\u7d22\u672c\u5730\u77e5\u8bc6\u5e93 \u2192 \u8bfb\u53d6CRS\u4f1a\u8bdd \u2192 \u751f\u6210\u63a8\u8350\u8f7d\u8377 \u2192 ASK-REC\u51b3\u7b56 \u2192 \u8c03\u7528\u8c46\u5305API \u2192 \u8fd4\u56de\u56de\u7b54+\u63a8\u8350\u5361+ASK\u9009\u9879",
         sz=12, c=BODY, fn=MONO)

    rect(s, Inches(0.5), Inches(5.1), Inches(12.2), Inches(1.8), BG_CARD, LIGHT_TAN, radius=True)
    rect(s, Inches(0.5), Inches(5.12), Inches(12.2), Inches(0.035), TAG_PURPLE)
    txt(s, Inches(0.8), Inches(5.3), Inches(5), Inches(0.3),
        "\u4e09\u5c42\u63d0\u793a\u8bcd\u7ea6\u675f", sz=14, c=TAG_PURPLE, bold=True, fn=SERIF)
    bullets(s, Inches(0.8), Inches(5.7), Inches(11.6), Inches(1), [
        "\u7b2c\u4e00\u5c42  \u5168\u5c40\u89d2\u8272\u63d0\u793a\uff1a\u786e\u5b9a\u9ed1\u5854\u8eab\u4efd\u3001\u8bed\u6c14\u3001\u670d\u52a1\u8303\u56f4\u548c\u77e5\u8bc6\u8fb9\u754c",
        "\u7b2c\u4e8c\u5c42  \u4efb\u52a1\u578b\u63d0\u793a\uff1a\u77e5\u8bc6\u547d\u4e2d\u65f6\u6da6\u8272\u7b80\u5316\uff0c\u672a\u547d\u4e2d\u65f6\u9650\u5b9a\u8f93\u51fa\u8303\u56f4",
        "\u7b2c\u4e09\u5c42  CRS\u611f\u77e5\u63d0\u793a\uff1a\u6839\u636ecold_start/mixed/precision\u52a8\u6001\u8c03\u6574\u8bed\u6c14\u6df1\u5ea6"], sz=12, sp=Pt(4), bc=TAG_PURPLE)

    # === 13 KG ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "05")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u77e5\u8bc6\u56fe\u8c31\u4e0e\u5185\u5bb9\u6cbb\u7406", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 13)

    card(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), TAG_BLUE)
    txt(s, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.35),
        "\u77e5\u8bc6\u56fe\u8c31\u8bbe\u8ba1", sz=17, c=TAG_BLUE, bold=True, fn=SERIF)
    bullets(s, Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5), [
        "\u5b9a\u4f4d\uff1a\u63a8\u8350\u548cAI\u5bf9\u8bdd\u7684\u52a0\u5f3a\u5668\uff0c\u975e\u4e3b\u6392\u5e8f\u5668",
        "\u6269\u5927\u63a8\u8350\u5019\u9009\uff1a\u76f8\u8fd1\u5b9e\u4f53\u3001\u7c7b\u522b\u5173\u7cfb\u3001\u8def\u5f84\u5173\u7cfb",
        "\u52a0\u5f3a\u63a8\u8350\u89e3\u91ca\uff1a\u8def\u5f84\u5173\u7cfb\u8f6c\u4e3a\u81ea\u7136\u8bed\u8a00\u63a8\u8350\u7406\u7531",
        "\u4ee5\u7528\u6237\u5df2\u77e5\u5b9e\u4f53\u4e3a\u8d77\u70b9\uff0c\u5411\u4e30\u5bcc\u65b9\u5411\u5ef6\u4f38",
        "CRS\u8054\u52a8\uff1a\u8bc6\u522b\u5b9e\u4f53\u2192\u6620\u5c04\u753b\u50cf\u2192\u9644\u52a0\u56fe\u8c31\u7406\u7531",
        "\u4e00\u77f3\u4e8c\u9e1f\uff1a\u589e\u5f3a\u89e3\u91ca + \u53c2\u4e0e\u5174\u8da3\u5efa\u6a21"], sz=12, sp=Pt(8), bc=TAG_BLUE)

    card(s, Inches(6.7), Inches(1.3), Inches(6), Inches(5.5), TAG_GREEN)
    txt(s, Inches(7.0), Inches(1.5), Inches(5.4), Inches(0.35),
        "\u5185\u5bb9\u6cbb\u7406\u4f53\u7cfb", sz=17, c=TAG_GREEN, bold=True, fn=SERIF)
    bullets(s, Inches(7.0), Inches(2.0), Inches(5.4), Inches(4.5), [
        "\u8d28\u91cf\u8bc4\u5206\uff08\u4e09\u7ef4\u5ea6\u91cf\u5316\uff09\uff1a",
        "  \u5185\u5bb9\u5b8c\u6574\u5ea6\uff1a\u6b63\u6587\u5b57\u6570\u3001\u5c01\u9762\u56fe\u3001\u6458\u8981",
        "  \u4fe1\u606f\u4e30\u5bcc\u5ea6\uff1a\u7ae0\u8282\u5f52\u5c5e\u3001\u6807\u7b7e\u8986\u76d6",
        "  \u539f\u521b\u8d28\u91cf\uff1a\u6765\u6e90\u53ef\u4fe1\u5ea6\u3001\u5185\u5bb9\u54c8\u5e0c\u53bb\u91cd",
        "\u5ba1\u6838\u53d1\u5e03\uff1a\u8d28\u91cf\u5206\u2265\u9608\u503c\u81ea\u52a8\u8fdb\u63a8\u8350\u6c60",
        "\u767d\u540d\u5355\u56de\u8865\uff1a\u5b9a\u671f\u4ece\u8fbe\u6807\u5185\u5bb9\u91cd\u65b0\u9009\u62d4\u7cbe\u9009"], sz=12, sp=Pt(6), bc=TAG_GREEN)

    # === 14-18 IMPLEMENTATION ===
    impl = [
        (14, "06", "\u9996\u9875\u4e0e\u4e2a\u6027\u5316\u63a8\u8350", [
            ("AI\u5bfc\u89c8\u6a2a\u5e45", "CRS\u6a21\u5f0f\u8fdb\u5ea6\u80f6\u56ca\n\u9ed1\u5854\u6570\u5b57\u4eba\u7ec4\u4ef6\n\u8868\u60c5\u968fCRS\u6a21\u5f0f\u53d8\u5316\n\u7acb\u5373\u5f00\u59cb\u6309\u94ae", CINNABAR),
            ("\u7cbe\u9009\u63a8\u8350\u533a", "\u7b2c\u4e00\u6761\u63a8\u8350\u5927\u5361\u7247\n\u4eca\u65e5\u63a8\u8350\uff1a\u5185\u5bb9/\u6d3b\u52a8/\u8ba8\u8bba\n\u63a8\u8350\u7406\u7531\u6807\u7b7e+\u65b0\u9c9c\u5ea6\u6807\u7b7e\n\u70b9\u51fb\u884c\u4e3a\u5373\u65f6\u56de\u6d41\u753b\u50cf", TAG_BLUE),
            ("\u5feb\u901f\u5165\u53e3+TabBar", "\u975e\u9057\u53d1\u5c55\u53f2/\u975e\u9057\u5730\u70b9\u5feb\u6377\u5165\u53e3\n\u81ea\u5b9a\u4e49TabBar\n\u80f6\u56ca\u5706\u89d2\u8bbe\u8ba1", TAG_GREEN)]),
        (15, "06", "\u975e\u9057\u5185\u5bb9\u4e0e\u6d3b\u52a8\u62a5\u540d", [
            ("\u975e\u9057\u5185\u5bb9\u6a21\u5757", "\u7011\u5e03\u6d41\u53cc\u5217\u5e03\u5c40\n\u7b56\u5c55\u7cbe\u9009\u663e\u773c\u4f4d\u7f6e\n\u6309\u5b50\u7ae0\u8282\u548c\u7c7b\u578b\u6253\u6807\u7b7e\n\u4e2a\u6027\u5316\u5185\u5bb9\u63a8\u8350\u6df7\u5408\u5c55\u793a\n\u5185\u5bb9\u8be6\u60c5\uff1a\u5c01\u9762+\u4e09\u5206\u949f\u770b\u70b9+\u5ef6\u4f38\u8bb2\u89e3\nAI\u6d6e\u7a97\u5165\u53e3", CINNABAR),
            ("\u6d3b\u52a8\u62a5\u540d\u6a21\u5757", "\u6d3b\u52a8\u5217\u8868\uff1a\u6309\u6708\u4efd\u7b5b\u9009\n\u524d\u4e24\u6761\u6807\u8bb0\u672c\u6708\u4e3b\u63a8\n\u4e2a\u6027\u5316\u6d3b\u52a8\u63a8\u8350\u6df7\u5408\u5c55\u793a\n\u6d3b\u52a8\u8be6\u60c5\uff1a\u5c01\u9762/\u65f6\u95f4\u5730\u70b9/\u7ec4\u7ec7\u8005\n\u5728\u7ebf\u62a5\u540d/\u53d6\u6d88\u62a5\u540d+AI\u6d6e\u7a97", TAG_GREEN)]),
        (16, "06", "\u793e\u533a\u8ba8\u8bba\u4e0e\u7528\u6237\u4e2d\u5fc3", [
            ("\u793e\u533a\u8ba8\u8bba\u6a21\u5757", "\u8ba8\u8bba\u5217\u8868\uff1a\u5173\u952e\u8bcd\u641c\u7d22+\u6807\u7b7e\u7b5b\u9009\n\u6392\u5e8f\u5207\u6362\uff08\u70ed\u95e8/\u6700\u65b0\uff09\u3001\u53ea\u770b\u6536\u85cf\n\u70ed\u95e8\u8bdd\u9898TOP3\u5c55\u793a\n\u53d1\u5e16\u4e09\u79cd\u6a21\u677f\uff1a\u63d0\u95ee\u578b/\u4f53\u9a8c\u5206\u4eab\u578b/\u6d3b\u52a8\u53cd\u9988\u578b\n\u8bdd\u9898\u8be6\u60c5\uff1a\u5e16\u5b50\u6b63\u6587+\u8bc4\u8bba\u5217\u8868", TAG_BLUE),
            ("\u7528\u6237\u4e2d\u5fc3\u6a21\u5757", "\u7528\u6237\u753b\u50cf\u5361\u7247\uff1a\u504f\u597d\u5173\u952e\u8bcd/\u6d3b\u8dc3\u573a\u666f/\u5173\u6ce8\u5730\u533a\n\u7edf\u8ba1\u5361\u7247\uff1a\u53d1\u5e16\u6570/\u6d3b\u52a8\u8bb0\u5f55/\u5df2\u9009\u504f\u597d\n\u504f\u597d\u8bbe\u7f6e\uff1a\u975e\u9057\u7c7b\u522b/\u573a\u666f\u7c7b\u578b/\u5730\u533a\n\u4fdd\u5b58\u540e\u4f5c\u4e3a\u663e\u5f0f\u4fe1\u53f7\u7eb3\u5165\u753b\u50cf", GOLD)]),
        (17, "06", 'AI\u6570\u5b57\u4eba\u201c\u9ed1\u5854\u201d\u6838\u5fc3\u4ea4\u4e92', [
            ("cold_start\n\u51b7\u542f\u52a8", "C < 28\n\u4e3b\u52a8\u63d0\u95ee\u4e3a\u4e3b\nASK\u5361\u7247\u6536\u96c6\u504f\u597d\n\u7c7b\u76ee\u2192\u5730\u533a\u2192\u573a\u666f\u2192\u7a0b\u5ea6\n\u4e0d\u76f4\u63a5\u7ed9\u51fa\u63a8\u8350\u7ed3\u679c\n\u6570\u5b57\u4eba\u597d\u5947\u5f20\u671b\u8868\u60c5", TAG_BLUE),
            ("mixed\n\u6df7\u5408\u6a21\u5f0f", "28 \u2264 C < 62\n\u8fb9\u63a8\u8350\u8fb9\u8ffd\u95ee\n1-2\u5f20AI\u63a8\u8350\u5361+B\u7ec4\u8ffd\u95ee\n\u4e0a\u65b9\u63a8\u8350+\u4e0b\u65b9ASK\u9009\u9879\n\u63a8\u8350\u7406\u7531\u8d34\u5408\u5f53\u524d\u9009\u62e9\n\u6570\u5b57\u4eba\u6258\u816e\u601d\u8003\u8868\u60c5", TAG_AMBER),
            ("precision\n\u7cbe\u51c6\u6a21\u5f0f", "C \u2265 62\n\u505c\u6b62\u8ffd\u95ee\u76f4\u63a5\u8f93\u51fa\n\u9ad8\u5ea6\u5339\u914d\u63a8\u8350\u7ed3\u679c\n\u7cbe\u51c6\u5339\u914d\u504f\u597d\u8868\u8fbe\n\u63a8\u8350\u5361\u5360\u4e3b\u8981\u533a\u57df\n\u6570\u5b57\u4eba\u81ea\u4fe1\u5fae\u7b11\u8868\u60c5", TAG_GREEN)]),
        (18, "06", "AI\u5bf9\u8bdd\u529f\u80fd\u4e0e\u4fe1\u606f\u5c42\u6b21", [
            ("\u8bed\u97f3\u64ad\u62a5", "\u8c46\u5305TTS\u4f18\u5148\nEdge-TTS\u964d\u7ea7\n\u53cc\u5f15\u64ce\u67b6\u6784", CINNABAR),
            ("\u5ef6\u4f38\u63a8\u8350", "\u56de\u7b54\u4e0b\u65b9\u5c55\u793a\n\u5185\u5bb9/\u6d3b\u52a8/\u8ba8\u8bba\n\u63a8\u8350\u5361\u7247\u5373\u65f6\u5237\u65b0", TAG_BLUE),
            ("\u884c\u52a8\u6e05\u5355", "\u6d4f\u89c8/\u62a5\u540d/\u8ba8\u8bba\n\u4e0b\u4e00\u6b65\u884c\u52a8\u5efa\u8bae\n\u95ed\u73af\u89e6\u53d1\u673a\u5236", TAG_GREEN),
            ("AI\u6d6e\u7a97", "\u5185\u5bb9\u9875/\u6d3b\u52a8\u9875\u5d4c\u5165\n\u573a\u666f\u5316\u95ee\u7b54\n\u968f\u65f6\u53d1\u8d77\u5bf9\u8bdd", GOLD)]),
    ]

    for sn, sec, title, content in impl:
        s = prs.slides.add_slide(bl)
        set_bg(s, BG_WARM)
        rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
        sec_num(s, sec)
        txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
            title, sz=26, c=DARK_RED, bold=True, fn=SERIF)
        divider(s)
        pnum(s, sn)

        if sn == 18:
            for i, (t, d, ac) in enumerate(content):
                x = Inches(0.5 + i * 3.15)
                card(s, x, Inches(1.4), Inches(2.9), Inches(2), ac)
                txt(s, x + Inches(0.15), Inches(1.6), Inches(2.6), Inches(0.3),
                    t, sz=13, c=ac, bold=True, fn=SERIF)
                txt(s, x + Inches(0.15), Inches(1.95), Inches(2.6), Inches(1.2),
                    d, sz=11, c=BODY, align=PP_ALIGN.CENTER)
            rect(s, Inches(0.5), Inches(3.7), Inches(12.2), Inches(3.2), BG_CARD, LIGHT_TAN, radius=True)
            rect(s, Inches(0.5), Inches(3.72), Inches(12.2), Inches(0.035), GOLD)
            txt(s, Inches(0.8), Inches(3.9), Inches(5), Inches(0.3),
                "AI\u5bf9\u8bdd\u9875\u4e09\u5c42\u4fe1\u606f\u7ed3\u6784", sz=15, c=GOLD, bold=True, fn=SERIF)
            bullets(s, Inches(0.8), Inches(4.4), Inches(11.6), Inches(2), [
                ("\u7b2c\u4e00\u5c42  \u4efb\u52a1\u4e3b\u7ebf\u5c42  ", "\u9ed1\u5854\u6570\u5b57\u4eba\u5f62\u8c61 | \u5bf9\u8bdd\u6d88\u606f\u533a | \u8f93\u5165\u6846 | \u5f53\u524d\u6a21\u5f0f\u63d0\u793a | ASK\u63d0\u95ee\u5361\u7247"),
                ("\u7b2c\u4e8c\u5c42  \u63a8\u8350\u7ed3\u679c\u5c42  ", "AI\u63a8\u8350\u5361 | \u63a8\u8350\u7406\u7531 | \u4e0b\u4e00\u6b65\u5efa\u8bae\uff08\u6d4f\u89c8/\u62a5\u540d/\u8ba8\u8bba\uff09"),
                ("\u7b2c\u4e09\u5c42  \u89e3\u91ca\u5c42  ", "\u7f6e\u4fe1\u5ea6\u8be6\u7ec6\u90e8\u5206 | \u7b56\u7565\u89e3\u91ca | \u56fe\u8c31\u4f9d\u636e")],
                sz=12, sp=Pt(10), bc=GOLD)
        elif sn == 17:
            for i, (t, d, ac) in enumerate(content):
                x = Inches(0.5 + i * 4.2)
                card(s, x, Inches(1.4), Inches(3.9), Inches(5.3), ac)
                rect(s, x + Inches(0.5), Inches(1.6), Inches(2.9), Inches(0.7), ac, radius=True)
                txt(s, x + Inches(0.5), Inches(1.6), Inches(2.9), Inches(0.7),
                    t, sz=13, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=SERIF)
                txt(s, x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(3.8),
                    d, sz=11, c=BODY)
        else:
            for i, (t, d, ac) in enumerate(content):
                if len(content) == 3:
                    x = Inches(0.5 + i * 4.2)
                    w = Inches(3.9)
                else:
                    x = Inches(0.5) if i == 0 else Inches(6.7)
                    w = Inches(5.8)
                card(s, x, Inches(1.4), w, Inches(5.3), ac)
                txt(s, x + Inches(0.3), Inches(1.6), w - Inches(0.6), Inches(0.3),
                    t, sz=15, c=ac, bold=True, fn=SERIF)
                txt(s, x + Inches(0.3), Inches(2.0), w - Inches(0.6), Inches(4.5),
                    d, sz=11, c=BODY)

    # === 19 ADMIN ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "06")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "Web\u7ba1\u7406\u7aef", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 19)

    for i, (t, ac) in enumerate([("\u5185\u5bb9\u7ba1\u7406", CINNABAR), ("\u8d28\u91cf\u68c0\u67e5", TAG_BLUE),
                                  ("\u6d3b\u52a8\u7ba1\u7406", TAG_GREEN), ("\u5e16\u5b50\u7ba1\u7406", GOLD),
                                  ("\u7528\u6237\u7ba1\u7406", TAG_PURPLE), ("\u7edf\u8ba1\u5bfc\u51fa", TAG_AMBER),
                                  ("\u77e5\u8bc6\u5e93\u7ba1\u7406", TAG_BLUE)]):
        col, row = i % 4, i // 4
        x, y = Inches(0.4 + col * 3.2), Inches(1.3 + row * 2.8)
        rect(s, x, y, Inches(3), Inches(2.4), BG_CARD, LIGHT_TAN, radius=True)
        rect(s, x, y + Inches(0.02), Inches(3), Inches(0.035), ac)
        txt(s, x + Inches(0.15), y + Inches(0.2), Inches(2.7), Inches(0.3),
            t, sz=14, c=ac, bold=True, fn=SERIF)

    # === 20 TESTING ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "07")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u7cfb\u7edf\u6d4b\u8bd5", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 20)

    for i, (title, accent, items) in enumerate([
        ("\u529f\u80fd\u6d4b\u8bd5", CINNABAR, [
            "\u9ed1\u76d2\u6d4b\u8bd5\uff1a10\u6761\u7528\u4f8b\uff0c\u901a\u8fc710\u6761",
            "\u8986\u76d68\u6a21\u5757\uff1a\u8ba4\u8bc1/\u5185\u5bb9/\u6d3b\u52a8/\u8ba8\u8bba/\u63a8\u8350/AI/\u56fe\u8c31/\u7ba1\u7406",
            "\u767d\u76d2\u6d4b\u8bd5\uff1a7\u6761\u7528\u4f8b\uff0c\u901a\u8fc77\u6761",
            "CRS\u51b3\u7b56\u5f15\u64ce/ASK\u6a21\u677f/\u77e5\u8bc6\u56fe\u8c31/\u6570\u636e\u6a21\u578b"]),
        ("\u6027\u80fd\u6d4b\u8bd5", TAG_BLUE, [
            "API\u5e73\u5747\u54cd\u5e94\u65f6\u95f4\uff1a43.16ms",
            "\u6700\u5927\u54cd\u5e94\u65f6\u95f4\uff1a51.33ms\uff08\u5747<100ms\uff09",
            "\u5e76\u53d110\u7528\u623750\u8bf7\u6c42\uff1a\u6210\u529f\u7387100%",
            "\u541e\u5410\u91cf225.5\u8bf7\u6c42/\u79d2",
            "AI\u5bf9\u8bdd\u5e73\u5747\u54cd\u5e94\uff1a7.9\u79d2"]),
        ("\u517c\u5bb9\u6027\u6d4b\u8bd5", TAG_GREEN, [
            "\u6d4f\u89c8\u5668\uff1aChrome/Firefox/Safari/Edge",
            "\u5fae\u4fe1\u57fa\u7840\u5e93\uff1a\u22652.20.0",
            "\u8bbe\u5907\uff1aiOS + Android",
            "\u6570\u636e\u5e93\uff1aSQLite/MySQL/PostgreSQL"])]):
        x = Inches(0.5 + i * 4.2)
        card(s, x, Inches(1.4), Inches(3.9), Inches(5.3), accent)
        txt(s, x + Inches(0.3), Inches(1.65), Inches(3.3), Inches(0.3),
            title, sz=15, c=accent, bold=True, fn=SERIF)
        bullets(s, x + Inches(0.3), Inches(2.1), Inches(3.3), Inches(4), items, sz=12, sp=Pt(8), bc=accent)

    # === 21 EXPERIMENT ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "07")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u5b9e\u9a8c\u4e0e\u5206\u6790", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 21)

    card(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), CINNABAR)
    txt(s, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.35),
        "AI\u5bf9\u8bdd\u8d28\u91cf\u8bc4\u4f30", sz=17, c=CINNABAR, bold=True, fn=SERIF)
    rect(s, Inches(0.8), Inches(2.0), Inches(3), Inches(0.42), CINNABAR, radius=True)
    txt(s, Inches(0.8), Inches(2.02), Inches(3), Inches(0.38),
         "\u77e5\u8bc6\u5e93\u547d\u4e2d\u7387 90.0%", sz=14, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=MONO)
    bullets(s, Inches(0.8), Inches(2.7), Inches(5.2), Inches(3.5), [
        "20\u4e2a\u6d4b\u8bd5\u95ee\u9898\uff0c15\u4e2a\u9884\u671f\u547d\u4e2d\u5168\u90e8\u547d\u4e2d",
        "5\u4e2a\u51b7\u95e8\u95ee\u9898\u4e2d3\u4e2a\u672a\u547d\u4e2d\u2192\u56de\u9000\u8c46\u5305/\u8054\u7f51",
        "\u547d\u4e2d\u9879\u7f6e\u4fe1\u5ea6\u96c6\u4e2d\u57280.80-0.85\u533a\u95f4",
        "P50=2.6\u79d2\uff0cP90=4.3\u79d2",
        "\u6d41\u5f0f\u8f93\u51fa3\u79d2\u5185\u53ef\u89c1\u9996\u4e2a\u56de\u7b54\u7247\u6bb5",
        "KB\u547d\u4e2d+\u6da6\u8272\u5360\u6bd4\u6700\u5927\uff08\u94fe\u8def\u6709\u6548\uff09",
        "\u8c46\u5305\u76f4\u7b54+\u7ec4\u5408\u5360\u6bd4\u8f83\u5927\uff08\u5927\u6a21\u578b\u4f53\u73b0\uff09"], sz=12, sp=Pt(5), bc=CINNABAR)

    card(s, Inches(6.7), Inches(1.3), Inches(6), Inches(5.5), TAG_BLUE)
    txt(s, Inches(7.0), Inches(1.5), Inches(5.4), Inches(0.35),
        "\u63a8\u8350\u8bc4\u4ef7\u6307\u6807", sz=17, c=TAG_BLUE, bold=True, fn=SERIF)
    for i, (m, d) in enumerate([("Precision@5", "\u63a8\u8350\u5217\u8868\u524d5\u9879\u4e2d\u4e0e\u7528\u6237\u504f\u597d\u76f8\u5173\u7684\u6bd4\u4f8b\n\u2192 \u53cd\u6620\u63a8\u8350\u51c6\u786e\u6027"),
                                 ("Diversity@5", "\u63a8\u8350\u5217\u8868\u8986\u76d6\u7684\u975e\u9057\u7c7b\u522b\u6570/\u5217\u8868\u957f\u5ea6\n\u2192 \u53cd\u6620\u63a8\u8350\u591a\u6837\u6027"),
                                 ("Coverage", "\u63a8\u8350\u7ed3\u679c\u8986\u76d6\u5185\u5bb9\u5360\u5168\u90e8\u5df2\u53d1\u5e03\u5185\u5bb9\u7684\u6bd4\u4f8b\n\u2192 \u53cd\u6620\u63a8\u8350\u8986\u76d6\u5e7f\u5ea6"),
                                 ("NDCG@5", "\u8003\u8651\u4f4d\u7f6e\u6743\u91cd\u7684\u6392\u5e8f\u8d28\u91cf\n\u2192 \u53cd\u6620\u63a8\u8350\u6392\u5e8f\u8d28\u91cf")]):
        y = Inches(2.1 + i * 1.1)
        rect(s, Inches(7.0), y, Inches(2.2), Inches(0.32), TAG_BLUE, radius=True)
        txt(s, Inches(7.0), y + Inches(0.01), Inches(2.2), Inches(0.3),
            m, sz=11, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=MONO)
        txt(s, Inches(9.4), y, Inches(3), Inches(0.9),
            d, sz=10, c=BODY)

    # === 22 CONCLUSION ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_WARM)
    rect(s, Inches(0), Inches(0), Inches(0.07), Inches(7.5), CINNABAR)
    sec_num(s, "08")
    txt(s, Inches(1.4), Inches(0.5), Inches(8), Inches(0.5),
        "\u7ed3\u8bba\u4e0e\u5c55\u671b", sz=26, c=DARK_RED, bold=True, fn=SERIF)
    divider(s)
    pnum(s, 22)

    card(s, Inches(0.5), Inches(1.3), Inches(7.5), Inches(5.5), CINNABAR)
    txt(s, Inches(0.8), Inches(1.5), Inches(7), Inches(0.35),
        "\u7814\u7a76\u7ed3\u8bba", sz=17, c=CINNABAR, bold=True, fn=SERIF)
    bullets(s, Inches(0.8), Inches(2.0), Inches(7), Inches(4), [
        "\u6784\u5efa\u4e86\u57fa\u4e8eCRS\u63a8\u8350\u4e0eAI\u6570\u5b57\u4eba\u7684\u975e\u9057\u6587\u5316\u4f20\u64ad\u5fae\u4fe1\u5c0f\u7a0b\u5e8f",
        "AI\u6570\u5b57\u4eba\u9ed1\u5854\u901a\u8fc7ASK\u8ffd\u95ee\u9010\u6b65\u6784\u5efa\u7528\u6237\u504f\u597d\u753b\u50cf",
        "\u4ece\u51b7\u542f\u52a8\u5230\u7cbe\u51c6\u63a8\u8350\u53ea\u97003-4\u8f6e\u4ea4\u4e92\u5373\u53ef\u6536\u655b",
        "\u672c\u5730\u77e5\u8bc6\u4f18\u5148+\u5927\u6a21\u578b\u56de\u9000\u7b56\u7565\u517c\u987e\u77e5\u8bc6\u51c6\u786e\u6027\u4e0e\u8bed\u4e49\u7406\u89e3",
        "ASK-REC\u51b3\u7b56\u5f15\u64ce\u6267\u884c\u4e09\u9636\u6bb5\u7f6e\u4fe1\u5ea6\u6536\u655b\u673a\u5236",
        "\u77e5\u8bc6\u56fe\u8c31\u5b9e\u73b0\u5b9e\u4f53\u504f\u597d\u6620\u5c04\u548c\u63a8\u8350\u89e3\u91ca\u751f\u6210",
        "\u7cfb\u7edf\u6027\u80fd\u6ee1\u8db3\u4e2d\u5c0f\u89c4\u6a21\u4f7f\u7528\u573a\u666f"], sz=13, sp=Pt(7), bc=CINNABAR)

    rect(s, Inches(8.5), Inches(1.3), Inches(4.3), Inches(5.5), LIGHT_GOLD, LIGHT_TAN, radius=True)
    rect(s, Inches(8.5), Inches(1.32), Inches(4.3), Inches(0.035), GOLD)
    txt(s, Inches(8.8), Inches(1.5), Inches(3.7), Inches(0.35),
        "\u672a\u6765\u5c55\u671b", sz=17, c=GOLD, bold=True, fn=SERIF)
    bullets(s, Inches(8.8), Inches(2.0), Inches(3.7), Inches(4), [
        "\u52a0\u5165\u66f4\u7ec6\u81f4\u7684\u77e5\u8bc6\u56fe\u8c31\u63a8\u7406\u80fd\u529b",
        "\u540c\u7ebf\u4e0b\u975e\u9057\u573a\u9986\u6df1\u5ea6\u5bf9\u63a5",
        "\u652f\u6301\u591a\u6a21\u6001\u4ea4\u4e92",
        "\u5927\u8bed\u8a00\u6a21\u578b\u6280\u672f\u53d1\u5c55\u63a8\u52a8\u5bf9\u8bdd\u66f4\u81ea\u7136"], sz=13, sp=Pt(14), bc=GOLD)

    # === 23 THANKS ===
    s = prs.slides.add_slide(bl)
    set_bg(s, BG_DARK)
    rect(s, Inches(0), Inches(0), Inches(0.1), Inches(7.5), CINNABAR)
    rect(s, Inches(0.1), Inches(0), Inches(0.03), Inches(7.5), GOLD)
    rect(s, Inches(2), Inches(2.3), Inches(9.3), Inches(0.02), GOLD)
    txt(s, Inches(2), Inches(2.6), Inches(9.3), Inches(0.8),
        "\u611f\u8c22\u5404\u4f4d\u8001\u5e08\u7684\u6307\u5bfc\u4e0e\u8bc4\u5ba1", sz=34, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=SERIF)
    rect(s, Inches(2), Inches(3.6), Inches(9.3), Inches(0.02), GOLD)
    txt(s, Inches(2), Inches(4.1), Inches(9.3), Inches(0.4),
        "\u6073\u8bf7\u5404\u4f4d\u8001\u5e08\u6279\u8bc4\u6307\u6b63", sz=18, c=GOLD, align=PP_ALIGN.CENTER, fn=SERIF)
    txt(s, Inches(2), Inches(5.3), Inches(9.3), Inches(0.35),
        "\u7b54\u8fa9\u4eba\uff1a\u738b\u5b50\u8f69    \u6307\u5bfc\u6559\u5e08\uff1a\u5f90\u9f99\u7434 \u6559\u6388  \u9ad8\u9759 \u9ad8\u7ea7\u5de5\u7a0b\u5e08",
        sz=13, c=MUTED, align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(5.7), Inches(9.3), Inches(0.3),
        "\u4ef2\u607a\u519c\u4e1a\u5de5\u7a0b\u5b66\u9662  \u4fe1\u606f\u79d1\u5b66\u4e0e\u6280\u672f\u5b66\u9662  2026.05.09",
        sz=11, c=MUTED, align=PP_ALIGN.CENTER, fn=MONO)

    out = r"d:\桌面\毕业设计\答辩演示文件夹\答辩PPT_v4_基于CRS推荐与AI数字人的非遗文化传播系统.pptx"
    prs.save(out)
    print(f"PPT\u751f\u6210\u5b8c\u6210: {out}")
    print(f"\u6587\u4ef6\u5927\u5c0f: {os.path.getsize(out)} \u5b57\u8282")
    print(f"\u603b\u9875\u6570: {len(prs.slides)}")


if __name__ == "__main__":
    create()
