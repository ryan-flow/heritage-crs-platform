# -*- coding: utf-8 -*-
import sys; sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

pptx_path = r"d:\桌面\毕业设计\答辩演示文件夹\智能推荐系统：从被动呈现到主动对话 (3).pptx"
prs = Presentation(pptx_path)

print(f"总页数: {len(prs.slides)}")
print("="*60)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*60}")
    print(f"Slide {i}")
    print(f"{'='*60}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(f"  {text[:120]}")
        if shape.has_table:
            table = shape.table
            print(f"  [表格 {len(table.rows)}行x{len(table.columns)}列]")
            for ri, row in enumerate(table.rows):
                cells = [cell.text.strip()[:30] for cell in row.cells]
                print(f"    R{ri}: {' | '.join(cells)}")
                if ri > 5:
                    print(f"    ...省略")
                    break
