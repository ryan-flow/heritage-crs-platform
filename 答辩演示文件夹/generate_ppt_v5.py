from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

W = Inches(13.333)
H = Inches(7.5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFA, 0xF8, 0xF5)
RED = RGBColor(0xC4, 0x1E, 0x3A)
DARK_RED = RGBColor(0x8B, 0x1A, 0x2B)
GOLD = RGBColor(0xC9, 0x9A, 0x4E)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT = RGBColor(0xF0, 0xEB, 0xE3)
SERIF = "Source Han Serif SC"
SANS = "Source Han Sans SC"
MONO = "Consolas"
TOTAL = 22


def bg(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c


def box(s, l, t, w, h, fill, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    return sh


def t(s, l, t_, w, h, text, sz=16, c=BLACK, bold=False, align=PP_ALIGN.LEFT, fn=SANS):
    tb = s.shapes.add_textbox(l, t_, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = bold
    p.font.name = fn
    p.alignment = align
    return tb


def ml(s, l, t_, w, h, lines_data, default_sz=13, default_c=BLACK):
    tb = s.shapes.add_textbox(l, t_, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ld in enumerate(lines_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(ld, str):
            ld = {"t": ld}
        p.space_after = ld.get("sp", Pt(4))
        text = ld.get("t", "")
        sz = ld.get("s", default_sz)
        color = ld.get("c", default_c)
        bold = ld.get("b", False)
        fn = ld.get("f", SANS)
        align = ld.get("a", PP_ALIGN.LEFT)
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = fn
        p.alignment = align
    return tb


def pn(s, n):
    t(s, Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.2),
      f"{n}/{TOTAL}", sz=8, c=GRAY, align=PP_ALIGN.RIGHT, fn=MONO)


def create():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    bl = prs.slide_layouts[6]

    # ── 1 COVER ──
    s = prs.slides.add_slide(bl)
    bg(s, WHITE)
    box(s, Inches(0), Inches(0), Inches(5.2), Inches(7.5), RED)
    t(s, Inches(0.6), Inches(1.0), Inches(4), Inches(0.3),
      "GRADUATION DEFENSE", sz=9, c=RGBColor(0xFF, 0xCC, 0xCC), fn=MONO)
    t(s, Inches(0.6), Inches(1.5), Inches(4.2), Inches(2.5),
      "\u57fa\u4e8eCRS\u63a8\u8350\n\u4e0eAI\u6570\u5b57\u4eba\u7684\n\u975e\u9057\u6587\u5316\u4f20\u64ad\n\u7cfb\u7edf\u5b9e\u73b0",
      sz=32, c=WHITE, bold=True, fn=SERIF)
    box(s, Inches(0.6), Inches(4.5), Inches(2), Inches(0.02), GOLD)
    t(s, Inches(0.6), Inches(4.8), Inches(4), Inches(2.2),
      "\u738b\u5b50\u8f69\n202210274225\n\u6570\u636e\u79d1\u5b66\u4e0e\u5927\u6570\u636e\u6280\u672f222\u73ed\n\u4fe1\u606f\u79d1\u5b66\u4e0e\u6280\u672f\u5b66\u9662\n\u6307\u5bfc\u6559\u5e08\uff1a\u5f90\u9f99\u7434 \u6559\u6388  \u9ad8\u9759 \u9ad8\u7ea7\u5de5\u7a0b\u5e08",
      sz=13, c=RGBColor(0xFF, 0xDD, 0xDD))
    t(s, Inches(0.6), Inches(6.8), Inches(4), Inches(0.3),
      "\u4ef2\u607a\u519c\u4e1a\u5de5\u7a0b\u5b66\u9662  2026.05.09", sz=11, c=GOLD, fn=SANS)

    t(s, Inches(6.0), Inches(2.0), Inches(6.5), Inches(4),
      "CRS\nAI Digital Human\nKnowledge Graph\nWeChat Mini-Program\nIntangible Cultural Heritage",
      sz=14, c=GRAY, fn=MONO, align=PP_ALIGN.RIGHT)

    # ── 2 TOC ──
    s = prs.slides.add_slide(bl)
    bg(s, WHITE)
    box(s, Inches(0), Inches(0), Inches(0.5), Inches(7.5), RED)
    t(s, Inches(1.0), Inches(0.5), Inches(3), Inches(0.3),
      "CONTENTS", sz=9, c=RED, fn=MONO)
    t(s, Inches(1.0), Inches(0.8), Inches(5), Inches(0.6),
      "\u6c47\u62a5\u63d0\u7eb2", sz=30, c=BLACK, bold=True, fn=SERIF)
    box(s, Inches(1.0), Inches(1.5), Inches(11), Inches(0.01), LIGHT)
    pn(s, 2)

    toc = ["01  \u7814\u7a76\u80cc\u666f\u4e0e\u610f\u4e49", "02  \u56fd\u5185\u5916\u7814\u7a76\u73b0\u72b6",
           "03  \u76f8\u5173\u6280\u672f\u4e0e\u7406\u8bba", "04  \u7cfb\u7edf\u5206\u6790",
           "05  \u7cfb\u7edf\u8bbe\u8ba1", "06  \u7cfb\u7edf\u5b9e\u73b0",
           "07  \u7cfb\u7edf\u6d4b\u8bd5\u4e0e\u5b9e\u9a8c", "08  \u7ed3\u8bba\u4e0e\u5c55\u671b"]
    for i, item in enumerate(toc):
        y = Inches(1.8 + i * 0.65)
        is_left = i < 4
        x = Inches(1.0) if is_left else Inches(7.0)
        t(s, x, y, Inches(5.5), Inches(0.5),
          item, sz=18, c=BLACK if is_left else GRAY, bold=is_left, fn=SERIF)
        if i < 4:
            box(s, Inches(1.0), y + Inches(0.45), Inches(5), Inches(0.008), LIGHT)

    # ── helper: content page with left red bar + big title ──
    def content_page(num, sec_title, sec_subtitle=""):
        s = prs.slides.add_slide(bl)
        bg(s, WHITE)
        box(s, Inches(0), Inches(0), Inches(0.5), Inches(7.5), RED)
        t(s, Inches(0.05), Inches(0.3), Inches(0.4), Inches(0.3),
          str(num).zfill(2), sz=10, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=MONO)
        t(s, Inches(1.0), Inches(0.4), Inches(10), Inches(0.6),
          sec_title, sz=28, c=BLACK, bold=True, fn=SERIF)
        if sec_subtitle:
            t(s, Inches(1.0), Inches(1.0), Inches(10), Inches(0.3),
              sec_subtitle, sz=12, c=GRAY)
        box(s, Inches(1.0), Inches(1.35), Inches(11.5), Inches(0.008), LIGHT)
        pn(s, num + 2)
        return s

    # ── 3 BACKGROUND ──
    s = content_page(1, "\u7814\u7a76\u80cc\u666f", "\u975e\u9057\u6570\u5b57\u5316\u4f20\u64ad\u7684\u65f6\u4ee3\u9700\u6c42\u4e0e\u73b0\u6709\u5e73\u53f0\u4e0d\u8db3")

    box(s, Inches(1.0), Inches(1.6), Inches(5.5), Inches(5.5), OFF_WHITE)
    t(s, Inches(1.3), Inches(1.8), Inches(5), Inches(0.4),
      "\u65f6\u4ee3\u9700\u6c42", sz=20, c=RED, bold=True, fn=SERIF)
    ml(s, Inches(1.3), Inches(2.4), Inches(4.8), Inches(4.2), [
        {"t": "\u56fd\u5bb6\u7ea7\u975e\u9057\u4ee3\u8868\u6027\u9879\u76ee1000+\uff0c\u8054\u5408\u56fd\u975e\u9057\u540d\u5f55\u6570\u91cf\u4e16\u754c\u7b2c\u4e00", "s": 13, "sp": Pt(8)},
        {"t": "\u77ed\u89c6\u9891\u5316\uff1a\u6296\u97f3\u975e\u9057\u89c6\u9891\u8d852.2\u4ebf\u6761\uff0c138\u4ebf\u6b21\u70b9\u8d5e", "s": 13, "sp": Pt(8)},
        {"t": "\u8d44\u6e90\u6570\u5b57\u5316\uff1a\u52a8\u753b\u3001\u6e38\u620f\u3001\u4ea7\u54c1\u8bbe\u8ba1\u9700\u6570\u5b57\u5316\u975e\u9057\u7d20\u6750", "s": 13, "sp": Pt(8)},
        {"t": "\u4f53\u9a8c\u5316\u8f6c\u578b\uff1a\u975e\u9057\u4ece\u535a\u7269\u9986\u9644\u5c5e\u8d70\u5411\u72ec\u7acb\u6d3b\u52a8\u4ea7\u54c1", "s": 13, "sp": Pt(8)},
        {"t": "\u4f20\u64ad\u91cd\u5fc3\u4ece\u4e13\u4e1a\u9886\u57df\u8f6c\u5411\u5927\u4f17\u89c6\u91ce", "s": 13},
    ])

    t(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.4),
      "\u73b0\u6709\u5e73\u53f0\u4e94\u5927\u4e0d\u8db3", sz=20, c=BLACK, bold=True, fn=SERIF)
    gaps = [("\u63a8\u8350\u80fd\u529b\u5f31", "\u70ed\u95e8\u6392\u5e8f\u4e3a\u4e3b\uff0c\u65e0\u4e2a\u6027\u5316\u63a8\u8350"),
            ("\u4ea4\u4e92\u65b9\u5f0f\u5355\u4e00", "\u56fe\u6587\u5c55\u793a\u4e3a\u4e3b\uff0c\u7f3a\u5c11\u8bed\u97f3\u5bf9\u8bdd\u4ea4\u4e92"),
            ("\u573a\u666f\u611f\u77e5\u7f3a\u5931", "\u4e0d\u80fd\u6355\u6349\u65f6\u95f4\u5730\u70b9\u6587\u5316\u60c5\u5883\u4fe1\u53f7"),
            ("\u51b7\u542f\u52a8\u4e25\u91cd", "\u7528\u6237\u884c\u4e3a\u7a00\u758f\uff0c\u65b0\u7528\u6237\u6d41\u5931\u7387\u9ad8"),
            ("\u63a8\u8350\u4e0d\u53ef\u89e3\u91ca", "\u4e0d\u7ed9\u51fa\u63a8\u8350\u7406\u7531\uff0c\u5f71\u54cd\u4fe1\u4efb\u5ea6")]
    for i, (label, desc) in enumerate(gaps):
        y = Inches(2.4 + i * 0.9)
        t(s, Inches(7.0), y, Inches(2), Inches(0.3),
          label, sz=14, c=RED, bold=True, fn=SERIF)
        t(s, Inches(7.0), y + Inches(0.3), Inches(5.5), Inches(0.5),
          desc, sz=12, c=GRAY)

    # ── 4 SIGNIFICANCE ──
    s = content_page(1, "\u7814\u7a76\u610f\u4e49")

    box(s, Inches(1.0), Inches(1.6), Inches(5.5), Inches(5.5), RED)
    t(s, Inches(1.3), Inches(1.8), Inches(5), Inches(0.4),
      "\u7406\u8bba\u610f\u4e49", sz=20, c=WHITE, bold=True, fn=SERIF)
    ml(s, Inches(1.3), Inches(2.4), Inches(4.8), Inches(4.2), [
        {"t": "\u5c06CRS\u5bf9\u8bdd\u63a8\u8350\u601d\u60f3\u5f15\u5165\u975e\u9057\u6587\u5316\u4f20\u64ad\u573a\u666f\uff0c\u63a2\u7d22SAUR\u4ea4\u4e92\u6a21\u5f0f", "s": 13, "c": WHITE, "sp": Pt(10)},
        {"t": "\u63d0\u51fa\u4ee5CRS\u4e3a\u5f15\u5bfc\u3001\u77e5\u8bc6\u56fe\u8c31\u4e3a\u589e\u5f3a\u3001AI\u6570\u5b57\u4eba\u4e3a\u8f7d\u4f53\u7684\u667a\u80fd\u5316\u670d\u52a1\u6846\u67b6", "s": 13, "c": WHITE, "sp": Pt(10)},
        {"t": "\u6784\u5efa\u4e09\u7ef4\u7f6e\u4fe1\u5ea6\u8bc4\u4f30\u6a21\u578b\uff0c\u5b9e\u73b0\u51b7\u542f\u52a8\u5230\u7cbe\u51c6\u63a8\u8350\u7684\u6e10\u8fdb\u6536\u655b", "s": 13, "c": WHITE, "sp": Pt(10)},
        {"t": "\u4e3a\u975e\u9057\u63a8\u8350\u9886\u57df\u63d0\u4f9b\u53ef\u89e3\u91ca\u7684\u63a8\u8350\u673a\u5236\u8bbe\u8ba1\u53c2\u8003", "s": 13, "c": WHITE},
    ])

    t(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.4),
      "\u5b9e\u8df5\u610f\u4e49", sz=20, c=BLACK, bold=True, fn=SERIF)
    ml(s, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.2), [
        {"t": "\u5f00\u53d1\u5b8c\u6574\u7684\u975e\u9057\u6587\u5316\u4f20\u64ad\u5fae\u4fe1\u5c0f\u7a0b\u5e8f\uff0c\u8986\u76d6\u6d4f\u89c8\u3001\u95ee\u7b54\u3001\u63a8\u8350\u3001\u6d3b\u52a8\u3001\u793e\u533a\u5168\u94fe\u8def", "s": 13, "sp": Pt(10)},
        {"t": "AI\u6570\u5b57\u4eba\u9ed1\u5854\u5b9e\u73b0\u591a\u8f6e\u95ee\u7b54\u3001\u8bed\u97f3\u64ad\u62a5\u3001CRS\u72b6\u6001\u663e\u793a\u3001\u5ef6\u4f38\u63a8\u8350\u3001\u884c\u52a8\u6e05\u5355\u751f\u6210", "s": 13, "sp": Pt(10)},
        {"t": "\u672c\u5730\u77e5\u8bc6\u5e93\u4f18\u5148+\u5927\u6a21\u578b\u56de\u9000\u7b56\u7565\uff0c\u517c\u987e\u77e5\u8bc6\u51c6\u786e\u6027\u4e0e\u8bed\u4e49\u7406\u89e3", "s": 13, "sp": Pt(10)},
        {"t": "\u4e3a\u975e\u9057\u6570\u5b57\u5316\u4f20\u64ad\u4e0e\u4e2a\u6027\u5316\u670d\u52a1\u63d0\u4f9b\u53ef\u843d\u5730\u7684\u5b9e\u8df5\u53c2\u8003", "s": 13},
    ])

    # ── 5 STATUS 1 ──
    s = content_page(2, "\u56fd\u5185\u5916\u7814\u7a76\u73b0\u72b6", "\u63a8\u8350\u7cfb\u7edf\u3001CRS\u3001\u5927\u8bed\u8a00\u6a21\u578b\u878d\u5408")

    cols = [("\u63a8\u8350\u7cfb\u7edf", RED, [
        "\u534f\u540c\u8fc7\u6ee4\uff1a\u4f9d\u8d56\u5386\u53f2\u884c\u4e3a\uff0c\u51b7\u542f\u52a8\u7a81\u51fa",
        "\u5185\u5bb9\u63a8\u8350\uff1a\u7279\u5f81\u5339\u914d\uff0c\u591a\u6837\u6027\u4e0d\u8db3",
        "\u6df7\u5408\u63a8\u8350\uff1a\u878d\u5408\u591a\u6e90\u4fe1\u53f7\uff0c\u9002\u5408\u7a00\u758f\u573a\u666f"]),
        ("\u5bf9\u8bdd\u63a8\u8350CRS", RGBColor(0x2A, 0x5D, 0x8F), [
        "Jannach\u7b49\uff1aCRS\u901a\u8fc7\u591a\u8f6e\u4ea4\u4e92\u83b7\u53d6\u53cd\u9988",
        "\u8d75\u68a6\u5a9b\u7b49\uff1a\u4f20\u7edf\u63a8\u8350\u4e09\u5927\u7f3a\u9677",
        "Wang\u7b49\uff1a\u77e5\u8bc6\u589e\u5f3a\u63d0\u793a\u5b66\u4e60\u7edf\u4e00\u6846\u67b6"]),
        ("LLM+\u63a8\u8350", RGBColor(0x3D, 0x7A, 0x4F), [
        "\u8c22\u5e7f\u660e\u7b49\uff1a\u5224\u522b\u5f0f\u4e0e\u751f\u6210\u5f0f\u8303\u5f0f",
        "Wu\u7b49\uff1aLLM for Recommendation\u7efc\u8ff0",
        "\u8d8b\u52bf\uff1a\u8d70\u5411LLM\u9a71\u52a8\u7684\u4ea4\u4e92\u5f0f\u63a8\u8350"])]
    for i, (title, accent, items) in enumerate(cols):
        x = Inches(1.0 + i * 4.0)
        box(s, x, Inches(1.6), Inches(0.06), Inches(5.2), accent)
        t(s, x + Inches(0.3), Inches(1.7), Inches(3.4), Inches(0.4),
          title, sz=18, c=accent, bold=True, fn=SERIF)
        ml(s, x + Inches(0.3), Inches(2.3), Inches(3.4), Inches(4), [
            {"t": it, "s": 12, "sp": Pt(10)} for it in items
        ])

    # ── 6 STATUS 2 ──
    s = content_page(2, "AI\u6570\u5b57\u4eba\u4e0e\u7814\u7a76\u4e0d\u8db3")

    t(s, Inches(1.0), Inches(1.6), Inches(5.5), Inches(0.4),
      "AI\u6570\u5b57\u4eba\u4e0e\u6587\u5316\u5bfc\u89c8", sz=17, c=RGBColor(0x6B, 0x4C, 0x8A), bold=True, fn=SERIF)
    ml(s, Inches(1.0), Inches(2.1), Inches(5.5), Inches(1.5), [
        {"t": "\u8fbd\u5c0f\u535a\uff1a\u77e5\u8bc6\u56fe\u8c31+\u5927\u6a21\u578b\u878d\u5408\u5fae\u8c03\u7684AI\u667a\u6167\u5bfc\u89c8\u7cfb\u7edf", "s": 12, "sp": Pt(6)},
        {"t": "InHeritage\uff1a\u6e38\u620f\u5316+AR\u7684\u6587\u5316\u9057\u4ea7\u4f20\u64ad\u5e94\u7528", "s": 12, "sp": Pt(6)},
        {"t": "\u8d8b\u52bf\uff1a\u4ece\u9759\u6001\u5c55\u793a\u8d70\u5411AI\u9a71\u52a8\u7684\u667a\u80fd\u5bfc\u89c8", "s": 12},
    ])

    t(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.4),
      "\u77e5\u8bc6\u56fe\u8c31\u589e\u5f3a\u63a8\u8350", sz=17, c=RGBColor(0xB8, 0x73, 0x1A), bold=True, fn=SERIF)
    ml(s, Inches(7.0), Inches(2.1), Inches(5.5), Inches(1.5), [
        {"t": "Guo\u7b49\uff1a\u77e5\u8bc6\u56fe\u8c31\u901a\u8fc7\u5b9e\u4f53\u5173\u7cfb\u6269\u5c55\u63a8\u8350\u5019\u9009\u96c6", "s": 12, "sp": Pt(6)},
        {"t": "\u738b\u654f\u7b49\uff1aKG+LLM\u589e\u5f3a\u63a8\u8350\u7cfb\u7edf\u7814\u7a76", "s": 12, "sp": Pt(6)},
        {"t": "\u6c6a\u5929\u96c4\u7b49\uff1a\u975e\u9057\u6570\u5b57\u5316\u53d1\u5c55\u4e2d\u77e5\u8bc6\u56fe\u8c31\u89c6\u89d2\u7684\u524d\u666f", "s": 12},
    ])

    box(s, Inches(1.0), Inches(3.8), Inches(11.5), Inches(0.01), LIGHT)
    t(s, Inches(1.0), Inches(4.0), Inches(8), Inches(0.4),
      "\u73b0\u6709\u7814\u7a76\u4e0d\u8db3\u4e0e\u672c\u8bfe\u9898\u5b9a\u4f4d", sz=17, c=BLACK, bold=True, fn=SERIF)
    ml(s, Inches(1.0), Inches(4.5), Inches(8), Inches(2.5), [
        {"t": "1.  \u73b0\u6709\u975e\u9057\u5e73\u53f0\u7f3a\u5c11CRS\u5bf9\u8bdd\u63a8\u8350\u80fd\u529b\uff0c\u65e0\u6cd5\u5728\u4ea4\u4e92\u4e2d\u4e3b\u52a8\u53d1\u73b0\u7528\u6237\u504f\u597d", "s": 12, "sp": Pt(6)},
        {"t": "2.  AI\u5bfc\u89c8\u7cfb\u7edf\u591a\u4fa7\u91cd\u77e5\u8bc6\u95ee\u7b54\uff0c\u7f3a\u5c11\u63a8\u8350\u5f15\u5bfc\u4e0e\u504f\u597d\u91c7\u96c6\u7684\u95ed\u73af\u8bbe\u8ba1", "s": 12, "sp": Pt(6)},
        {"t": "3.  \u77e5\u8bc6\u56fe\u8c31\u5728\u975e\u9057\u63a8\u8350\u4e2d\u7684\u5e94\u7528\u5c1a\u5904\u8d77\u6b65\u9636\u6bb5\uff0c\u63a8\u8350\u89e3\u91ca\u6027\u4e0d\u8db3", "s": 12, "sp": Pt(6)},
        {"t": "4.  CRS+AI\u6570\u5b57\u4eba+\u77e5\u8bc6\u56fe\u8c31\u534f\u540c\u65b9\u6848\u5c1a\u672a\u89c1\u62a5\u9053", "s": 12},
    ])

    box(s, Inches(9.5), Inches(4.0), Inches(3.2), Inches(3.0), RED)
    t(s, Inches(9.5), Inches(4.3), Inches(3.2), Inches(2.5),
      "\u672c\u8bfe\u9898\u5b9a\u4f4d\n\nCRS\u4e3a\u5f15\u5bfc\n\u56fe\u8c31\u4e3a\u589e\u5f3a\n\u6570\u5b57\u4eba\u4e3a\u8f7d\u4f53",
      sz=14, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=SERIF)

    # ── 7 TECH ──
    s = content_page(3, "\u76f8\u5173\u6280\u672f\u4e0e\u7406\u8bba")

    techs = [("\u524d\u7aef", "\u5fae\u4fe1\u5c0f\u7a0b\u5e8f\nWXML/WXSS/JS", RED),
             ("\u540e\u7aef", "FastAPI\nSQLAlchemy", RGBColor(0x2A, 0x5D, 0x8F)),
             ("\u6570\u636e", "SQLite", RGBColor(0x3D, 0x7A, 0x4F)),
             ("AI", "\u8c46\u5305\u5927\u6a21\u578b\n\u672c\u5730\u77e5\u8bc6\u5e93", GOLD),
             ("\u63a8\u8350", "CRS\u5bf9\u8bdd\u63a8\u8350\n\u77e5\u8bc6\u56fe\u8c31", DARK_RED)]
    for i, (cat, name, accent) in enumerate(techs):
        x = Inches(1.0 + i * 2.45)
        box(s, x, Inches(1.6), Inches(2.2), Inches(5.2), OFF_WHITE)
        box(s, x, Inches(1.6), Inches(2.2), Inches(0.06), accent)
        t(s, x + Inches(0.2), Inches(1.85), Inches(1.8), Inches(0.3),
          cat, sz=11, c=accent, bold=True, fn=MONO)
        t(s, x + Inches(0.2), Inches(2.3), Inches(1.8), Inches(1.2),
          name, sz=16, c=BLACK, bold=True, fn=SERIF)

    # ── 8 ANALYSIS ──
    s = content_page(4, "\u7cfb\u7edf\u5206\u6790", "\u4e09\u7c7b\u89d2\u8272\u4e0e\u516d\u5927\u529f\u80fd\u6a21\u5757")

    roles = [("\u666e\u901a\u7528\u6237", "\u6d4f\u89c8\u5185\u5bb9\u3001AI\u5bf9\u8bdd\u3001\u4e2a\u6027\u5316\u63a8\u8350\u3001\u6d3b\u52a8\u62a5\u540d\u3001\u793e\u533a\u8ba8\u8bba\u3001\u504f\u597d\u7ba1\u7406"),
             ("\u5185\u5bb9\u7ba1\u7406\u5458", "\u5185\u5bb9\u5ba1\u6838\u53d1\u5e03\u3001\u8d28\u91cf\u8bc4\u5b9a\u3001\u6570\u636e\u6536\u96c6\u3001\u57fa\u672c\u7528\u6237\u7ba1\u7406"),
             ("\u7cfb\u7edf\u7ba1\u7406\u5458", "\u7cfb\u7edf\u8fd0\u7ef4\u914d\u7f6e\u3001\u8fd0\u884c\u76d1\u63a7\u3001\u6743\u9650\u63a7\u5236")]
    for i, (role, desc) in enumerate(roles):
        y = Inches(1.6 + i * 0.7)
        t(s, Inches(1.0), y, Inches(2), Inches(0.3),
          role, sz=14, c=RED, bold=True, fn=SERIF)
        t(s, Inches(3.2), y, Inches(9), Inches(0.5),
          desc, sz=12, c=GRAY)

    box(s, Inches(1.0), Inches(3.7), Inches(11.5), Inches(0.01), LIGHT)
    t(s, Inches(1.0), Inches(3.9), Inches(5), Inches(0.4),
      "\u516d\u5927\u529f\u80fd\u6a21\u5757", sz=17, c=BLACK, bold=True, fn=SERIF)

    mods = ["\u975e\u9057\u5185\u5bb9\u5c55\u793a", "\u4e2a\u6027\u5316\u63a8\u8350", "AI\u6570\u5b57\u4eba\u5bf9\u8bdd",
            "\u6d3b\u52a8\u62a5\u540d", "\u793e\u533a\u4e92\u52a8", "\u7528\u6237\u4e2d\u5fc3"]
    for i, mod in enumerate(mods):
        col, row = i % 3, i // 3
        x = Inches(1.0 + col * 4.0)
        y = Inches(4.5 + row * 1.2)
        box(s, x, y, Inches(3.6), Inches(0.9), OFF_WHITE)
        t(s, x + Inches(0.2), y + Inches(0.15), Inches(3.2), Inches(0.3),
          mod, sz=14, c=BLACK, bold=True, fn=SERIF)

    # ── 9 FEASIBILITY ──
    s = content_page(4, "\u53ef\u884c\u6027\u5206\u6790")

    feas = [("\u6280\u672f", "\u5fae\u4fe1\u5c0f\u7a0b\u5e8f\u6210\u719f\u6846\u67b6 | FastAPI+SQLAlchemy ORM | \u77e5\u8bc6\u5e93-\u641c\u7d22-\u5927\u6a21\u578b\u4e09\u7ea7\u56de\u9000 | \u89c4\u5219\u52a0\u6743\u6df7\u5408\u63a8\u8350"),
            ("\u7ecf\u6d4e", "SQLite\u96f6\u6210\u672c | \u4e91\u670d\u52a1\u5668\u6708\u8d39\u767e\u5143\u5185 | \u8c46\u5305\u5927\u6a21\u578b\u514d\u8d39\u989d\u5ea6 | \u5c0f\u7a0b\u5e8f\u6ce8\u518c\u53d1\u5e03\u5e73\u4ef7"),
            ("\u64cd\u4f5c", "\u5c0f\u7a0b\u5e8f\u514d\u5b89\u88c5\u626b\u7801\u5373\u7528 | AI\u6570\u5b57\u4eba\u96f6\u5b66\u4e60\u6210\u672c | \u7ba1\u7406\u7aefWeb\u754c\u9762\u76f4\u89c2\u64cd\u4f5c"),
            ("\u793e\u4f1a\u6cd5\u5f8b", "\u672c\u5730\u77e5\u8bc6\u5e93\u964d\u4f4eAI\u5e7b\u89c9 | \u7b26\u5408\u4e2a\u4eba\u4fe1\u606f\u4fdd\u62a4\u6cd5 | \u4f20\u64ad\u975e\u9057\u6709\u79ef\u6781\u793e\u4f1a\u610f\u4e49 | \u7b26\u5408\u975e\u7269\u8d28\u6587\u5316\u9057\u4ea7\u6cd5")]
    for i, (title, desc) in enumerate(feas):
        y = Inches(1.6 + i * 1.4)
        t(s, Inches(1.0), y, Inches(1.5), Inches(0.3),
          title, sz=16, c=RED, bold=True, fn=SERIF)
        t(s, Inches(2.8), y, Inches(9.5), Inches(1.0),
          desc, sz=12, c=GRAY)

    # ── 10 ARCH ──
    s = content_page(5, "\u7cfb\u7edf\u6574\u4f53\u67b6\u6784")

    layers = [
        ("\u524d\u7aef\u4ea4\u4e92\u5c42", "\u5fae\u4fe1\u5c0f\u7a0b\u5e8f\u7528\u6237\u7aef  |  Web\u7ba1\u7406\u7aef", OFF_WHITE),
        ("\u63a5\u53e3\u670d\u52a1\u5c42", "AI\u95ee\u7b54 | CRS\u72b6\u6001 | \u7528\u6237\u753b\u50cf | \u5185\u5bb9\u6d3b\u52a8 | \u8ba8\u8bba\u533a | \u767b\u5f55", WHITE),
        ("\u6838\u5fc3\u80fd\u529b\u5c42", "AI\u670d\u52a1 | \u63a8\u8350\u670d\u52a1 | \u672c\u5730\u77e5\u8bc6\u68c0\u7d22 | \u77e5\u8bc6\u56fe\u8c31\u589e\u5f3a | \u8054\u7f51\u8865\u5145 | TTS\u8bed\u97f3", OFF_WHITE),
        ("\u6570\u636e\u4e0e\u5916\u90e8\u670d\u52a1\u5c42", "SQLite | \u672c\u5730\u77e5\u8bc6 | \u63a8\u8350\u65e5\u5fd7 | CRS\u4f1a\u8bdd | \u8c46\u5305API | \u8054\u7f51\u641c\u7d22 | \u5fae\u4fe1\u767b\u5f55", WHITE)]
    for i, (name, content, bg_c) in enumerate(layers):
        y = Inches(1.6 + i * 1.4)
        box(s, Inches(1.0), y, Inches(11.5), Inches(1.15), bg_c, LIGHT)
        t(s, Inches(1.3), y + Inches(0.1), Inches(2.5), Inches(0.3),
          name, sz=14, c=RED, bold=True, fn=SERIF)
        t(s, Inches(1.3), y + Inches(0.45), Inches(10.5), Inches(0.5),
          content, sz=11, c=GRAY)

    # ── 11 CRS ──
    s = content_page(5, "CRS\u5bf9\u8bdd\u63a8\u8350\u673a\u5236", "\u4e09\u7ef4\u7f6e\u4fe1\u5ea6\u8bc4\u4f30 + ASK-REC\u72b6\u6001\u673a")

    box(s, Inches(1.0), Inches(1.6), Inches(5.5), Inches(0.6), BLACK)
    t(s, Inches(1.3), Inches(1.65), Inches(5), Inches(0.5),
      "C = 0.40\u00d7Se + 0.35\u00d7Si + 0.25\u00d7Sd",
      sz=18, c=GOLD, bold=True, align=PP_ALIGN.CENTER, fn=MONO)

    ml(s, Inches(1.0), Inches(2.4), Inches(5.5), Inches(4), [
        {"t": "Se \u663e\u5f0f\u504f\u597d", "s": 14, "c": RED, "b": True, "sp": Pt(4)},
        {"t": "\u7528\u6237\u504f\u597d\u8bbe\u7f6e\u3001ASK\u9009\u9879\u56de\u7b54", "s": 12, "sp": Pt(12)},
        {"t": "Si \u9690\u5f0f\u884c\u4e3a", "s": 14, "c": RED, "b": True, "sp": Pt(4)},
        {"t": "\u66dd\u5149\u3001\u70b9\u51fb\u3001\u6d4f\u89c8\u3001\u6d3b\u52a8\u62a5\u540d\u3001\u8ba8\u8bba\u4e92\u52a8", "s": 12, "sp": Pt(12)},
        {"t": "Sd \u5bf9\u8bdd\u8bed\u4e49", "s": 14, "c": RED, "b": True, "sp": Pt(4)},
        {"t": "AI\u63d0\u95ee\u4e2d\u7684\u975e\u9057\u5173\u952e\u8bcd\u3001\u5730\u533a\u7279\u5f81\u3001\u573a\u666f\u504f\u597d", "s": 12},
    ])

    t(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.4),
      "ASK-REC \u72b6\u6001\u673a", sz=18, c=BLACK, bold=True, fn=SERIF)

    modes = [("cold_start", "C < 28", "\u4f18\u5148ASK\u91c7\u96c6\u504f\u597d", RGBColor(0x2A, 0x5D, 0x8F)),
             ("mixed", "28 \u2264 C < 62", "\u8fb9\u63a8\u8350\u8fb9\u8ffd\u95ee", RGBColor(0xB8, 0x73, 0x1A)),
             ("precision", "C \u2265 62", "\u7cbe\u51c6\u63a8\u8350\u4e3a\u4e3b", RGBColor(0x3D, 0x7A, 0x4F))]
    for i, (mode, th, desc, color) in enumerate(modes):
        y = Inches(2.2 + i * 1.5)
        box(s, Inches(7.0), y, Inches(0.08), Inches(1.2), color)
        t(s, Inches(7.3), y, Inches(2), Inches(0.3),
          mode, sz=14, c=color, bold=True, fn=MONO)
        t(s, Inches(7.3), y + Inches(0.3), Inches(2), Inches(0.25),
          th, sz=11, c=GRAY, fn=MONO)
        t(s, Inches(7.3), y + Inches(0.6), Inches(5), Inches(0.4),
          desc, sz=12, c=BLACK)

    # ── 12 AI ──
    s = content_page(5, 'AI\u6570\u5b57\u4eba\u201c\u9ed1\u5854\u201d')

    duties = [("\u77e5\u8bc6\u4f20\u6388", "\u975e\u9057\u95ee\u9898\u89e3\u91ca\uff0c\u5bfc\u89c8\u5165\u53e3"),
              ("\u63a8\u8350\u5f15\u5bfc", "\u63d0\u95ee\u5f15\u5411\u63a8\u8350\uff0c\u4e0d\u505c\u7559\u5728\u95ee\u7b54"),
              ("\u91c7\u96c6\u504f\u597d", "ASK\u6a21\u677f\u6536\u96c6\u7c7b\u522b/\u5730\u533a/\u4f53\u9a8c"),
              ("\u95ed\u73af\u89e6\u53d1", "\u884c\u4e3a\u53cd\u9988\u753b\u50cf\uff0cCRS\u72b6\u6001\u6df1\u5316")]
    for i, (duty, desc) in enumerate(duties):
        y = Inches(1.6 + i * 0.7)
        t(s, Inches(1.0), y, Inches(2), Inches(0.3),
          duty, sz=14, c=RED, bold=True, fn=SERIF)
        t(s, Inches(3.2), y, Inches(4), Inches(0.5),
          desc, sz=12, c=GRAY)

    box(s, Inches(1.0), Inches(4.4), Inches(11.5), Inches(0.01), LIGHT)
    t(s, Inches(1.0), Inches(4.6), Inches(5), Inches(0.3),
      "AI\u56de\u7b54\u94fe\u8def", sz=16, c=BLACK, bold=True, fn=SERIF)
    t(s, Inches(1.0), Inches(5.0), Inches(11.5), Inches(0.4),
      "\u63a5\u6536\u95ee\u9898 \u2192 \u68c0\u7d22\u672c\u5730\u77e5\u8bc6\u5e93 \u2192 \u8bfb\u53d6CRS\u4f1a\u8bdd \u2192 \u751f\u6210\u63a8\u8350\u8f7d\u8377 \u2192 ASK-REC\u51b3\u7b56 \u2192 \u8c03\u7528\u8c46\u5305API \u2192 \u8fd4\u56de\u56de\u7b54+\u63a8\u8350\u5361+ASK\u9009\u9879",
      sz=11, c=GRAY, fn=MONO)

    t(s, Inches(1.0), Inches(5.6), Inches(5), Inches(0.3),
      "\u4e09\u5c42\u63d0\u793a\u8bcd\u7ea6\u675f", sz=16, c=BLACK, bold=True, fn=SERIF)
    ml(s, Inches(1.0), Inches(6.0), Inches(11.5), Inches(1.2), [
        {"t": "\u7b2c\u4e00\u5c42  \u5168\u5c40\u89d2\u8272\u63d0\u793a\uff1a\u786e\u5b9a\u9ed1\u5854\u8eab\u4efd\u3001\u8bed\u6c14\u3001\u670d\u52a1\u8303\u56f4\u548c\u77e5\u8bc6\u8fb9\u754c", "s": 11, "sp": Pt(4)},
        {"t": "\u7b2c\u4e8c\u5c42  \u4efb\u52a1\u578b\u63d0\u793a\uff1a\u77e5\u8bc6\u547d\u4e2d\u65f6\u6da6\u8272\u7b80\u5316\uff0c\u672a\u547d\u4e2d\u65f6\u9650\u5b9a\u8f93\u51fa\u8303\u56f4", "s": 11, "sp": Pt(4)},
        {"t": "\u7b2c\u4e09\u5c42  CRS\u611f\u77e5\u63d0\u793a\uff1a\u6839\u636ecold_start/mixed/precision\u52a8\u6001\u8c03\u6574\u8bed\u6c14\u6df1\u5ea6", "s": 11},
    ])

    # ── 13 KG ──
    s = content_page(5, "\u77e5\u8bc6\u56fe\u8c31\u4e0e\u5185\u5bb9\u6cbb\u7406")

    t(s, Inches(1.0), Inches(1.6), Inches(5.5), Inches(0.4),
      "\u77e5\u8bc6\u56fe\u8c31\u8bbe\u8ba1", sz=18, c=RGBColor(0x2A, 0x5D, 0x8F), bold=True, fn=SERIF)
    box(s, Inches(1.0), Inches(2.1), Inches(0.06), Inches(4.5), RGBColor(0x2A, 0x5D, 0x8F))
    ml(s, Inches(1.3), Inches(2.2), Inches(5), Inches(4), [
        {"t": "\u5b9a\u4f4d\uff1a\u63a8\u8350\u548cAI\u5bf9\u8bdd\u7684\u52a0\u5f3a\u5668\uff0c\u975e\u4e3b\u6392\u5e8f\u5668", "s": 12, "sp": Pt(8)},
        {"t": "\u6269\u5927\u63a8\u8350\u5019\u9009\uff1a\u76f8\u8fd1\u5b9e\u4f53\u3001\u7c7b\u522b\u5173\u7cfb\u3001\u8def\u5f84\u5173\u7cfb", "s": 12, "sp": Pt(8)},
        {"t": "\u52a0\u5f3a\u63a8\u8350\u89e3\u91ca\uff1a\u8def\u5f84\u5173\u7cfb\u8f6c\u4e3a\u81ea\u7136\u8bed\u8a00\u63a8\u8350\u7406\u7531", "s": 12, "sp": Pt(8)},
        {"t": "CRS\u8054\u52a8\uff1a\u8bc6\u522b\u5b9e\u4f53\u2192\u6620\u5c04\u753b\u50cf\u2192\u9644\u52a0\u56fe\u8c31\u7406\u7531", "s": 12, "sp": Pt(8)},
        {"t": "\u4e00\u77f3\u4e8c\u9e1f\uff1a\u589e\u5f3a\u89e3\u91ca + \u53c2\u4e0e\u5174\u8da3\u5efa\u6a21", "s": 12},
    ])

    t(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.4),
      "\u5185\u5bb9\u6cbb\u7406\u4f53\u7cfb", sz=18, c=RGBColor(0x3D, 0x7A, 0x4F), bold=True, fn=SERIF)
    box(s, Inches(7.0), Inches(2.1), Inches(0.06), Inches(4.5), RGBColor(0x3D, 0x7A, 0x4F))
    ml(s, Inches(7.3), Inches(2.2), Inches(5), Inches(4), [
        {"t": "\u8d28\u91cf\u8bc4\u5206\uff08\u4e09\u7ef4\u5ea6\u91cf\u5316\uff09\uff1a", "s": 12, "b": True, "sp": Pt(4)},
        {"t": "\u5185\u5bb9\u5b8c\u6574\u5ea6\uff1a\u6b63\u6587\u5b57\u6570\u3001\u5c01\u9762\u56fe\u3001\u6458\u8981", "s": 12, "sp": Pt(6)},
        {"t": "\u4fe1\u606f\u4e30\u5bcc\u5ea6\uff1a\u7ae0\u8282\u5f52\u5c5e\u3001\u6807\u7b7e\u8986\u76d6", "s": 12, "sp": Pt(6)},
        {"t": "\u539f\u521b\u8d28\u91cf\uff1a\u6765\u6e90\u53ef\u4fe1\u5ea6\u3001\u5185\u5bb9\u54c8\u5e0c\u53bb\u91cd", "s": 12, "sp": Pt(8)},
        {"t": "\u5ba1\u6838\u53d1\u5e03\uff1a\u8d28\u91cf\u5206\u2265\u9608\u503c\u81ea\u52a8\u8fdb\u63a8\u8350\u6c60", "s": 12, "sp": Pt(8)},
        {"t": "\u767d\u540d\u5355\u56de\u8865\uff1a\u5b9a\u671f\u4ece\u8fbe\u6807\u5185\u5bb9\u91cd\u65b0\u9009\u62d4\u7cbe\u9009", "s": 12},
    ])

    # ── 14-18 IMPLEMENTATION (simplified) ──
    impl_pages = [
        (14, "\u9996\u9875\u4e0e\u4e2a\u6027\u5316\u63a8\u8350",
         [("AI\u5bfc\u89c8\u6a2a\u5e45", "CRS\u6a21\u5f0f\u8fdb\u5ea6\u80f6\u56ca | \u9ed1\u5854\u6570\u5b57\u4eba\u7ec4\u4ef6 | \u8868\u60c5\u968fCRS\u6a21\u5f0f\u53d8\u5316 | \u7acb\u5373\u5f00\u59cb\u6309\u94ae"),
          ("\u7cbe\u9009\u63a8\u8350\u533a", "\u7b2c\u4e00\u6761\u63a8\u8350\u5927\u5361\u7247 | \u4eca\u65e5\u63a8\u8350\uff1a\u5185\u5bb9/\u6d3b\u52a8/\u8ba8\u8bba | \u63a8\u8350\u7406\u7531\u6807\u7b7e | \u70b9\u51fb\u884c\u4e3a\u5373\u65f6\u56de\u6d41\u753b\u50cf"),
          ("\u5feb\u901f\u5165\u53e3+TabBar", "\u975e\u9057\u53d1\u5c55\u53f2/\u975e\u9057\u5730\u70b9\u5feb\u6377\u5165\u53e3 | \u81ea\u5b9a\u4e49TabBar | \u80f6\u56ca\u5706\u89d2\u8bbe\u8ba1")]),
        (15, "\u975e\u9057\u5185\u5bb9\u4e0e\u6d3b\u52a8\u62a5\u540d",
         [("\u975e\u9057\u5185\u5bb9\u6a21\u5757", "\u7011\u5e03\u6d41\u53cc\u5217\u5e03\u5c40 | \u7b56\u5c55\u7cbe\u9009 | \u5b50\u7ae0\u8282\u6807\u7b7e | \u4e2a\u6027\u5316\u6df7\u5408\u5c55\u793a | \u4e09\u5206\u949f\u770b\u70b9+\u5ef6\u4f38\u8bb2\u89e3 | AI\u6d6e\u7a97"),
          ("\u6d3b\u52a8\u62a5\u540d\u6a21\u5757", "\u6309\u6708\u4efd\u7b5b\u9009 | \u672c\u6708\u4e3b\u63a8 | \u4e2a\u6027\u5316\u6d3b\u52a8\u63a8\u8350 | \u5728\u7ebf\u62a5\u540d/\u53d6\u6d88 | AI\u6d6e\u7a97")]),
        (16, "\u793e\u533a\u8ba8\u8bba\u4e0e\u7528\u6237\u4e2d\u5fc3",
         [("\u793e\u533a\u8ba8\u8bba", "\u5173\u952e\u8bcd\u641c\u7d22+\u6807\u7b7e\u7b5b\u9009 | \u6392\u5e8f\u5207\u6362 | \u70ed\u95e8\u8bdd\u9898TOP3 | \u4e09\u79cd\u53d1\u5e16\u6a21\u677f | \u8bc4\u8bba\u70b9\u8d5e\u6536\u85cf"),
          ("\u7528\u6237\u4e2d\u5fc3", "\u753b\u50cf\u5361\u7247\uff1a\u504f\u597d/\u573a\u666f/\u5730\u533a | \u7edf\u8ba1\u5361\u7247 | \u504f\u597d\u8bbe\u7f6e\uff1a\u7c7b\u522b/\u573a\u666f/\u5730\u533a | \u4fdd\u5b58\u540e\u7eb3\u5165\u753b\u50cf")]),
        (17, 'AI\u6570\u5b57\u4eba\u201c\u9ed1\u5854\u201d\u6838\u5fc3\u4ea4\u4e92',
         [("cold_start  C<28", "\u4e3b\u52a8\u63d0\u95ee\u4e3a\u4e3b | ASK\u5361\u7247\u6536\u96c6\u504f\u597d | \u7c7b\u76ee\u2192\u5730\u533a\u2192\u573a\u666f\u2192\u7a0b\u5ea6 | \u4e0d\u76f4\u63a5\u7ed9\u51fa\u63a8\u8350 | \u597d\u5947\u5f20\u671b\u8868\u60c5"),
          ("mixed  28\u2264C<62", "\u8fb9\u63a8\u8350\u8fb9\u8ffd\u95ee | 1-2\u5f20AI\u63a8\u8350\u5361+B\u7ec4\u8ffd\u95ee | \u4e0a\u65b9\u63a8\u8350+\u4e0b\u65b9ASK | \u6258\u816e\u601d\u8003\u8868\u60c5"),
          ("precision  C\u226562", "\u505c\u6b62\u8ffd\u95ee\u76f4\u63a5\u8f93\u51fa | \u9ad8\u5ea6\u5339\u914d\u63a8\u8350 | \u63a8\u8350\u5361\u5360\u4e3b\u8981\u533a\u57df | \u81ea\u4fe1\u5fae\u7b11\u8868\u60c5")]),
        (18, "AI\u5bf9\u8bdd\u529f\u80fd\u4e0e\u4fe1\u606f\u5c42\u6b21",
         [("\u8bed\u97f3\u64ad\u62a5", "\u8c46\u5305TTS\u4f18\u5148 | Edge-TTS\u964d\u7ea7 | \u53cc\u5f15\u64ce\u67b6\u6784"),
          ("\u5ef6\u4f38\u63a8\u8350", "\u56de\u7b54\u4e0b\u65b9\u5c55\u793a\u76f8\u5173\u5185\u5bb9/\u6d3b\u52a8/\u8ba8\u8bba | \u63a8\u8350\u5361\u7247\u5373\u65f6\u5237\u65b0"),
          ("\u884c\u52a8\u6e05\u5355", "\u6d4f\u89c8/\u62a5\u540d/\u8ba8\u8bba\u4e0b\u4e00\u6b65\u884c\u52a8\u5efa\u8bae | \u95ed\u73af\u89e6\u53d1\u673a\u5236"),
          ("AI\u6d6e\u7a97", "\u5185\u5bb9\u9875/\u6d3b\u52a8\u9875\u5d4c\u5165 | \u573a\u666f\u5316\u95ee\u7b54 | \u968f\u65f6\u53d1\u8d77\u5bf9\u8bdd")]),
    ]

    for sn, title, items in impl_pages:
        s = content_page(sn - 2, title)
        for i, (label, desc) in enumerate(items):
            y = Inches(1.6 + i * 1.5)
            box(s, Inches(1.0), y, Inches(0.06), Inches(1.2), RED)
            t(s, Inches(1.3), y, Inches(3), Inches(0.3),
              label, sz=15, c=RED, bold=True, fn=SERIF)
            t(s, Inches(1.3), y + Inches(0.35), Inches(11), Inches(0.8),
              desc, sz=12, c=GRAY)
            if i < len(items) - 1:
                box(s, Inches(1.0), y + Inches(1.3), Inches(11.5), Inches(0.005), LIGHT)

    # ── 19 ADMIN ──
    s = content_page(6, "Web\u7ba1\u7406\u7aef")

    admins = ["\u5185\u5bb9\u7ba1\u7406", "\u8d28\u91cf\u68c0\u67e5", "\u6d3b\u52a8\u7ba1\u7406", "\u5e16\u5b50\u7ba1\u7406", "\u7528\u6237\u7ba1\u7406", "\u7edf\u8ba1\u5bfc\u51fa", "\u77e5\u8bc6\u5e93\u7ba1\u7406"]
    for i, name in enumerate(admins):
        col, row = i % 4, i // 4
        x = Inches(1.0 + col * 3.0)
        y = Inches(1.6 + row * 2.5)
        box(s, x, y, Inches(2.7), Inches(2.0), OFF_WHITE, LIGHT)
        t(s, x + Inches(0.2), y + Inches(0.2), Inches(2.3), Inches(0.3),
          name, sz=15, c=BLACK, bold=True, fn=SERIF)

    # ── 20 TESTING ──
    s = content_page(7, "\u7cfb\u7edf\u6d4b\u8bd5")

    tests = [("\u529f\u80fd\u6d4b\u8bd5", "\u9ed1\u76d210\u6761\u7528\u4f8b\u5168\u901a\u8fc7 | \u767d\u76d27\u6761\u7528\u4f8b\u5168\u901a\u8fc7 | \u8986\u76d68\u6a21\u5757"),
             ("\u6027\u80fd\u6d4b\u8bd5", "API\u5e73\u5747\u54cd\u5e9443.16ms | \u5e76\u53d110\u7528\u6237\u6210\u529f\u7387100% | \u541e\u5410\u91cf225.5\u8bf7\u6c42/\u79d2 | AI\u5bf9\u8bdd\u5e73\u57477.9\u79d2"),
             ("\u517c\u5bb9\u6027\u6d4b\u8bd5", "Chrome/Firefox/Safari/Edge | \u5fae\u4fe1\u57fa\u7840\u5e93\u22652.20.0 | iOS+Android | SQLite/MySQL/PostgreSQL")]
    for i, (title, desc) in enumerate(tests):
        y = Inches(1.6 + i * 1.8)
        t(s, Inches(1.0), y, Inches(2), Inches(0.3),
          title, sz=16, c=RED, bold=True, fn=SERIF)
        t(s, Inches(3.2), y, Inches(9.5), Inches(1.5),
          desc, sz=12, c=GRAY)

    # ── 21 EXPERIMENT ──
    s = content_page(7, "\u5b9e\u9a8c\u4e0e\u5206\u6790")

    box(s, Inches(1.0), Inches(1.6), Inches(2.5), Inches(0.5), RED)
    t(s, Inches(1.0), Inches(1.63), Inches(2.5), Inches(0.45),
      "\u77e5\u8bc6\u5e93\u547d\u4e2d\u7387 90.0%", sz=16, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=MONO)

    ml(s, Inches(1.0), Inches(2.3), Inches(5.5), Inches(4), [
        {"t": "20\u4e2a\u6d4b\u8bd5\u95ee\u9898\uff0c15\u4e2a\u9884\u671f\u547d\u4e2d\u5168\u90e8\u547d\u4e2d", "s": 12, "sp": Pt(6)},
        {"t": "5\u4e2a\u51b7\u95e8\u95ee\u9898\u4e2d3\u4e2a\u672a\u547d\u4e2d\u2192\u56de\u9000\u8c46\u5305/\u8054\u7f51", "s": 12, "sp": Pt(6)},
        {"t": "\u547d\u4e2d\u9879\u7f6e\u4fe1\u5ea6\u96c6\u4e2d\u57280.80-0.85\u533a\u95f4", "s": 12, "sp": Pt(6)},
        {"t": "P50=2.6\u79d2\uff0cP90=4.3\u79d2", "s": 12, "sp": Pt(6)},
        {"t": "\u6d41\u5f0f\u8f93\u51fa3\u79d2\u5185\u53ef\u89c1\u9996\u4e2a\u56de\u7b54\u7247\u6bb5", "s": 12, "sp": Pt(6)},
        {"t": "KB\u547d\u4e2d+\u6da6\u8272\u5360\u6bd4\u6700\u5927\uff08\u94fe\u8def\u6709\u6548\uff09", "s": 12, "sp": Pt(6)},
        {"t": "\u8c46\u5305\u76f4\u7b54+\u7ec4\u5408\u5360\u6bd4\u8f83\u5927\uff08\u5927\u6a21\u578b\u4f53\u73b0\uff09", "s": 12},
    ])

    t(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.4),
      "\u63a8\u8350\u8bc4\u4ef7\u6307\u6807", sz=18, c=BLACK, bold=True, fn=SERIF)
    metrics = [("Precision@5", "\u63a8\u8350\u524d5\u9879\u4e2d\u4e0e\u504f\u597d\u76f8\u5173\u7684\u6bd4\u4f8b \u2192 \u51c6\u786e\u6027"),
               ("Diversity@5", "\u63a8\u8350\u8986\u76d6\u975e\u9057\u7c7b\u522b\u6570/\u5217\u8868\u957f\u5ea6 \u2192 \u591a\u6837\u6027"),
               ("Coverage", "\u63a8\u8350\u8986\u76d6\u5185\u5bb9\u5360\u5168\u90e8\u5df2\u53d1\u5e03\u6bd4\u4f8b \u2192 \u5e7f\u5ea6"),
               ("NDCG@5", "\u8003\u8651\u4f4d\u7f6e\u6743\u91cd\u7684\u6392\u5e8f\u8d28\u91cf \u2192 \u6392\u5e8f\u8d28\u91cf")]
    for i, (m, d) in enumerate(metrics):
        y = Inches(2.2 + i * 1.1)
        t(s, Inches(7.0), y, Inches(2.5), Inches(0.3),
          m, sz=13, c=RED, bold=True, fn=MONO)
        t(s, Inches(7.0), y + Inches(0.3), Inches(5.5), Inches(0.6),
          d, sz=11, c=GRAY)

    # ── 22 CONCLUSION ──
    s = content_page(8, "\u7ed3\u8bba\u4e0e\u5c55\u671b")

    ml(s, Inches(1.0), Inches(1.6), Inches(7.5), Inches(5), [
        {"t": "\u7814\u7a76\u7ed3\u8bba", "s": 18, "c": RED, "b": True, "sp": Pt(10)},
        {"t": "\u6784\u5efa\u4e86\u57fa\u4e8eCRS\u63a8\u8350\u4e0eAI\u6570\u5b57\u4eba\u7684\u975e\u9057\u6587\u5316\u4f20\u64ad\u5fae\u4fe1\u5c0f\u7a0b\u5e8f", "s": 13, "sp": Pt(6)},
        {"t": "AI\u6570\u5b57\u4eba\u9ed1\u5854\u901a\u8fc7ASK\u8ffd\u95ee\u9010\u6b65\u6784\u5efa\u7528\u6237\u504f\u597d\u753b\u50cf", "s": 13, "sp": Pt(6)},
        {"t": "\u4ece\u51b7\u542f\u52a8\u5230\u7cbe\u51c6\u63a8\u8350\u53ea\u97003-4\u8f6e\u4ea4\u4e92\u5373\u53ef\u6536\u655b", "s": 13, "sp": Pt(6)},
        {"t": "\u672c\u5730\u77e5\u8bc6\u4f18\u5148+\u5927\u6a21\u578b\u56de\u9000\u517c\u987e\u51c6\u786e\u6027\u4e0e\u8bed\u4e49\u7406\u89e3", "s": 13, "sp": Pt(6)},
        {"t": "ASK-REC\u51b3\u7b56\u5f15\u64ce\u6267\u884c\u4e09\u9636\u6bb5\u7f6e\u4fe1\u5ea6\u6536\u655b\u673a\u5236", "s": 13, "sp": Pt(6)},
        {"t": "\u77e5\u8bc6\u56fe\u8c31\u5b9e\u73b0\u5b9e\u4f53\u504f\u597d\u6620\u5c04\u548c\u63a8\u8350\u89e3\u91ca\u751f\u6210", "s": 13, "sp": Pt(6)},
        {"t": "\u7cfb\u7edf\u6027\u80fd\u6ee1\u8db3\u4e2d\u5c0f\u89c4\u6a21\u4f7f\u7528\u573a\u666f", "s": 13},
    ])

    box(s, Inches(9.0), Inches(1.6), Inches(3.8), Inches(5.5), OFF_WHITE, LIGHT)
    t(s, Inches(9.3), Inches(1.8), Inches(3.2), Inches(0.4),
      "\u672a\u6765\u5c55\u671b", sz=18, c=GOLD, bold=True, fn=SERIF)
    ml(s, Inches(9.3), Inches(2.4), Inches(3.2), Inches(4), [
        {"t": "\u52a0\u5165\u66f4\u7ec6\u81f4\u7684\u77e5\u8bc6\u56fe\u8c31\u63a8\u7406\u80fd\u529b", "s": 13, "sp": Pt(12)},
        {"t": "\u540c\u7ebf\u4e0b\u975e\u9057\u573a\u9986\u6df1\u5ea6\u5bf9\u63a5", "s": 13, "sp": Pt(12)},
        {"t": "\u652f\u6301\u591a\u6a21\u6001\u4ea4\u4e92", "s": 13, "sp": Pt(12)},
        {"t": "\u5927\u8bed\u8a00\u6a21\u578b\u53d1\u5c55\u63a8\u52a8\u5bf9\u8bdd\u66f4\u81ea\u7136", "s": 13},
    ])

    # ── 23 THANKS ──
    s = prs.slides.add_slide(bl)
    bg(s, RED)
    t(s, Inches(2), Inches(2.5), Inches(9.3), Inches(1),
      "\u611f\u8c22\u5404\u4f4d\u8001\u5e08\u7684\u6307\u5bfc\u4e0e\u8bc4\u5ba1",
      sz=36, c=WHITE, bold=True, align=PP_ALIGN.CENTER, fn=SERIF)
    box(s, Inches(4), Inches(3.7), Inches(5.3), Inches(0.02), GOLD)
    t(s, Inches(2), Inches(4.0), Inches(9.3), Inches(0.4),
      "\u6073\u8bf7\u5404\u4f4d\u8001\u5e08\u6279\u8bc4\u6307\u6b63", sz=18, c=GOLD, align=PP_ALIGN.CENTER, fn=SERIF)
    t(s, Inches(2), Inches(5.3), Inches(9.3), Inches(0.3),
      "\u738b\u5b50\u8f69  |  \u5f90\u9f99\u7434 \u6559\u6388  \u9ad8\u9759 \u9ad8\u7ea7\u5de5\u7a0b\u5e08  |  \u4ef2\u607a\u519c\u4e1a\u5de5\u7a0b\u5b66\u9662  |  2026.05.09",
      sz=12, c=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

    out = r"d:\桌面\毕业设计\答辩演示文件夹\答辩PPT_v5_基于CRS推荐与AI数字人的非遗文化传播系统.pptx"
    prs.save(out)
    print(f"PPT\u751f\u6210\u5b8c\u6210: {out}")
    print(f"\u6587\u4ef6\u5927\u5c0f: {os.path.getsize(out)} \u5b57\u8282")
    print(f"\u603b\u9875\u6570: {len(prs.slides)}")


if __name__ == "__main__":
    create()
